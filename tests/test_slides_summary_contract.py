from __future__ import annotations

import csv
import json
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation

from lab_sidecar.slides.service import (
    MAX_CAPTION_CHARS,
    MAX_FIGURES,
    MAX_KEY_COLUMNS,
    MAX_LOG_LINES,
    MAX_NUMERIC_COLUMNS,
    MAX_REPORT_BULLETS,
    MAX_TABLE_COLUMNS,
    MAX_TABLE_ROWS,
)
from tests.test_cli_smoke import copy_examples, extract_task_id, invoke, wait_for_output


REQUIRED_TOP_LEVEL_KEYS = {
    "schema_version",
    "task_id",
    "task_status",
    "template",
    "font_family",
    "font_fallbacks",
    "generated_at",
    "pptx_path",
    "summary_path",
    "generated_from",
    "slide_count",
    "included_figures",
    "included_metrics",
    "warnings",
    "figure_warnings",
    "figure_skipped_candidates",
    "text_truncations",
    "table_truncations",
    "key_comparisons",
    "caption_truncations",
    "slide_evidence",
    "claim_traces",
    "qa_checks",
    "slides",
    "report_excerpt",
    "source_artifacts",
}
COMPAT_ALIAS_KEYS = {"slide_titles", "metrics", "metrics_table", "figures"}
OPTIONAL_UNSTABLE_KEYS = {"project_goal", "full_text_fields"}
TASK_STATUS_VALUES = {"pending", "running", "completed", "failed", "cancelled"}
REQUIRED_INCLUDED_FIGURE_KEYS = {
    "figure_id",
    "chart_type",
    "path",
    "x",
    "y",
    "group_by",
    "source_metrics",
}
REQUIRED_INCLUDED_METRICS_KEYS = {
    "present",
    "path",
    "scenario_summary_path",
    "scenario",
    "row_count",
    "key_columns",
    "numeric",
}
REQUIRED_SLIDE_KEYS = {"slide_index", "title", "purpose", "source_artifacts", "evidence", "empty_source_reason"}
REQUIRED_SLIDE_EVIDENCE_KEYS = {"path", "artifact_id", "role"}
REQUIRED_CLAIM_TRACE_KEYS = {"claim_id", "surface", "claim_type", "value", "evidence"}
REQUIRED_SLIDE_TRACE_KEYS = {"slide_index", "title", "purpose", "evidence_count", "empty_source_reason"}
REQUIRED_QA_GROUPS = {
    "slide_count",
    "empty_slide_check",
    "title_check",
    "artifact_duplicate_check",
    "table_overflow_guard",
    "caption_overflow_guard",
}
FORBIDDEN_EMBEDDED_FRAGMENTS = (
    '"prompt":"secret"',
    '"response":"secret"',
    "ppt/presentation.xml",
    "<p:sld",
)


def read_slides_summary(workspace: Path, task_id: str) -> dict[str, Any]:
    path = workspace / ".lab-sidecar" / "tasks" / task_id / "slides" / "slides-summary.json"
    return json.loads(path.read_text(encoding="utf-8"))


def assert_slides_summary_contract(
    summary: dict[str, Any],
    *,
    workspace: Path,
    expected_task_id: str,
    expected_status: str,
) -> None:
    task_path = workspace / ".lab-sidecar" / "tasks" / expected_task_id
    pptx_path = task_path / "slides" / "presentation-draft.pptx"
    summary_path = task_path / "slides" / "slides-summary.json"

    assert REQUIRED_TOP_LEVEL_KEYS <= set(summary)
    assert COMPAT_ALIAS_KEYS <= set(summary)
    assert summary["schema_version"] == "1"
    assert summary["task_id"] == expected_task_id
    assert summary["task_status"] == expected_status
    assert summary["task_status"] in TASK_STATUS_VALUES
    assert isinstance(summary["template"], str) and summary["template"]
    assert isinstance(summary["font_family"], str) and summary["font_family"]
    assert isinstance(summary["font_fallbacks"], list)
    assert summary["font_fallbacks"]
    assert all(isinstance(item, str) and item for item in summary["font_fallbacks"])
    assert isinstance(summary["generated_at"], str) and summary["generated_at"]
    assert summary["pptx_path"] == "slides/presentation-draft.pptx"
    assert summary["summary_path"] == "slides/slides-summary.json"
    assert pptx_path.is_file()
    assert summary_path.is_file()
    assert pptx_path.stat().st_size > 0
    assert summary_path.stat().st_size > 0

    deck = Presentation(pptx_path)
    assert summary["slide_count"] == len(deck.slides)
    assert isinstance(summary["slide_count"], int)
    assert summary["slide_count"] >= 1

    _assert_source_artifact_refs(summary["generated_from"], task_path=task_path)
    _assert_source_artifact_refs(summary["source_artifacts"], task_path=task_path)
    assert summary["generated_from"] == summary["source_artifacts"]

    _assert_included_figures(summary["included_figures"], task_path=task_path)
    _assert_included_metrics(summary["included_metrics"], task_path=task_path)

    for key in ["warnings", "figure_warnings", "report_excerpt"]:
        assert isinstance(summary[key], list)
        assert all(isinstance(item, str) and item for item in summary[key])
    assert len(summary["report_excerpt"]) <= MAX_REPORT_BULLETS

    _assert_text_truncations(summary["text_truncations"])
    _assert_table_truncations(summary["table_truncations"])
    _assert_key_comparisons(summary["key_comparisons"])
    _assert_caption_truncations(summary["caption_truncations"])
    _assert_slide_evidence(summary["slide_evidence"], slide_count=summary["slide_count"])
    _assert_claim_traces(summary["claim_traces"], expected_status=expected_status)
    _assert_qa_checks(summary["qa_checks"], slide_count=summary["slide_count"])
    _assert_slides(summary["slides"], slide_count=summary["slide_count"], top_level_sources=summary["source_artifacts"])
    _assert_compat_aliases(summary, task_path=task_path)
    _assert_optional_unstable_keys(summary, task_path=task_path)
    _assert_summary_is_reference_only(summary, task_path=task_path)


def _assert_source_artifact_refs(items: list[str], *, task_path: Path) -> None:
    assert isinstance(items, list)
    assert items
    assert len(items) == len(set(items))
    assert "slides/presentation-draft.pptx" not in items
    assert "slides/slides-summary.json" not in items
    assert not any(item.startswith("slides/") for item in items)
    for item in items:
        assert isinstance(item, str) and item and "\n" not in item
        assert (task_path / item).exists()


def _assert_included_figures(items: list[dict[str, Any]], *, task_path: Path) -> None:
    assert isinstance(items, list)
    assert len(items) <= MAX_FIGURES
    for item in items:
        assert set(item) == REQUIRED_INCLUDED_FIGURE_KEYS
        assert isinstance(item["figure_id"], str) and item["figure_id"]
        assert isinstance(item["chart_type"], str) and item["chart_type"]
        assert isinstance(item["path"], str) and item["path"].endswith(".png")
        assert (task_path / item["path"]).is_file()
        assert isinstance(item["x"], str)
        assert isinstance(item["y"], str)
        assert isinstance(item["group_by"], str)
        assert isinstance(item["source_metrics"], str) and item["source_metrics"]
        assert item["source_metrics"].endswith("/metrics/normalized_metrics.csv") or item["source_metrics"] == "metrics/normalized_metrics.csv"


def _assert_included_metrics(metrics: dict[str, Any], *, task_path: Path) -> None:
    assert set(metrics) == REQUIRED_INCLUDED_METRICS_KEYS
    assert isinstance(metrics["present"], bool)
    assert metrics["path"] == "metrics/normalized_metrics.csv"
    assert isinstance(metrics["row_count"], int)
    assert metrics["row_count"] >= 0
    assert isinstance(metrics["key_columns"], list)
    assert len(metrics["key_columns"]) <= MAX_KEY_COLUMNS
    assert all(isinstance(item, str) and item for item in metrics["key_columns"])
    assert isinstance(metrics["numeric"], list)
    assert len(metrics["numeric"]) <= MAX_NUMERIC_COLUMNS
    for item in metrics["numeric"]:
        assert {"column", "count", "mean", "min", "max"} <= set(item)
        assert isinstance(item["column"], str) and item["column"]
        assert isinstance(item["count"], int)
        assert item["count"] >= 0
        for key in ["mean", "min", "max"]:
            assert item[key] is None or isinstance(item[key], int | float)
        if "final" in item:
            assert item["final"] is None or isinstance(item["final"], int | float)

    if metrics["present"]:
        assert metrics["row_count"] > 0
        assert (task_path / metrics["path"]).is_file()
        if metrics["scenario_summary_path"] is not None:
            assert metrics["scenario_summary_path"] == "metrics/scenario-summary.json"
            assert (task_path / metrics["scenario_summary_path"]).is_file()
    else:
        assert metrics["row_count"] == 0
        assert metrics["key_columns"] == []
        assert metrics["numeric"] == []
        assert metrics["scenario_summary_path"] is None

    _assert_compact_scenario(metrics["scenario"])


def _assert_compact_scenario(scenario: dict[str, Any]) -> None:
    assert isinstance(scenario, dict)
    if not scenario.get("present"):
        assert scenario == {"present": False}
        return

    assert scenario["path"] == "metrics/scenario-summary.json"
    assert isinstance(scenario.get("scenario_type"), str) and scenario["scenario_type"]

    primary_metric = scenario["primary_metric"]
    assert isinstance(primary_metric, dict)
    assert {"name", "direction", "unit", "selection_reason"} <= set(primary_metric)
    assert primary_metric["name"] is None or isinstance(primary_metric["name"], str)
    assert primary_metric["direction"] in {"max", "min", None}
    assert primary_metric["unit"] is None or isinstance(primary_metric["unit"], str | int | float | bool)
    assert primary_metric["selection_reason"] is None or isinstance(primary_metric["selection_reason"], str | int | float | bool)

    groups = scenario["groups"]
    assert isinstance(groups, dict)
    assert {"configured", "primary", "secondary", "seed", "context", "inferred"} <= set(groups)
    assert isinstance(groups["configured"], dict)
    assert groups["primary"] is None or isinstance(groups["primary"], str)
    assert groups["secondary"] is None or isinstance(groups["secondary"], str)
    assert groups["seed"] is None or isinstance(groups["seed"], str)
    assert isinstance(groups["context"], list)
    assert isinstance(groups["inferred"], list)

    assert isinstance(scenario["best_rows"], list)
    assert isinstance(scenario["last_rows"], list)
    assert len(scenario["best_rows"]) <= 3
    assert len(scenario["last_rows"]) <= 3
    for item in scenario["best_rows"]:
        _assert_compact_best_row(item)
    for item in scenario["last_rows"]:
        _assert_compact_last_row(item)

    seed_aggregates = scenario["seed_aggregates"]
    assert isinstance(seed_aggregates, dict)
    assert isinstance(seed_aggregates["present"], bool)
    assert len(seed_aggregates.get("items", [])) <= 3
    if seed_aggregates["present"]:
        assert isinstance(seed_aggregates["metric"], str) and seed_aggregates["metric"]
        assert seed_aggregates["direction"] in {"max", "min"}
        assert isinstance(seed_aggregates["claim_limit"], str) and seed_aggregates["claim_limit"]
        for item in seed_aggregates["items"]:
            _assert_compact_seed_aggregate(item)
    else:
        assert seed_aggregates["items"] == []

    warnings = scenario["warnings"]
    assert isinstance(warnings, list)
    assert len(warnings) <= 5
    assert all(isinstance(item, str) and item for item in warnings)

    omitted = scenario["omitted"]
    assert isinstance(omitted, dict)
    assert all(value == "omitted_by_default" for value in omitted.values())


def _assert_compact_best_row(item: dict[str, Any]) -> None:
    assert {"metric", "direction", "selection_reason", "value", "row_number", "selected_fields", "omitted_field_count", "evidence"} <= set(item)
    assert isinstance(item["metric"], str) and item["metric"]
    assert item["direction"] in {"max", "min"}
    assert item["value"] is None or isinstance(item["value"], int | float)
    assert isinstance(item["row_number"], int) and item["row_number"] >= 1
    assert isinstance(item["selected_fields"], dict)
    assert isinstance(item["omitted_field_count"], int) and item["omitted_field_count"] >= 0
    _assert_row_evidence(item["evidence"], expected_row_number=item["row_number"])


def _assert_compact_last_row(item: dict[str, Any]) -> None:
    assert {"group", "row_number", "checkpoint_field", "selected_fields", "evidence"} <= set(item)
    assert isinstance(item["group"], dict)
    assert isinstance(item["row_number"], int) and item["row_number"] >= 1
    assert item["checkpoint_field"] is None or isinstance(item["checkpoint_field"], str)
    assert isinstance(item["selected_fields"], dict)
    _assert_row_evidence(item["evidence"], expected_row_number=item["row_number"])


def _assert_compact_seed_aggregate(item: dict[str, Any]) -> None:
    assert {"group", "metric", "direction", "row_count", "seed_count", "mean", "min", "max", "evidence"} <= set(item)
    assert isinstance(item["group"], dict)
    assert isinstance(item["metric"], str) and item["metric"]
    assert item["direction"] in {"max", "min"}
    assert isinstance(item["row_count"], int) and item["row_count"] >= 1
    assert isinstance(item["seed_count"], int) and item["seed_count"] >= 1
    for key in ["mean", "min", "max"]:
        assert item[key] is None or isinstance(item[key], int | float)
    evidence = item["evidence"]
    assert evidence["artifact_id"] == "metrics_normalized_csv"
    assert evidence["path"] == "metrics/normalized_metrics.csv"
    assert isinstance(evidence["row_numbers"], list)
    assert len(evidence["row_numbers"]) <= 6
    assert evidence["body"] == "omitted"


def _assert_row_evidence(evidence: dict[str, Any], *, expected_row_number: int) -> None:
    assert evidence == {
        "artifact_id": "metrics_normalized_csv",
        "path": "metrics/normalized_metrics.csv",
        "row_number": expected_row_number,
        "body": "omitted",
    }


def _assert_text_truncations(items: list[dict[str, Any]]) -> None:
    assert isinstance(items, list)
    for item in items:
        assert {"key", "display", "limit", "truncated"} <= set(item)
        assert isinstance(item["key"], str) and item["key"]
        assert isinstance(item["display"], str)
        assert isinstance(item["limit"], int) and item["limit"] >= 1
        assert item["truncated"] is True
        if item["key"] in {"stdout_tail", "stderr_tail"}:
            assert "full" not in item
            assert item["max_lines"] == MAX_LOG_LINES
            assert isinstance(item["full_omitted_reason"], str) and item["full_omitted_reason"]
            assert isinstance(item["omitted_char_count"], int) and item["omitted_char_count"] > 0
            assert isinstance(item["omitted_line_count"], int) and item["omitted_line_count"] > 0


def _assert_table_truncations(items: list[dict[str, Any]]) -> None:
    assert isinstance(items, list)
    for item in items:
        assert set(item) == {
            "source_metrics",
            "shown_columns",
            "hidden_columns",
            "displayed_rows",
            "total_rows",
            "displayed_columns",
            "total_columns",
            "omitted_row_count",
            "omitted_column_count",
            "truncated_cells_count",
        }
        assert item["source_metrics"] == "metrics/normalized_metrics.csv"
        assert isinstance(item["shown_columns"], list)
        assert isinstance(item["hidden_columns"], list)
        assert len(item["shown_columns"]) <= MAX_TABLE_COLUMNS
        assert item["displayed_rows"] <= MAX_TABLE_ROWS
        assert item["displayed_columns"] <= MAX_TABLE_COLUMNS
        for key in [
            "displayed_rows",
            "total_rows",
            "displayed_columns",
            "total_columns",
            "omitted_row_count",
            "omitted_column_count",
            "truncated_cells_count",
        ]:
            assert isinstance(item[key], int)
            assert item[key] >= 0


def _assert_key_comparisons(items: list[dict[str, Any]]) -> None:
    assert isinstance(items, list)
    for item in items:
        assert {"present", "metric", "direction", "group_column", "best_item", "baseline_item", "delta", "top_items", "source_metrics", "reason"} <= set(item)
        assert item["present"] is True
        assert isinstance(item["metric"], str) and item["metric"]
        assert item["direction"] in {"higher", "lower"}
        assert isinstance(item["group_column"], str) and item["group_column"]
        _assert_ranked_comparison_item(item["best_item"])
        if item["baseline_item"] is not None:
            _assert_ranked_comparison_item(item["baseline_item"])
        assert item["delta"] is None or isinstance(item["delta"], int | float)
        assert isinstance(item["top_items"], list)
        assert len(item["top_items"]) <= 3
        for top_item in item["top_items"]:
            _assert_ranked_comparison_item(top_item)
        assert item["source_metrics"] == "metrics/normalized_metrics.csv"
        assert item["reason"] is None or isinstance(item["reason"], str)
        if "display_delta" in item:
            assert item["display_delta"] is None or isinstance(item["display_delta"], str)


def _assert_ranked_comparison_item(item: dict[str, Any]) -> None:
    assert isinstance(item, dict)
    assert set(item) == {"label", "value", "display_value"}
    assert isinstance(item["label"], str) and item["label"]
    assert isinstance(item["value"], int | float)
    assert isinstance(item["display_value"], str) and item["display_value"]


def _assert_caption_truncations(items: list[dict[str, Any]]) -> None:
    assert isinstance(items, list)
    for item in items:
        assert {"figure_id", "display", "limit"} <= set(item)
        assert isinstance(item["figure_id"], str) and item["figure_id"]
        assert isinstance(item["display"], str) and item["display"]
        assert item["limit"] == MAX_CAPTION_CHARS


def _assert_slide_evidence(items: list[dict[str, Any]], *, slide_count: int) -> None:
    assert isinstance(items, list)
    assert len(items) == slide_count
    for item in items:
        assert set(item) == REQUIRED_SLIDE_TRACE_KEYS
        assert isinstance(item["slide_index"], int) and 1 <= item["slide_index"] <= slide_count
        assert isinstance(item["title"], str)
        assert isinstance(item["purpose"], str) and item["purpose"]
        assert isinstance(item["evidence_count"], int) and item["evidence_count"] >= 0
        assert item["empty_source_reason"] is None or isinstance(item["empty_source_reason"], str)


def _assert_claim_traces(items: list[dict[str, Any]], *, expected_status: str) -> None:
    assert isinstance(items, list)
    assert items
    by_id = {item["claim_id"]: item for item in items}
    for item in items:
        assert REQUIRED_CLAIM_TRACE_KEYS <= set(item)
        assert item["surface"] == "slides"
        assert isinstance(item["claim_id"], str) and item["claim_id"]
        assert isinstance(item["claim_type"], str) and item["claim_type"]
        assert isinstance(item["evidence"], list)
        assert item["evidence"]
        for evidence in item["evidence"]:
            _assert_claim_evidence(evidence)

    assert "slides.key_comparison.status" in by_id
    if expected_status == "completed":
        assert "slides.metrics.row_count" in by_id
        assert "slides.metrics.table_preview" in by_id
        assert "slides.diagnostic.failed_status" not in by_id
        assert "slides.diagnostic.cancelled_status" not in by_id
    elif expected_status == "failed":
        assert "slides.metrics.unavailable" in by_id
        assert "slides.diagnostic.failed_status" in by_id
        assert not any(item["claim_type"] == "numeric_summary" for item in items)
    elif expected_status == "cancelled":
        assert "slides.metrics.unavailable" in by_id
        assert "slides.diagnostic.cancelled_status" in by_id
        assert not any(item["claim_type"] == "numeric_summary" for item in items)


def _assert_claim_evidence(evidence: dict[str, Any]) -> None:
    assert isinstance(evidence, dict)
    assert isinstance(evidence.get("artifact_id"), str) and evidence["artifact_id"]
    assert isinstance(evidence.get("path"), str) and evidence["path"]
    assert "rows" not in evidence
    if "body" in evidence:
        assert evidence["body"] == "omitted"


def _assert_qa_checks(qa_checks: dict[str, Any], *, slide_count: int) -> None:
    assert set(qa_checks) == REQUIRED_QA_GROUPS

    slide_count_check = qa_checks["slide_count"]
    assert slide_count_check == {"value": slide_count, "passed": True}

    empty_slide_check = qa_checks["empty_slide_check"]
    assert empty_slide_check["passed"] is True
    assert empty_slide_check["empty_slide_indices"] == []

    title_check = qa_checks["title_check"]
    assert title_check["passed"] is True
    assert title_check["missing_title_slide_indices"] == []

    artifact_duplicate_check = qa_checks["artifact_duplicate_check"]
    assert artifact_duplicate_check["passed"] is True
    assert isinstance(artifact_duplicate_check["artifact_count_after_upsert"], int)
    assert artifact_duplicate_check["artifact_count_after_upsert"] >= 2

    table_guard = qa_checks["table_overflow_guard"]
    assert table_guard["passed"] is True
    assert len(table_guard["shown_columns"]) <= MAX_TABLE_COLUMNS
    assert isinstance(table_guard["hidden_columns"], list)
    assert isinstance(table_guard["truncated_cells_count"], int)
    assert table_guard["truncated_cells_count"] >= 0

    caption_guard = qa_checks["caption_overflow_guard"]
    assert caption_guard["passed"] is True
    assert caption_guard["max_caption_chars"] == MAX_CAPTION_CHARS
    assert isinstance(caption_guard["truncated_caption_count"], int)
    assert caption_guard["truncated_caption_count"] >= 0


def _assert_slides(items: list[dict[str, Any]], *, slide_count: int, top_level_sources: list[str]) -> None:
    assert isinstance(items, list)
    assert len(items) == slide_count
    for item in items:
        assert set(item) == REQUIRED_SLIDE_KEYS
        assert isinstance(item["slide_index"], int) and 1 <= item["slide_index"] <= slide_count
        assert isinstance(item["title"], str)
        assert isinstance(item["purpose"], str) and item["purpose"]
        assert isinstance(item["source_artifacts"], list)
        assert all(isinstance(source, str) and source for source in item["source_artifacts"])
        assert set(item["source_artifacts"]).issubset(set(top_level_sources))
        assert isinstance(item["evidence"], list)
        assert item["evidence"] or item["empty_source_reason"]
        assert item["empty_source_reason"] is None or isinstance(item["empty_source_reason"], str)
        for evidence in item["evidence"]:
            assert REQUIRED_SLIDE_EVIDENCE_KEYS <= set(evidence)
            assert isinstance(evidence["path"], str) and evidence["path"]
            assert isinstance(evidence["artifact_id"], str) and evidence["artifact_id"]
            assert isinstance(evidence["role"], str) and evidence["role"]
            if "body" in evidence:
                assert evidence["body"] == "omitted"


def _assert_compat_aliases(summary: dict[str, Any], *, task_path: Path) -> None:
    slide_titles = summary["slide_titles"]
    assert isinstance(slide_titles, list)
    assert slide_titles == [slide["title"] for slide in summary["slides"]]

    figures = summary["figures"]
    assert isinstance(figures, list)
    assert figures == [item["path"] for item in summary["included_figures"]]
    for path in figures:
        assert isinstance(path, str) and path.endswith(".png")
        assert (task_path / path).is_file()

    metrics = summary["metrics"]
    assert isinstance(metrics, dict)
    assert {"present", "path", "scenario_summary_path", "scenario", "row_count", "columns", "key_columns", "omitted_key_column_count", "numeric", "numeric_omitted_count"} <= set(metrics)
    assert metrics["path"] == "metrics/normalized_metrics.csv"
    assert isinstance(metrics["columns"], list)
    assert isinstance(metrics["omitted_key_column_count"], int)
    assert metrics["omitted_key_column_count"] >= 0
    assert isinstance(metrics["numeric_omitted_count"], int)
    assert metrics["numeric_omitted_count"] >= 0

    metrics_table = summary["metrics_table"]
    assert isinstance(metrics_table, dict)
    assert {
        "present",
        "source_metrics",
        "columns",
        "shown_columns",
        "hidden_columns",
        "rows",
        "displayed_row_count",
        "total_row_count",
        "total_column_count",
        "truncated",
        "omitted_row_count",
        "omitted_column_count",
        "truncated_cells_count",
    } <= set(metrics_table)
    assert metrics_table["source_metrics"] == "metrics/normalized_metrics.csv"
    assert len(metrics_table["shown_columns"]) <= MAX_TABLE_COLUMNS
    assert metrics_table["displayed_row_count"] <= MAX_TABLE_ROWS
    assert isinstance(metrics_table["rows"], list)
    assert len(metrics_table["rows"]) == metrics_table["displayed_row_count"]


def _assert_optional_unstable_keys(summary: dict[str, Any], *, task_path: Path) -> None:
    assert set(summary) >= OPTIONAL_UNSTABLE_KEYS

    project_goal = summary["project_goal"]
    assert isinstance(project_goal, dict)
    assert {"present", "path", "excerpt"} <= set(project_goal)
    assert isinstance(project_goal["present"], bool)
    if project_goal["present"]:
        assert isinstance(project_goal["path"], str) and project_goal["path"]
        assert (task_path / project_goal["path"]).exists() or (Path(project_goal["path"]).is_absolute() and Path(project_goal["path"]).exists())
        assert isinstance(project_goal["excerpt"], str) and project_goal["excerpt"]
    else:
        assert project_goal["path"] is None
        assert isinstance(project_goal["excerpt"], str) and project_goal["excerpt"]

    full_text_fields = summary["full_text_fields"]
    assert isinstance(full_text_fields, dict)
    assert {"command", "source_path", "working_dir", "failure_summary", "artifact_dir"} <= set(full_text_fields)


def _assert_summary_is_reference_only(summary: dict[str, Any], *, task_path: Path) -> None:
    serialized = json.dumps(summary, ensure_ascii=False)
    for fragment in FORBIDDEN_EMBEDDED_FRAGMENTS:
        assert fragment not in serialized

    for relative_path in [
        "stdout.log",
        "stderr.log",
        "reports/report-fragment.md",
        "metrics/normalized_metrics.csv",
    ]:
        path = task_path / relative_path
        if path.exists():
            artifact_text = path.read_text(encoding="utf-8", errors="replace")
            if artifact_text:
                assert artifact_text not in serialized


def _write_wide_metrics_csv(path: Path) -> None:
    headers = ["model", "variant", "accuracy", "f1", "latency_ms", "loss", "memory_mb", "runtime_s", "extra_a", "extra_b"]
    rows = []
    for index in range(12):
        rows.append(
            {
                "model": f"model_{index}",
                "variant": f"v{index}",
                "accuracy": str(0.70 + index / 100),
                "f1": str(0.60 + index / 100),
                "latency_ms": str(30 + index),
                "loss": str(1.0 - index / 100),
                "memory_mb": str(100 + index),
                "runtime_s": str(5 + index),
                "extra_a": str(index),
                "extra_b": str(index * 2),
            }
        )
    with path.open("w", newline="", encoding="utf-8") as fh:
        writer = csv.DictWriter(fh, fieldnames=headers)
        writer.writeheader()
        writer.writerows(rows)


def test_slides_summary_completed_after_collect_figures_report_matches_public_contract(tmp_path: Path) -> None:
    workspace = tmp_path / "slides-summary-contract-completed"
    workspace.mkdir()
    copy_examples(workspace)
    assert invoke(workspace, ["init"]).exit_code == 0

    task_id = extract_task_id(invoke(workspace, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(workspace, ["collect", task_id]).exit_code == 0
    assert invoke(workspace, ["figures", task_id]).exit_code == 0
    assert invoke(workspace, ["report", task_id]).exit_code == 0
    assert invoke(workspace, ["slides", task_id]).exit_code == 0

    summary = read_slides_summary(workspace, task_id)
    assert_slides_summary_contract(
        summary,
        workspace=workspace,
        expected_task_id=task_id,
        expected_status="completed",
    )
    assert summary["included_figures"]
    assert summary["included_metrics"]["present"] is True
    assert any(item["claim_id"] == "slides.metrics.scenario_summary" for item in summary["claim_traces"])
    assert "reports/report-fragment.md" in summary["generated_from"]
    assert "figures/figure-summary.json" in summary["generated_from"]


def test_slides_summary_without_figures_but_with_report_and_metrics_stays_bounded(tmp_path: Path) -> None:
    workspace = tmp_path / "slides-summary-contract-no-figures"
    workspace.mkdir()
    _write_wide_metrics_csv(workspace / "wide_metrics.csv")
    assert invoke(workspace, ["init"]).exit_code == 0

    task_id = extract_task_id(invoke(workspace, ["ingest", "wide_metrics.csv"]).output)
    assert invoke(workspace, ["collect", task_id]).exit_code == 0
    assert invoke(workspace, ["report", task_id]).exit_code == 0
    assert invoke(workspace, ["slides", task_id]).exit_code == 0

    summary = read_slides_summary(workspace, task_id)
    assert_slides_summary_contract(
        summary,
        workspace=workspace,
        expected_task_id=task_id,
        expected_status="completed",
    )
    assert summary["included_figures"] == []
    assert summary["figures"] == []
    assert summary["metrics"]["numeric_omitted_count"] > 0
    assert summary["metrics"]["omitted_key_column_count"] > 0
    assert summary["table_truncations"]
    assert summary["qa_checks"]["table_overflow_guard"]["passed"] is True
    assert any("omitted" in warning for warning in summary["warnings"])


def test_slides_summary_failed_diagnostic_omits_full_logs_reports_and_pptx_bytes(tmp_path: Path) -> None:
    workspace = tmp_path / "slides-summary-contract-failed"
    workspace.mkdir()
    copy_examples(workspace)
    assert invoke(workspace, ["init"]).exit_code == 0

    command = f'"{sys.executable}" examples/simple-failure/fail.py'
    task_id = extract_task_id(invoke(workspace, ["run", command]).output)
    assert invoke(workspace, ["report", task_id]).exit_code == 0
    assert invoke(workspace, ["slides", task_id]).exit_code == 0

    summary = read_slides_summary(workspace, task_id)
    assert_slides_summary_contract(
        summary,
        workspace=workspace,
        expected_task_id=task_id,
        expected_status="failed",
    )
    assert summary["included_metrics"]["present"] is False
    assert "reports/report-fragment.md" in summary["generated_from"]
    assert "slides.diagnostic.failed_status" in {item["claim_id"] for item in summary["claim_traces"]}


def test_slides_summary_cancelled_diagnostic_omits_full_logs_reports_and_pptx_bytes(tmp_path: Path) -> None:
    workspace = tmp_path / "slides-summary-contract-cancelled"
    workspace.mkdir()
    script = workspace / "cancel_me.py"
    script.write_text(
        "\n".join(
            [
                "import sys",
                "import time",
                "print('slides-summary-contract-ready', flush=True)",
                "print('slides-summary-contract-stderr', file=sys.stderr, flush=True)",
                "time.sleep(30)",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    assert invoke(workspace, ["init"]).exit_code == 0

    command = f'"{sys.executable}" cancel_me.py'
    task_id = extract_task_id(invoke(workspace, ["run", command, "--background"]).output)
    wait_for_output(workspace, task_id, "slides-summary-contract-ready")
    assert invoke(workspace, ["cancel", task_id]).exit_code == 0
    assert invoke(workspace, ["report", task_id]).exit_code == 0
    assert invoke(workspace, ["slides", task_id]).exit_code == 0

    summary = read_slides_summary(workspace, task_id)
    assert_slides_summary_contract(
        summary,
        workspace=workspace,
        expected_task_id=task_id,
        expected_status="cancelled",
    )
    assert summary["included_metrics"]["present"] is False
    assert "reports/report-fragment.md" in summary["generated_from"]
    assert "slides.diagnostic.cancelled_status" in {item["claim_id"] for item in summary["claim_traces"]}
