from __future__ import annotations

import csv
import json
import os
import shutil
import sys
import time
from pathlib import Path

import yaml
from PIL import Image
from pptx import Presentation
from typer.testing import CliRunner

from lab_sidecar.cli.app import app
from lab_sidecar.core.provenance import git_snapshot
from lab_sidecar.runner.process import terminate_process_tree


runner = CliRunner()
PROJECT_ROOT = Path(__file__).resolve().parents[1]


def copy_examples(workspace: Path) -> None:
    shutil.copytree(PROJECT_ROOT / "examples", workspace / "examples")


def invoke(workspace: Path, args: list[str]):
    old_cwd = Path.cwd()
    os.chdir(workspace)
    try:
        return runner.invoke(app, args, prog_name="labsidecar", env={}, catch_exceptions=False)
    finally:
        os.chdir(old_cwd)


def extract_task_id(output: str) -> str:
    for line in output.splitlines():
        if line.startswith("Task created: "):
            return line.split(": ", 1)[1]
        if line.startswith("Imported as task: "):
            return line.split(": ", 1)[1]
    raise AssertionError(f"No task id in output:\n{output}")


def read_manifest(workspace: Path, task_id: str) -> dict:
    path = workspace / ".lab-sidecar" / "tasks" / task_id / "manifest.json"
    return json.loads(path.read_text(encoding="utf-8"))


def assert_manifest_json_valid(workspace: Path, task_id: str) -> dict:
    path = workspace / ".lab-sidecar" / "tasks" / task_id / "manifest.json"
    raw = path.read_text(encoding="utf-8")
    parsed = json.loads(raw)
    assert parsed["task_id"] == task_id
    return parsed


def read_csv_rows(path: Path) -> list[dict[str, str]]:
    with path.open("r", newline="", encoding="utf-8") as fh:
        return list(csv.DictReader(fh))


def assert_non_empty_file(path: Path) -> None:
    assert path.is_file()
    assert path.stat().st_size > 0


def assert_nonblank_png(path: Path, min_width: int = 320, min_height: int = 180) -> None:
    assert_non_empty_file(path)
    with Image.open(path) as image:
        rgb = image.convert("RGB")
        assert rgb.size[0] >= min_width
        assert rgb.size[1] >= min_height
        colors = rgb.getcolors(maxcolors=1_000_000)
        assert colors is not None
        assert len(colors) > 1


def write_stage3_messy_results(workspace: Path) -> Path:
    source = workspace / "messy-results"
    (source / "baseline" / "seed_1").mkdir(parents=True)
    (source / "baseline" / "seed_2").mkdir(parents=True)
    (source / "baseline" / "seed_3").mkdir(parents=True)
    (source / "candidate" / "seed_1").mkdir(parents=True)
    (source / "candidate" / "seed_2").mkdir(parents=True)
    (source / "candidate" / "seed_3").mkdir(parents=True)
    (source / "debug").mkdir(parents=True)
    (source / "scratch").mkdir(parents=True)

    for method, start in [("baseline", 70), ("candidate", 78)]:
        for seed in [1, 2, 3]:
            (source / method / f"seed_{seed}" / "metrics.csv").write_text(
                "\n".join(
                    [
                        "iter,algo,trial,score_pct,runtime_ms",
                        f"1,{method},{seed},0.{start + seed},5{seed}",
                        f"2,{method},{seed},0.{start + seed + 4},4{seed}",
                    ]
                )
                + "\n",
                encoding="utf-8",
            )

    (source / "debug" / "debug_metrics.csv").write_text(
        "iter,algo,trial,score_pct,runtime_ms\n1,debug,0,0.01,999\n",
        encoding="utf-8",
    )
    (source / "scratch" / "scratch.csv").write_text(
        "iter,algo,trial,score_pct,runtime_ms\n1,scratch,0,0.02,999\n",
        encoding="utf-8",
    )
    return source


def assert_readable_deck(path: Path, min_slides: int = 5, max_slides: int = 8) -> Presentation:
    assert_non_empty_file(path)
    deck = Presentation(path)
    assert min_slides <= len(deck.slides) <= max_slides
    for index, slide in enumerate(deck.slides, start=1):
        text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        assert text.strip(), f"slide {index} appears blank"
        assert any(shape.has_text_frame and shape.text.strip() for shape in slide.shapes), f"slide {index} has no title/text"
    return deck


def read_traceability(workspace: Path, task_id: str) -> dict:
    path = workspace / ".lab-sidecar" / "tasks" / task_id / "provenance" / "traceability.json"
    assert_non_empty_file(path)
    return json.loads(path.read_text(encoding="utf-8"))


def assert_traceability_is_bounded(traceability: dict) -> None:
    serialized = json.dumps(traceability, ensure_ascii=False)
    forbidden_fragments = [
        "worker transcript body\n",
        '"prompt":"secret"',
        '"response":"secret"',
        "FileNotFoundError: simulated missing dataset file",
        "epoch,train_loss,val_loss",
        "ppt/presentation.xml",
        "<p:sld",
    ]
    for fragment in forbidden_fragments:
        assert fragment not in serialized
    for trace in traceability.get("claim_traces", []):
        for evidence in trace.get("evidence", []):
            assert "rows" not in evidence
            assert evidence.get("body") in {None, "omitted"}


def test_git_snapshot_records_repository_commit_and_status() -> None:
    snapshot = git_snapshot(PROJECT_ROOT)

    assert snapshot["is_repository"] is True
    assert snapshot["commit"]
    assert isinstance(snapshot["status_short"], list)
    assert "dirty" in snapshot


def wait_for_output(workspace: Path, task_id: str, needle: str, timeout: float = 10.0) -> None:
    path = workspace / ".lab-sidecar" / "tasks" / task_id / "stdout.log"
    deadline = time.time() + timeout
    while time.time() < deadline:
        if path.exists() and needle in path.read_text(encoding="utf-8", errors="replace"):
            return
        time.sleep(0.1)
    raise AssertionError(f"Timed out waiting for {needle!r} in {path}")


def test_init_creates_workspace(tmp_path: Path) -> None:
    result = invoke(tmp_path, ["init"])

    assert result.exit_code == 0
    assert (tmp_path / ".lab-sidecar").is_dir()
    assert (tmp_path / ".lab-sidecar" / "config.yaml").is_file()
    assert (tmp_path / ".lab-sidecar" / "tasks").is_dir()
    assert (tmp_path / ".lab-sidecar" / "index.sqlite").is_file()
    assert "Next:" in result.output
    assert 'labsidecar run "<command>"' in result.output
    assert "labsidecar ingest <path>" in result.output

    second = invoke(tmp_path, ["init"])
    assert second.exit_code == 2

    forced = invoke(tmp_path, ["init", "--force"])
    assert forced.exit_code == 0


def test_doctor_reports_workspace_health(tmp_path: Path) -> None:
    before_init = invoke(tmp_path, ["doctor"])

    assert before_init.exit_code == 0
    assert "Lab-Sidecar doctor" in before_init.output
    assert "[ok] Python:" in before_init.output
    assert "[ok] Writable workspace" in before_init.output
    assert "[warn] Workspace config: not initialized" in before_init.output
    assert "labsidecar init" in before_init.output

    assert invoke(tmp_path, ["init"]).exit_code == 0
    after_init = invoke(tmp_path, ["doctor"])

    assert after_init.exit_code == 0
    assert "[ok] Config: .lab-sidecar/config.yaml" in after_init.output
    assert "[ok] Task directory: .lab-sidecar/tasks" in after_init.output


def test_simple_success_task_and_queries(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0

    command = f'"{sys.executable}" examples/simple-success/train.py --output metrics.csv'
    result = invoke(tmp_path, ["run", command])

    assert result.exit_code == 0
    task_id = extract_task_id(result.output)
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    manifest = read_manifest(tmp_path, task_id)

    assert f"Artifacts: .lab-sidecar/tasks/{task_id}" in result.output
    assert f"Stdout: .lab-sidecar/tasks/{task_id}/stdout.log" in result.output
    assert f"Stderr: .lab-sidecar/tasks/{task_id}/stderr.log" in result.output
    assert "Next:" in result.output
    assert f"labsidecar collect {task_id}" in result.output
    assert f"labsidecar artifacts {task_id}" in result.output

    assert manifest["schema_version"] == "1"
    assert manifest["task_id"] == task_id
    assert manifest["mode"] == "run"
    assert manifest["status"] == "completed"
    assert manifest["working_dir"] == "."
    assert manifest["command"] == command
    assert manifest["run_mode"] == "shell"
    assert manifest["argv"] is None
    assert manifest["safe_profile"] is None
    assert manifest["source_path"] is None
    assert manifest["exit_code"] == 0
    assert set(["task_dir", "stdout", "stderr"]).issubset(manifest["paths"])
    assert manifest["artifacts"]
    assert (task_path / "manifest.json").is_file()
    assert (task_path / "stdout.log").is_file()
    assert (task_path / "stderr.log").is_file()
    assert (task_path / "reproduce" / "command.txt").is_file()
    assert (task_path / "reproduce" / "run.json").is_file()
    assert (task_path / "reproduce" / "env.json").is_file()
    assert (task_path / "reproduce" / "git.json").is_file()
    assert (task_path / "reproduce" / "dependencies.json").is_file()
    env_snapshot = json.loads((task_path / "reproduce" / "env.json").read_text(encoding="utf-8"))
    run_snapshot = json.loads((task_path / "reproduce" / "run.json").read_text(encoding="utf-8"))
    git_snapshot = json.loads((task_path / "reproduce" / "git.json").read_text(encoding="utf-8"))
    dependencies = json.loads((task_path / "reproduce" / "dependencies.json").read_text(encoding="utf-8"))
    assert run_snapshot["run_mode"] == "shell"
    assert run_snapshot["command_text"] == command
    assert run_snapshot["argv"] is None
    assert env_snapshot["python_version"]
    assert env_snapshot["python_executable"]
    assert env_snapshot["platform"]
    assert env_snapshot["working_dir"] == str(tmp_path)
    assert env_snapshot["git_path"] == "reproduce/git.json"
    assert env_snapshot["dependencies_path"] == "reproduce/dependencies.json"
    assert "is_repository" in git_snapshot
    assert isinstance(dependencies, dict)

    status = invoke(tmp_path, ["status", task_id])
    assert status.exit_code == 0
    assert "Status: completed" in status.output
    assert "Run mode: shell" in status.output
    assert "Exit code: 0" in status.output
    assert f"Artifact dir: .lab-sidecar/tasks/{task_id}" in status.output

    logs = invoke(tmp_path, ["logs", task_id, "--tail", "20"])
    assert logs.exit_code == 0
    assert "Best val_accuracy=0.86" in logs.output
    assert "== stderr" in logs.output

    artifacts = invoke(tmp_path, ["artifacts", task_id])
    assert artifacts.exit_code == 0
    assert "[log]" in artifacts.output
    assert "stdout.log" in artifacts.output
    assert "reproduce/run.json" in artifacts.output
    assert "reproduce/git.json" in artifacts.output
    assert "reproduce/dependencies.json" in artifacts.output

    (tmp_path / ".lab-sidecar" / "index.sqlite").unlink()
    assert invoke(tmp_path, ["status", task_id]).exit_code == 0
    assert invoke(tmp_path, ["logs", task_id, "--stream", "stdout", "--tail", "5"]).exit_code == 0
    assert invoke(tmp_path, ["artifacts", task_id]).exit_code == 0


def test_simple_failure_task(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0

    command = f'"{sys.executable}" examples/simple-failure/fail.py'
    result = invoke(tmp_path, ["run", command])

    assert result.exit_code == 0
    task_id = extract_task_id(result.output)
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    manifest = read_manifest(tmp_path, task_id)

    assert manifest["status"] == "failed"
    assert manifest["exit_code"] != 0
    assert manifest["failure_summary"]
    assert "FileNotFoundError" in manifest["failure_summary"]
    assert (task_path / "stderr.log").is_file()
    assert "FileNotFoundError" in (task_path / "stderr.log").read_text(encoding="utf-8")

    status = invoke(tmp_path, ["status", task_id])
    assert status.exit_code == 0
    assert "Status: failed" in status.output
    assert "Failure summary:" in status.output

    logs = invoke(tmp_path, ["logs", task_id, "--stream", "stderr", "--tail", "20"])
    assert logs.exit_code == 0
    assert "FileNotFoundError" in logs.output


def test_no_shell_argv_foreground_preserves_literal_arguments_and_metadata(tmp_path: Path) -> None:
    script = tmp_path / "argv_probe.py"
    script.write_text(
        "\n".join(
            [
                "import json",
                "import sys",
                "from pathlib import Path",
                "Path('argv-seen.json').write_text(json.dumps(sys.argv[1:]), encoding='utf-8')",
                "Path('metrics.csv').write_text('epoch,val_accuracy\\n1,0.91\\n', encoding='utf-8')",
                "print('ARGV_PROBE_OK', flush=True)",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    assert invoke(tmp_path, ["init"]).exit_code == 0

    argv = [
        sys.executable,
        "argv_probe.py",
        "value with spaces",
        "SECRET_LOOKING_ARG_" + ("x" * 120),
        "&&",
        "touch",
        "argv-sentinel.txt",
        "*.csv",
        "$HOME",
        "--child-opt",
        "child value",
    ]
    result = invoke(tmp_path, ["run", "--no-shell", "--", *argv])

    assert result.exit_code == 0
    assert "Run mode: argv" in result.output
    task_id = extract_task_id(result.output)
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    manifest = read_manifest(tmp_path, task_id)
    run_snapshot = json.loads((task_path / "reproduce" / "run.json").read_text(encoding="utf-8"))

    assert manifest["status"] == "completed"
    assert manifest["run_mode"] == "argv"
    assert manifest["argv"] == argv
    assert manifest["safe_profile"] is None
    assert run_snapshot["run_mode"] == "argv"
    assert run_snapshot["argv"] == argv
    assert "shell=False" in run_snapshot["execution_note"]
    assert not (tmp_path / "argv-sentinel.txt").exists()
    assert json.loads((tmp_path / "argv-seen.json").read_text(encoding="utf-8")) == argv[2:]

    status = invoke(tmp_path, ["status", task_id])
    assert status.exit_code == 0
    assert "Run mode: argv" in status.output
    artifacts = invoke(tmp_path, ["artifacts", task_id])
    assert artifacts.exit_code == 0
    assert "reproduce/run.json" in artifacts.output
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    package_path = tmp_path / f"lab-sidecar-package-{task_id}"
    assert invoke(tmp_path, ["package", task_id, "--output", package_path.as_posix()]).exit_code == 0

    package_summary = json.loads((package_path / "package-summary.json").read_text(encoding="utf-8"))
    package_index = json.loads((package_path / "artifact-index.json").read_text(encoding="utf-8"))
    package_traceability = json.loads((package_path / "provenance" / "traceability.json").read_text(encoding="utf-8"))
    assert package_summary["task"]["run_mode"] == "argv"
    assert package_summary["task"]["argv_count"] == len(argv)
    assert package_summary["task"]["run_spec_path"] == "reproduce/run.json"
    assert "SECRET_LOOKING_ARG_" not in json.dumps(package_summary, ensure_ascii=False)
    assert "reproduce/run.json" in {item["package_path"] for item in package_index["included"]}
    assert package_traceability["task"]["run_mode"] == "argv"
    assert package_traceability["task"]["argv"] == argv
    assert package_traceability["task"]["run_spec_path"] == "reproduce/run.json"
    assert_traceability_is_bounded(package_traceability)


def test_no_shell_argv_foreground_failure_records_failed_task(tmp_path: Path) -> None:
    script = tmp_path / "argv_failure.py"
    script.write_text(
        "\n".join(
            [
                "import sys",
                "print('argv-fail-start', file=sys.stderr, flush=True)",
                "raise SystemExit(7)",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    assert invoke(tmp_path, ["init"]).exit_code == 0

    argv = [sys.executable, "argv_failure.py"]
    result = invoke(tmp_path, ["run", "--no-shell", "--", *argv])

    assert result.exit_code == 0
    task_id = extract_task_id(result.output)
    manifest = read_manifest(tmp_path, task_id)
    assert manifest["status"] == "failed"
    assert manifest["exit_code"] == 7
    assert manifest["run_mode"] == "argv"
    assert manifest["argv"] == argv
    assert "argv-fail-start" in manifest["failure_summary"]


def test_no_shell_missing_executable_records_clear_failed_task(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0

    missing_executable = "lab-sidecar-missing-executable-v014"
    result = invoke(tmp_path, ["run", "--no-shell", "--", missing_executable, "--version"])

    assert result.exit_code == 0
    task_id = extract_task_id(result.output)
    manifest = read_manifest(tmp_path, task_id)
    stderr_text = (tmp_path / ".lab-sidecar" / "tasks" / task_id / "stderr.log").read_text(encoding="utf-8")
    assert manifest["status"] == "failed"
    assert manifest["run_mode"] == "argv"
    assert manifest["argv"] == [missing_executable, "--version"]
    assert manifest["exit_code"] is None
    assert "command could not be started" in stderr_text
    assert missing_executable in stderr_text


def test_no_shell_child_flags_are_not_parsed_as_sidecar_options(tmp_path: Path) -> None:
    script = tmp_path / "argv_child_flags.py"
    script.write_text(
        "\n".join(
            [
                "import json",
                "import sys",
                "from pathlib import Path",
                "Path('child-flags.json').write_text(json.dumps(sys.argv[1:]), encoding='utf-8')",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    assert invoke(tmp_path, ["init"]).exit_code == 0

    result = invoke(tmp_path, ["run", "--no-shell", sys.executable, "argv_child_flags.py", "--name", "child"])

    assert result.exit_code == 0
    task_id = extract_task_id(result.output)
    manifest = read_manifest(tmp_path, task_id)
    assert manifest["name"] is None
    assert manifest["argv"] == [sys.executable, "argv_child_flags.py", "--name", "child"]
    assert json.loads((tmp_path / "child-flags.json").read_text(encoding="utf-8")) == ["--name", "child"]


def test_no_shell_argv_background_success_and_cancel(tmp_path: Path) -> None:
    success_script = tmp_path / "argv_background_success.py"
    success_script.write_text(
        "\n".join(
            [
                "import time",
                "print('argv-background-start', flush=True)",
                "time.sleep(0.2)",
                "print('argv-background-done', flush=True)",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    cancel_script = tmp_path / "argv_background_cancel.py"
    cancel_script.write_text(
        "\n".join(
            [
                "import time",
                "print('argv-cancel-ready', flush=True)",
                "time.sleep(30)",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    assert invoke(tmp_path, ["init"]).exit_code == 0

    success_argv = [sys.executable, "argv_background_success.py"]
    success_task = extract_task_id(
        invoke(tmp_path, ["run", "--background", "--no-shell", "--", *success_argv]).output
    )
    deadline = time.time() + 10
    status = None
    while time.time() < deadline:
        status = invoke(tmp_path, ["status", success_task])
        assert status.exit_code == 0
        if "Status: completed" in status.output:
            break
        time.sleep(0.2)
    assert status is not None
    assert "Status: completed" in status.output
    success_manifest = read_manifest(tmp_path, success_task)
    assert success_manifest["run_mode"] == "argv"
    assert success_manifest["argv"] == success_argv
    assert success_manifest["worker_pid"] is None

    cancel_argv = [sys.executable, "argv_background_cancel.py"]
    cancel_task = extract_task_id(
        invoke(tmp_path, ["run", "--background", "--no-shell", "--", *cancel_argv]).output
    )
    wait_for_output(tmp_path, cancel_task, "argv-cancel-ready")
    cancel = invoke(tmp_path, ["cancel", cancel_task])
    assert cancel.exit_code == 0
    cancel_manifest = read_manifest(tmp_path, cancel_task)
    assert cancel_manifest["status"] == "cancelled"
    assert cancel_manifest["run_mode"] == "argv"
    assert cancel_manifest["argv"] == cancel_argv
    assert cancel_manifest["pid"] is None
    assert cancel_manifest["worker_pid"] is None


def test_no_shell_without_argv_fails_before_task_creation(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0

    result = invoke(tmp_path, ["run", "--no-shell", "--"])

    assert result.exit_code == 2
    assert "--no-shell requires argv" in result.output
    assert list((tmp_path / ".lab-sidecar" / "tasks").iterdir()) == []


def test_each_task_has_independent_directory(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0

    command = f'"{sys.executable}" examples/simple-success/train.py --output metrics.csv'
    first = extract_task_id(invoke(tmp_path, ["run", command]).output)
    second = extract_task_id(invoke(tmp_path, ["run", command]).output)

    assert first != second
    assert (tmp_path / ".lab-sidecar" / "tasks" / first).is_dir()
    assert (tmp_path / ".lab-sidecar" / "tasks" / second).is_dir()


def test_background_run_status_logs_and_cancel(tmp_path: Path) -> None:
    script = tmp_path / "long_task.py"
    script.write_text(
        "\n".join(
            [
                "import time",
                "for i in range(30):",
                "    print(f'tick={i}', flush=True)",
                "    time.sleep(0.5)",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    assert invoke(tmp_path, ["init"]).exit_code == 0

    command = f'"{sys.executable}" long_task.py'
    result = invoke(tmp_path, ["run", command, "--background"])

    assert result.exit_code == 0
    task_id = extract_task_id(result.output)
    manifest = assert_manifest_json_valid(tmp_path, task_id)
    assert manifest["status"] == "running"
    assert manifest["worker_pid"]

    wait_for_output(tmp_path, task_id, "tick=0")

    status = invoke(tmp_path, ["status", task_id])
    assert status.exit_code == 0
    assert "Status: running" in status.output
    assert_manifest_json_valid(tmp_path, task_id)

    logs = invoke(tmp_path, ["logs", task_id, "--stream", "stdout", "--tail", "5"])
    assert logs.exit_code == 0
    assert "tick=0" in logs.output
    assert_manifest_json_valid(tmp_path, task_id)

    cancel = invoke(tmp_path, ["cancel", task_id])
    assert cancel.exit_code == 0
    assert "Status: cancelled" in cancel.output
    assert_manifest_json_valid(tmp_path, task_id)

    status_after_cancel = invoke(tmp_path, ["status", task_id])
    assert status_after_cancel.exit_code == 0
    assert "Status: cancelled" in status_after_cancel.output
    cancelled_manifest = assert_manifest_json_valid(tmp_path, task_id)
    assert cancelled_manifest["status"] == "cancelled"
    assert cancelled_manifest["pid"] is None
    assert cancelled_manifest["worker_pid"] is None


def test_background_completed_task_refreshes_to_completed(tmp_path: Path) -> None:
    script = tmp_path / "quick_success.py"
    script.write_text(
        "\n".join(
            [
                "import time",
                "print('quick-start', flush=True)",
                "time.sleep(0.2)",
                "print('quick-done', flush=True)",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    assert invoke(tmp_path, ["init"]).exit_code == 0

    task_id = extract_task_id(invoke(tmp_path, ["run", f'"{sys.executable}" quick_success.py', "--background"]).output)

    deadline = time.time() + 10
    status = None
    while time.time() < deadline:
        status = invoke(tmp_path, ["status", task_id])
        assert status.exit_code == 0
        if "Status: completed" in status.output:
            break
        time.sleep(0.2)

    assert status is not None
    assert "Status: completed" in status.output
    manifest = read_manifest(tmp_path, task_id)
    assert manifest["status"] == "completed"
    assert manifest["exit_code"] == 0
    assert manifest["pid"] is None
    assert manifest["worker_pid"] is None
    logs = invoke(tmp_path, ["logs", task_id, "--stream", "stdout", "--tail", "5"])
    assert logs.exit_code == 0
    assert "quick-done" in logs.output


def test_background_failed_task_refreshes_to_failed(tmp_path: Path) -> None:
    script = tmp_path / "quick_failure.py"
    script.write_text(
        "\n".join(
            [
                "import sys",
                "print('about to fail', file=sys.stderr, flush=True)",
                "raise RuntimeError('phase2 background failure')",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    assert invoke(tmp_path, ["init"]).exit_code == 0

    task_id = extract_task_id(invoke(tmp_path, ["run", f'"{sys.executable}" quick_failure.py', "--background"]).output)

    deadline = time.time() + 10
    status = None
    while time.time() < deadline:
        status = invoke(tmp_path, ["status", task_id])
        assert status.exit_code == 0
        if "Status: failed" in status.output:
            break
        time.sleep(0.2)

    assert status is not None
    assert "Status: failed" in status.output
    assert "phase2 background failure" in status.output
    manifest = read_manifest(tmp_path, task_id)
    assert manifest["status"] == "failed"
    assert manifest["exit_code"] != 0
    assert manifest["pid"] is None
    assert manifest["worker_pid"] is None
    assert "phase2 background failure" in manifest["failure_summary"]


def test_background_completed_task_after_deleting_sqlite_uses_manifest(tmp_path: Path) -> None:
    script = tmp_path / "quick_success.py"
    script.write_text("print('sqlite-free-complete', flush=True)\n", encoding="utf-8")
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["run", f'"{sys.executable}" quick_success.py', "--background"]).output)

    deadline = time.time() + 10
    while time.time() < deadline:
        if read_manifest(tmp_path, task_id)["status"] == "completed":
            break
        time.sleep(0.2)

    index_path = tmp_path / ".lab-sidecar" / "index.sqlite"
    index_path.unlink()
    status = invoke(tmp_path, ["status", task_id])

    assert status.exit_code == 0
    assert "Status: completed" in status.output
    assert invoke(tmp_path, ["logs", task_id, "--tail", "5"]).exit_code == 0
    assert invoke(tmp_path, ["artifacts", task_id]).exit_code == 0
    assert index_path.is_file()


def test_stale_background_worker_is_marked_failed_with_diagnostics(tmp_path: Path) -> None:
    script = tmp_path / "long_task.py"
    script.write_text(
        "\n".join(
            [
                "import time",
                "print('stale-ready', flush=True)",
                "time.sleep(30)",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["run", f'"{sys.executable}" long_task.py', "--background"]).output)
    wait_for_output(tmp_path, task_id, "stale-ready")

    manifest = read_manifest(tmp_path, task_id)
    worker_pid = manifest["worker_pid"]
    assert worker_pid
    terminate_process_tree(worker_pid)
    child_pid = manifest["pid"]
    if child_pid:
        terminate_process_tree(child_pid)

    manifest_path = tmp_path / ".lab-sidecar" / "tasks" / task_id / "manifest.json"
    manifest["pid"] = None
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    deadline = time.time() + 10
    status = None
    while time.time() < deadline:
        status = invoke(tmp_path, ["status", task_id])
        assert status.exit_code == 0
        if "Status: failed" in status.output:
            break
        time.sleep(0.2)

    assert status is not None
    assert "Status: failed" in status.output
    manifest = read_manifest(tmp_path, task_id)
    assert manifest["status"] == "failed"
    assert manifest["pid"] is None
    assert manifest["worker_pid"] is None
    assert "background task recovery marked this task failed" in manifest["failure_summary"]
    stderr_text = (tmp_path / ".lab-sidecar" / "tasks" / task_id / "stderr.log").read_text(encoding="utf-8")
    assert "background task recovery marked this task failed" in stderr_text


def test_list_and_open_use_manifest_task_state(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    command = f'"{sys.executable}" examples/simple-success/train.py --output metrics.csv'
    task_id = extract_task_id(invoke(tmp_path, ["run", command, "--name", "list-open-smoke"]).output)

    listed = invoke(tmp_path, ["list"])
    opened = invoke(tmp_path, ["open", task_id])

    assert listed.exit_code == 0
    assert task_id in listed.output
    assert "completed" in listed.output
    assert "list-open-smoke" in listed.output
    assert opened.exit_code == 0
    assert str(tmp_path / ".lab-sidecar" / "tasks" / task_id) in opened.output


def test_list_limit_empty_and_open_missing_task(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0

    empty = invoke(tmp_path, ["list"])
    assert empty.exit_code == 0
    assert "No tasks found." in empty.output

    command = f'"{sys.executable}" examples/simple-success/train.py --output metrics.csv'
    first = extract_task_id(invoke(tmp_path, ["run", command, "--name", "first"]).output)
    second = extract_task_id(invoke(tmp_path, ["run", command, "--name", "second"]).output)

    limited = invoke(tmp_path, ["list", "--limit", "1"])
    assert limited.exit_code == 0
    assert second in limited.output
    assert "second" in limited.output
    assert first not in limited.output
    assert "first" not in limited.output

    missing = invoke(tmp_path, ["open", "task_missing"])
    assert missing.exit_code == 3
    assert "task 'task_missing' was not found" in missing.output


def test_list_status_filter_and_missing_manifest_are_safe(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    success_command = f'"{sys.executable}" examples/simple-success/train.py --output metrics.csv'
    failure_command = f'"{sys.executable}" examples/simple-failure/fail.py'

    completed_task = extract_task_id(invoke(tmp_path, ["run", success_command, "--name", "completed-list"]).output)
    failed_task = extract_task_id(invoke(tmp_path, ["run", failure_command, "--name", "failed-list"]).output)
    stale_dir = tmp_path / ".lab-sidecar" / "tasks" / "task_missing_manifest"
    stale_dir.mkdir()

    listed = invoke(tmp_path, ["list", "--limit", "5"])
    completed = invoke(tmp_path, ["list", "--status", "completed"])
    failed = invoke(tmp_path, ["list", "--status", "failed"])
    invalid = invoke(tmp_path, ["list", "--status", "unknown"])

    assert listed.exit_code == 0
    assert "task_id" in listed.output
    assert "finished_at" in listed.output
    assert "artifacts" in listed.output
    assert completed_task in listed.output
    assert failed_task in listed.output
    assert "task_missing_manifest" not in listed.output

    assert completed.exit_code == 0
    assert completed_task in completed.output
    assert failed_task not in completed.output

    assert failed.exit_code == 0
    assert failed_task in failed.output
    assert completed_task not in failed.output

    assert invalid.exit_code == 2


def test_status_dashboard_for_completed_ingested_failed_and_cancelled_tasks(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0

    success_command = f'"{sys.executable}" examples/simple-success/train.py --output metrics.csv'
    completed_task = extract_task_id(invoke(tmp_path, ["run", success_command, "--name", "status-complete"]).output)
    completed_status = invoke(tmp_path, ["status", completed_task])
    assert completed_status.exit_code == 0
    assert "Name: status-complete" in completed_status.output
    assert "Mode: run" in completed_status.output
    assert f"Artifact dir: .lab-sidecar/tasks/{completed_task}" in completed_status.output
    assert "Artifact types:" in completed_status.output
    assert f"labsidecar summarize {completed_task}" in completed_status.output

    ingested_task = extract_task_id(
        invoke(tmp_path, ["ingest", "examples/csv-comparison", "--name", "status-ingest"]).output
    )
    ingested_status = invoke(tmp_path, ["status", ingested_task])
    assert ingested_status.exit_code == 0
    assert "Mode: ingest" in ingested_status.output
    assert "Source: examples/csv-comparison" in ingested_status.output

    failure_command = f'"{sys.executable}" examples/simple-failure/fail.py'
    failed_task = extract_task_id(invoke(tmp_path, ["run", failure_command]).output)
    failed_status = invoke(tmp_path, ["status", failed_task])
    assert failed_status.exit_code == 0
    assert "Status: failed" in failed_status.output
    assert "Failure summary:" in failed_status.output
    assert "Diagnostic log:" in failed_status.output
    assert f"labsidecar logs {failed_task} --stream stderr --tail 40" in failed_status.output

    script = tmp_path / "long_task.py"
    script.write_text(
        "\n".join(
            [
                "import time",
                "for i in range(30):",
                "    print(f'tick={i}', flush=True)",
                "    time.sleep(0.5)",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    running_task = extract_task_id(invoke(tmp_path, ["run", f'"{sys.executable}" long_task.py', "--background"]).output)
    wait_for_output(tmp_path, running_task, "tick=0")
    assert invoke(tmp_path, ["cancel", running_task]).exit_code == 0
    cancelled_status = invoke(tmp_path, ["status", running_task])
    assert cancelled_status.exit_code == 0
    assert "Status: cancelled" in cancelled_status.output
    assert f"labsidecar artifacts {running_task}" in cancelled_status.output
    assert f"labsidecar collect {running_task}" not in cancelled_status.output


def test_summarize_before_and_after_artifacts_stays_bounded(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    command = f'"{sys.executable}" examples/simple-success/train.py --output metrics.csv'
    task_id = extract_task_id(invoke(tmp_path, ["run", command, "--name", "bounded-summary"]).output)

    before = invoke(tmp_path, ["summarize", task_id])
    assert before.exit_code == 0
    assert "Name: bounded-summary" in before.output
    assert "collection summary: (not generated)" in before.output
    assert "normalized table: (not generated)" in before.output
    assert "Best val_accuracy=0.86" not in before.output
    assert "epoch,train_loss,val_loss" not in before.output

    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    scenario_summary = json.loads((tmp_path / ".lab-sidecar" / "tasks" / task_id / "metrics" / "scenario-summary.json").read_text(encoding="utf-8"))
    assert scenario_summary["scenario_type"] == "training-run"
    assert scenario_summary["primary_metric"]["name"] == "val_accuracy"
    assert scenario_summary["omitted"]["full_metric_rows"] == "omitted_by_default"
    assert invoke(tmp_path, ["figures", task_id]).exit_code == 0
    assert invoke(tmp_path, ["report", task_id]).exit_code == 0
    assert invoke(tmp_path, ["slides", task_id]).exit_code == 0

    after = invoke(tmp_path, ["summarize", task_id])
    assert after.exit_code == 0
    assert "rows: 5" in after.output
    assert "Scenario:" in after.output
    assert "type: training-run" in after.output
    assert "primary metric: val_accuracy" in after.output
    assert "metrics/scenario-summary.json" in after.output
    assert "metrics/normalized_metrics.csv" in after.output
    assert "figures/figure-summary.json" in after.output
    assert "reports/report-fragment.md" in after.output
    assert "slides/presentation-draft.pptx" in after.output
    assert "Best val_accuracy=0.86" not in after.output
    assert "epoch,train_loss,val_loss" not in after.output


def test_summarize_failed_task_and_long_command_are_bounded(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    long_arg = "x" * 180
    command = f'"{sys.executable}" examples/simple-failure/fail.py --note {long_arg}'
    task_id = extract_task_id(invoke(tmp_path, ["run", command]).output)

    summary = invoke(tmp_path, ["summarize", task_id])

    assert summary.exit_code == 0
    assert "Status: failed" in summary.output
    assert "Failure summary:" in summary.output
    assert long_arg not in summary.output
    assert "..." in summary.output


def test_package_completed_task_exports_allowlisted_artifacts_only(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    command = f'"{sys.executable}" examples/simple-success/train.py --output metrics.csv'
    task_id = extract_task_id(invoke(tmp_path, ["run", command, "--name", "package-success"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    assert invoke(tmp_path, ["figures", task_id]).exit_code == 0
    assert invoke(tmp_path, ["report", task_id]).exit_code == 0
    assert invoke(tmp_path, ["slides", task_id]).exit_code == 0

    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    (task_path / "worker.log").write_text("worker transcript body\n", encoding="utf-8")
    worker_run_dir = task_path / "intelligence" / "worker_run_test"
    (worker_run_dir / "sandbox").mkdir(parents=True)
    (worker_run_dir / "sandbox" / "scratch.txt").write_text("scratch\n", encoding="utf-8")
    (worker_run_dir / "ai-provider-prompt.json").write_text('{"prompt":"secret"}\n', encoding="utf-8")
    (worker_run_dir / "ai-provider-response.json").write_text('{"response":"secret"}\n', encoding="utf-8")

    package_path = tmp_path / f"lab-sidecar-package-{task_id}"
    result = invoke(tmp_path, ["package", task_id, "--output", package_path.as_posix()])

    assert result.exit_code == 0
    assert "Package created:" in result.output
    assert "Type: result" in result.output
    for relative in [
        "README.md",
        "manifest.json",
        "package-summary.json",
        "artifact-index.json",
        "redaction-notes.md",
        "reproduce/command.txt",
        "reproduce/run.json",
        "reproduce/env.json",
        "reproduce/git.json",
        "reproduce/dependencies.json",
        "metrics/normalized_metrics.csv",
        "metrics/normalized_metrics.json",
        "metrics/collection-summary.json",
        "metrics/scenario-summary.json",
        "figures/figure-spec.yaml",
        "figures/figure-summary.json",
        "reports/report-fragment.md",
        "reports/report-summary.json",
        "slides/presentation-draft.pptx",
        "slides/slides-summary.json",
        "provenance/traceability.json",
    ]:
        assert (package_path / relative).is_file(), relative
    assert list((package_path / "figures").glob("*.png"))
    assert list((package_path / "figures").glob("*.svg"))

    summary = json.loads((package_path / "package-summary.json").read_text(encoding="utf-8"))
    index = json.loads((package_path / "artifact-index.json").read_text(encoding="utf-8"))
    readme = (package_path / "README.md").read_text(encoding="utf-8")
    redaction_notes = (package_path / "redaction-notes.md").read_text(encoding="utf-8")
    included_paths = {item["package_path"] for item in index["included"]}
    omitted_paths = {item["path"] for item in index["omitted"]}
    package_metadata = {item["package_path"]: item for item in index["package_metadata"]}

    assert summary["package_type"] == "result"
    assert summary["task"]["task_id"] == task_id
    assert summary["task"]["name"] == "package-success"
    assert "metrics/normalized_metrics.csv" in included_paths
    assert "slides/presentation-draft.pptx" in included_paths
    assert "provenance/traceability.json" in included_paths
    trace_entry = next(item for item in index["included"] if item["package_path"] == "provenance/traceability.json")
    assert trace_entry["category"] == "provenance"
    assert trace_entry["sha256"]
    assert trace_entry["size_bytes"] > 0
    assert package_metadata["README.md"]["sha256"]
    assert package_metadata["README.md"]["size_bytes"] > 0
    assert package_metadata["package-summary.json"]["sha256"]
    assert package_metadata["package-summary.json"]["size_bytes"] > 0
    assert package_metadata["redaction-notes.md"]["sha256"]
    assert package_metadata["artifact-index.json"]["sha256"] is None
    assert "self-referential" in package_metadata["artifact-index.json"]["digest_omitted_reason"]
    assert "stdout.log" in omitted_paths
    assert "stderr.log" in omitted_paths
    assert "worker.log" in omitted_paths
    assert "intelligence/worker_run_test/ai-provider-prompt.json" in omitted_paths
    assert "intelligence/worker_run_test/ai-provider-response.json" in omitted_paths
    assert "intelligence/worker_run_test/sandbox" in omitted_paths
    assert ".lab-sidecar/index.sqlite" in omitted_paths
    assert "Lab-Sidecar Result Package" in readme
    assert "Failed Task Diagnostic Package" not in readme
    assert "full stdout/stderr logs" in redaction_notes
    assert not (package_path / "stdout.log").exists()
    assert not (package_path / "stderr.log").exists()
    assert not (package_path / "worker.log").exists()
    assert not (package_path / "intelligence").exists()
    assert not (package_path / ".lab-sidecar" / "index.sqlite").exists()
    assert not (package_path / "examples").exists()
    assert not (package_path / "metrics.csv").exists()
    package_traceability = json.loads((package_path / "provenance" / "traceability.json").read_text(encoding="utf-8"))
    assert package_traceability["task_id"] == task_id
    assert any(artifact["path"] == "metrics/normalized_metrics.csv" and artifact["sha256"] for artifact in package_traceability["artifacts"])
    assert any(source["sha256"] for source in package_traceability["sources"])
    package_omitted = {item["path"] for item in package_traceability["omitted"]}
    assert "stdout.log" in package_omitted
    assert "stderr.log" in package_omitted
    assert_traceability_is_bounded(package_traceability)


def test_package_failed_task_is_diagnostic_and_omits_full_logs(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    command = f'"{sys.executable}" examples/simple-failure/fail.py'
    task_id = extract_task_id(invoke(tmp_path, ["run", command, "--name", "package-failure"]).output)

    package_path = tmp_path / f"lab-sidecar-package-{task_id}"
    result = invoke(tmp_path, ["package", task_id, "--output", package_path.as_posix()])

    assert result.exit_code == 0
    assert "Type: diagnostic" in result.output
    assert (package_path / "manifest.json").is_file()
    assert (package_path / "README.md").is_file()
    assert (package_path / "redaction-notes.md").is_file()
    assert (package_path / "reproduce" / "command.txt").is_file()
    assert (package_path / "reproduce" / "run.json").is_file()
    assert (package_path / "provenance" / "traceability.json").is_file()
    assert not (package_path / "stderr.log").exists()
    assert not (package_path / "stdout.log").exists()

    readme = (package_path / "README.md").read_text(encoding="utf-8")
    summary = json.loads((package_path / "package-summary.json").read_text(encoding="utf-8"))
    index = json.loads((package_path / "artifact-index.json").read_text(encoding="utf-8"))

    assert "Failed Task Diagnostic Package" in readme
    assert "not a successful experiment summary" in readme
    assert "FileNotFoundError" in readme
    assert summary["package_type"] == "diagnostic"
    assert summary["task"]["status"] == "failed"
    assert summary["task"]["failure_summary"]
    unavailable_paths = {item["path"] for item in index["unavailable"]}
    omitted_paths = {item["path"] for item in index["omitted"]}
    assert "metrics/normalized_metrics.csv" in unavailable_paths
    assert "figures/figure-summary.json" in unavailable_paths
    assert "provenance/traceability.json" not in unavailable_paths
    assert "stdout.log" in omitted_paths
    assert "stderr.log" in omitted_paths
    traceability = json.loads((package_path / "provenance" / "traceability.json").read_text(encoding="utf-8"))
    claim_ids = {item["claim_id"] for item in traceability["claim_traces"]}
    assert "slides.diagnostic.failed_status" not in claim_ids
    assert traceability["task"]["status"] == "failed"
    assert any(artifact["artifact_id"] == "log_stderr" for artifact in traceability["artifacts"])
    assert_traceability_is_bounded(traceability)


def test_package_ingested_task_omits_raw_source_refs_and_source_files(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(
        invoke(tmp_path, ["ingest", "examples/csv-comparison", "--name", "package-ingest"]).output
    )
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0

    package_path = tmp_path / f"lab-sidecar-package-{task_id}"
    result = invoke(tmp_path, ["package", task_id, "--output", package_path.as_posix()])

    assert result.exit_code == 0
    assert "Type: result" in result.output
    assert (package_path / "metrics" / "normalized_metrics.csv").is_file()
    assert (package_path / "provenance" / "traceability.json").is_file()
    assert not (package_path / "raw" / "source_refs.json").exists()
    assert not (package_path / "examples").exists()

    index = json.loads((package_path / "artifact-index.json").read_text(encoding="utf-8"))
    omitted_paths = {item["path"] for item in index["omitted"]}
    unavailable_paths = {item["path"] for item in index["unavailable"]}

    assert "raw/source_refs.json" in omitted_paths
    assert "examples/csv-comparison" in omitted_paths
    assert "reproduce/command.txt" in unavailable_paths
    assert "reproduce/run.json" in unavailable_paths
    traceability = json.loads((package_path / "provenance" / "traceability.json").read_text(encoding="utf-8"))
    assert "examples/csv-comparison" in {item["path"] for item in traceability["omitted"]}
    assert any(source["path"].endswith("baseline.csv") and source["sha256"] for source in traceability["sources"])
    assert_traceability_is_bounded(traceability)


def test_package_missing_task_and_invalid_output_path(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0

    missing = invoke(tmp_path, ["package", "task_missing", "--output", "missing-package"])
    assert missing.exit_code == 3
    assert "task 'task_missing' was not found" in missing.output

    task_id = extract_task_id(
        invoke(tmp_path, ["run", f'"{sys.executable}" examples/simple-success/train.py --output metrics.csv']).output
    )
    output_file = tmp_path / "not-a-package-dir"
    output_file.write_text("occupied\n", encoding="utf-8")

    invalid = invoke(tmp_path, ["package", task_id, "--output", output_file.as_posix()])
    assert invalid.exit_code == 2
    assert "package output path is not usable" in invalid.output
    assert "not a directory" in invalid.output

    non_empty_dir = tmp_path / "occupied-package-dir"
    non_empty_dir.mkdir()
    (non_empty_dir / "existing.txt").write_text("keep me\n", encoding="utf-8")

    invalid_non_empty = invoke(tmp_path, ["package", task_id, "--output", non_empty_dir.as_posix()])
    assert invalid_non_empty.exit_code == 2
    assert "package output path is not usable" in invalid_non_empty.output
    assert "not empty" in invalid_non_empty.output


def test_compare_two_completed_tasks_with_shared_metrics(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    first_source = tmp_path / "run-a"
    second_source = tmp_path / "run-b"
    first_source.mkdir()
    second_source.mkdir()
    (first_source / "metrics.csv").write_text(
        "epoch,model,val_accuracy,val_loss\n1,a,0.71,0.52\n2,a,0.82,0.42\n",
        encoding="utf-8",
    )
    (second_source / "metrics.csv").write_text(
        "epoch,model,val_accuracy,val_loss\n1,b,0.75,0.49\n2,b,0.86,0.37\n",
        encoding="utf-8",
    )
    first_task = extract_task_id(invoke(tmp_path, ["ingest", "run-a", "--name", "run a"]).output)
    second_task = extract_task_id(invoke(tmp_path, ["ingest", "run-b", "--name", "run b"]).output)
    assert invoke(tmp_path, ["collect", first_task]).exit_code == 0
    assert invoke(tmp_path, ["collect", second_task]).exit_code == 0

    result = invoke(tmp_path, ["compare", first_task, second_task])

    assert result.exit_code == 0
    assert "Compared tasks: 2" in result.output
    assert "Common numeric fields:" in result.output
    assert "val_accuracy" in result.output
    assert "0.82" in result.output
    assert "0.86" in result.output
    assert "metrics/normalized_metrics.csv" in result.output
    assert "a,0.82,0.42" not in result.output
    assert "b,0.86,0.37" not in result.output


def test_algorithm_benchmark_scenario_summary_from_ingest_config(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    config_path = tmp_path / "algorithm-benchmark.yaml"
    config_path.write_text(
        "\n".join(
            [
                "sources:",
                "  - examples/algorithm-benchmark/results.json",
                "fields:",
                "  algorithm: algorithm",
                "  seed: seed",
                "  input_size: input_size",
                "  runtime_ms:",
                "    source: runtime_ms",
                "    unit: ms",
                "  memory_mb:",
                "    source: memory_mb",
                "    unit: MB",
                "groups:",
                "  primary: algorithm",
                "  secondary: seed",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/algorithm-benchmark"]).output)

    collect_result = invoke(tmp_path, ["collect", task_id, "--config", config_path.as_posix()])
    assert collect_result.exit_code == 0
    assert "Scenario summary:" in collect_result.output
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    scenario = json.loads((task_path / "metrics" / "scenario-summary.json").read_text(encoding="utf-8"))

    assert scenario["scenario_type"] == "algorithm-benchmark"
    assert scenario["primary_metric"]["name"] == "runtime_ms"
    assert scenario["primary_metric"]["direction"] == "min"
    assert scenario["groups"]["primary"] == "algorithm"
    assert scenario["groups"]["seed"] == "seed"
    assert scenario["seed_aggregates"]["present"] is True
    assert scenario["seed_aggregates"]["items"][0]["group"]["algorithm"] == "quick_sort"
    assert scenario["seed_aggregates"]["claim_limit"] == "descriptive aggregate only; no statistical significance is inferred"
    serialized = json.dumps(scenario, ensure_ascii=False)
    assert "full_metric_rows" in serialized
    assert '"runs"' not in serialized


def test_training_run_scenario_summary_warns_when_primary_metric_missing(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    source = tmp_path / "measurement-run"
    source.mkdir()
    secret_comment = "SECRET-MEASUREMENT-COMMENT-" + "x" * 180
    (source / "metrics.csv").write_text(
        "\n".join(
            [
                "method,epoch,measurement,private_comment",
                f"baseline,1,10.1,{secret_comment}",
                f"baseline,2,10.3,{secret_comment}",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    (tmp_path / "metrics.yaml").write_text(
        "\n".join(
            [
                "sources:",
                "  - measurement-run/metrics.csv",
                "fields:",
                "  method: method",
                "  epoch: epoch",
                "  measurement: measurement",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "measurement-run"]).output)

    result = invoke(tmp_path, ["collect", task_id, "--config", "metrics.yaml"])

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    scenario = json.loads((task_path / "metrics" / "scenario-summary.json").read_text(encoding="utf-8"))
    serialized = json.dumps(scenario, ensure_ascii=False).lower()
    assert scenario["scenario_type"] == "training-run"
    assert scenario["primary_metric"]["name"] is None
    assert scenario["primary_metric"]["direction"] is None
    assert scenario["best_rows"] == []
    assert any("primary metric was not detected" in warning for warning in scenario["warnings"])
    assert "secret-measurement-comment" not in serialized
    assert "superior" not in serialized
    assert "winner" not in serialized


def test_nested_json_config_collects_algorithm_benchmark_scenario(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    source = tmp_path / "nested-results"
    source.mkdir()
    (source / "results.json").write_text(
        json.dumps(
            {
                "experiment": "nested-benchmark",
                "runs": [
                    {
                        "algorithm": "baseline",
                        "seed": 1,
                        "input": {"size": 100},
                        "metrics": {"runtime_ms": "55.2", "memory_mb": "120"},
                        "prompt": "SECRET-NESTED-PROMPT-BASELINE",
                    },
                    {
                        "algorithm": "candidate",
                        "seed": 1,
                        "input": {"size": 100},
                        "metrics": {"runtime_ms": "41.7", "memory_mb": "130"},
                        "private_comment": "SECRET-NESTED-PRIVATE-COMMENT",
                    },
                    {
                        "algorithm": "candidate",
                        "seed": 2,
                        "input": {"size": 100},
                        "metrics": {"runtime_ms": "39.9", "memory_mb": "131"},
                        "error_message": "SECRET-NESTED-ERROR-MESSAGE",
                    },
                ],
            }
        )
        + "\n",
        encoding="utf-8",
    )
    (tmp_path / "metrics.yaml").write_text(
        "\n".join(
            [
                "sources:",
                "  - nested-results/results.json",
                "fields:",
                "  algorithm: algorithm",
                "  seed: seed",
                "  input_size: input_size",
                "  runtime_ms:",
                "    source: metrics_runtime_ms",
                "    unit: ms",
                "  memory_mb:",
                "    source: metrics_memory_mb",
                "    unit: MB",
                "groups:",
                "  primary: algorithm",
                "  secondary: seed",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "nested-results"]).output)

    result = invoke(tmp_path, ["collect", task_id, "--config", "metrics.yaml"])

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    rows = read_csv_rows(task_path / "metrics" / "normalized_metrics.csv")
    scenario = json.loads((task_path / "metrics" / "scenario-summary.json").read_text(encoding="utf-8"))
    serialized = json.dumps(scenario, ensure_ascii=False)
    assert len(rows) == 3
    assert {"source_file", "algorithm", "seed", "input_size", "runtime_ms", "memory_mb"}.issubset(rows[0])
    assert scenario["scenario_type"] == "algorithm-benchmark"
    assert scenario["primary_metric"]["name"] == "runtime_ms"
    assert scenario["primary_metric"]["direction"] == "min"
    assert scenario["best_rows"][0]["value"] == 39.9
    assert scenario["seed_aggregates"]["present"] is True
    assert "SECRET-NESTED" not in serialized
    assert '"runs"' not in serialized


def test_multi_source_scenario_summary_stays_bounded(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    source = tmp_path / "many-results"
    source.mkdir()
    secret_note = "SECRET-MULTI-SOURCE-NOTE-" + "x" * 180
    for index in range(45):
        (source / f"run_{index:02d}.csv").write_text(
            "\n".join(
                [
                    "algorithm,seed,runtime_ms,notes",
                    f"algo_{index % 3},{index},{100 - index},{secret_note}",
                ]
            )
            + "\n",
            encoding="utf-8",
        )
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "many-results"]).output)

    result = invoke(tmp_path, ["collect", task_id])

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    scenario = json.loads((task_path / "metrics" / "scenario-summary.json").read_text(encoding="utf-8"))
    serialized = json.dumps(scenario, ensure_ascii=False)
    assert scenario["primary_metric"]["name"] == "runtime_ms"
    assert scenario["primary_metric"]["direction"] == "min"
    assert len(scenario["evidence"]["source_files"]) == 20
    assert scenario["evidence"]["omitted_source_file_count"] == 25
    assert len(scenario["last_rows"]) <= 6
    assert len(scenario["best_rows"]) <= 4
    assert "SECRET-MULTI-SOURCE-NOTE" not in serialized


def test_compare_missing_metrics_no_common_numeric_and_too_many_tasks(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    metric_tasks: list[str] = []
    for index in range(5):
        source = tmp_path / f"run-{index}"
        source.mkdir()
        (source / "metrics.csv").write_text(
            f"epoch,score\n1,{0.7 + index / 100:.2f}\n",
            encoding="utf-8",
        )
        task_id = extract_task_id(invoke(tmp_path, ["ingest", f"run-{index}"]).output)
        assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
        metric_tasks.append(task_id)

    missing_source = tmp_path / "missing-metrics"
    missing_source.mkdir()
    (missing_source / "notes.txt").write_text("no metrics\n", encoding="utf-8")
    missing_task = extract_task_id(invoke(tmp_path, ["ingest", "missing-metrics"]).output)

    missing = invoke(tmp_path, ["compare", metric_tasks[0], missing_task])
    too_few = invoke(tmp_path, ["compare", metric_tasks[0]])
    too_many = invoke(tmp_path, ["compare", *metric_tasks, missing_task])

    assert missing.exit_code == 5
    assert "metrics are missing" in missing.output
    assert "labsidecar collect <task_id>" in missing.output
    assert too_few.exit_code == 2
    assert "at least 2 task ids" in too_few.output
    assert too_many.exit_code == 2
    assert "at most 5 task ids" in too_many.output

    first_text = tmp_path / "text-a"
    second_text = tmp_path / "text-b"
    first_text.mkdir()
    second_text.mkdir()
    (first_text / "metrics.csv").write_text("model,label\nalpha,good\n", encoding="utf-8")
    (second_text / "metrics.csv").write_text("model,label\nbeta,better\n", encoding="utf-8")
    first_text_task = extract_task_id(invoke(tmp_path, ["ingest", "text-a"]).output)
    second_text_task = extract_task_id(invoke(tmp_path, ["ingest", "text-b"]).output)
    assert invoke(tmp_path, ["collect", first_text_task]).exit_code == 0
    assert invoke(tmp_path, ["collect", second_text_task]).exit_code == 0

    no_common_numeric = invoke(tmp_path, ["compare", first_text_task, second_text_task])

    assert no_common_numeric.exit_code == 5
    assert "no common numeric metric fields" in no_common_numeric.output


def test_artifact_chain_prints_next_steps(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0

    ingest_result = invoke(tmp_path, ["ingest", "examples/csv-comparison"])

    assert ingest_result.exit_code == 0
    task_id = extract_task_id(ingest_result.output)
    assert f"Artifacts: .lab-sidecar/tasks/{task_id}" in ingest_result.output
    assert "Next:" in ingest_result.output
    assert f"labsidecar collect {task_id}" in ingest_result.output
    assert f"labsidecar artifacts {task_id}" in ingest_result.output

    collect_result = invoke(tmp_path, ["collect", task_id])
    assert collect_result.exit_code == 0
    assert "Next:" in collect_result.output
    assert f"labsidecar figures {task_id}" in collect_result.output
    assert f"labsidecar report {task_id}" in collect_result.output

    figures_result = invoke(tmp_path, ["figures", task_id])
    assert figures_result.exit_code == 0
    assert "Next:" in figures_result.output
    assert f"labsidecar report {task_id}" in figures_result.output
    assert f"labsidecar slides {task_id}" in figures_result.output

    report_result = invoke(tmp_path, ["report", task_id])
    assert report_result.exit_code == 0
    assert "Next:" in report_result.output
    assert f"labsidecar slides {task_id}" in report_result.output

    slides_result = invoke(tmp_path, ["slides", task_id])
    assert slides_result.exit_code == 0
    assert "Next:" in slides_result.output
    assert f"labsidecar artifacts {task_id}" in slides_result.output
    assert f"labsidecar open {task_id}" in slides_result.output


def test_cancel_non_running_task_returns_state_error(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0

    command = f'"{sys.executable}" examples/simple-success/train.py --output metrics.csv'
    task_id = extract_task_id(invoke(tmp_path, ["run", command]).output)

    result = invoke(tmp_path, ["cancel", task_id])
    assert result.exit_code == 4
    assert "is not running" in result.output
    assert "Current status: completed" in result.output


def test_ingest_existing_directory(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    source = tmp_path / "examples" / "csv-comparison"
    before = sorted(path.name for path in source.iterdir())

    result = invoke(tmp_path, ["ingest", "examples/csv-comparison", "--name", "csv import"])

    assert result.exit_code == 0
    task_id = extract_task_id(result.output)
    manifest = read_manifest(tmp_path, task_id)
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    refs_path = task_path / "raw" / "source_refs.json"
    refs = json.loads(refs_path.read_text(encoding="utf-8"))

    assert manifest["schema_version"] == "1"
    assert manifest["mode"] == "ingest"
    assert manifest["status"] == "completed"
    assert manifest["command"] is None
    assert manifest["source_path"] == "examples/csv-comparison"
    assert manifest["working_dir"] == "."
    assert manifest["exit_code"] == 0
    assert manifest["name"] == "csv import"
    assert (task_path / "stdout.log").is_file()
    assert (task_path / "stderr.log").is_file()
    assert refs_path.is_file()
    assert refs["source_path"] == "examples/csv-comparison"
    assert refs["path_kind"] == "relative"
    assert refs["source_type"] == "directory"
    assert refs["child_count"] == len(before)
    assert "examples/csv-comparison/baseline.csv" in refs["candidate_files"]
    baseline_ref = next(child for child in refs["children"] if child["name"] == "baseline.csv")
    assert baseline_ref["sha256"]
    assert baseline_ref["size_bytes"] > 0
    assert sorted(path.name for path in source.iterdir()) == before

    status = invoke(tmp_path, ["status", task_id])
    assert status.exit_code == 0
    assert "Status: completed" in status.output

    artifacts = invoke(tmp_path, ["artifacts", task_id])
    assert artifacts.exit_code == 0
    assert "[raw]" in artifacts.output
    assert "raw/source_refs.json" in artifacts.output

    (tmp_path / ".lab-sidecar" / "index.sqlite").unlink()
    assert invoke(tmp_path, ["status", task_id]).exit_code == 0
    assert invoke(tmp_path, ["artifacts", task_id]).exit_code == 0


def test_ingest_existing_file(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    source = tmp_path / "examples" / "algorithm-benchmark" / "results.json"
    before = source.read_text(encoding="utf-8")

    result = invoke(tmp_path, ["ingest", "examples/algorithm-benchmark/results.json"])

    assert result.exit_code == 0
    task_id = extract_task_id(result.output)
    manifest = read_manifest(tmp_path, task_id)
    refs = json.loads(
        (
            tmp_path
            / ".lab-sidecar"
            / "tasks"
            / task_id
            / "raw"
            / "source_refs.json"
        ).read_text(encoding="utf-8")
    )

    assert manifest["mode"] == "ingest"
    assert manifest["source_path"] == "examples/algorithm-benchmark/results.json"
    assert manifest["exit_code"] == 0
    assert refs["source_type"] == "file"
    assert refs["path_kind"] == "relative"
    assert refs["size_bytes"] == source.stat().st_size
    assert refs["sha256"]
    assert refs["is_candidate"] is True
    assert source.read_text(encoding="utf-8") == before


def test_ingest_missing_path_and_uninitialized_workspace(tmp_path: Path) -> None:
    missing = invoke(tmp_path, ["ingest", "missing-results"])
    assert missing.exit_code == 2
    assert "workspace is not initialized" in missing.output

    assert invoke(tmp_path, ["init"]).exit_code == 0
    missing_after_init = invoke(tmp_path, ["ingest", "missing-results"])
    assert missing_after_init.exit_code == 2
    assert "does not exist" in missing_after_init.output


def test_cli_exit_codes_for_missing_tasks_and_uninitialized_run(tmp_path: Path) -> None:
    run_without_init = invoke(tmp_path, ["run", f'"{sys.executable}" -c "print(1)"'])
    assert run_without_init.exit_code == 2
    assert "workspace is not initialized" in run_without_init.output

    assert invoke(tmp_path, ["init"]).exit_code == 0

    missing_status = invoke(tmp_path, ["status", "task_missing"])
    assert missing_status.exit_code == 3

    missing_logs = invoke(tmp_path, ["logs", "task_missing"])
    assert missing_logs.exit_code == 3

    missing_artifacts = invoke(tmp_path, ["artifacts", "task_missing"])
    assert missing_artifacts.exit_code == 3

    missing_cancel = invoke(tmp_path, ["cancel", "task_missing"])
    assert missing_cancel.exit_code == 3


def test_logs_missing_artifact_returns_exit_code_5(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    command = f'"{sys.executable}" examples/simple-success/train.py --output metrics.csv'
    task_id = extract_task_id(invoke(tmp_path, ["run", command]).output)

    (tmp_path / ".lab-sidecar" / "tasks" / task_id / "stdout.log").unlink()

    result = invoke(tmp_path, ["logs", task_id, "--stream", "stdout"])
    assert result.exit_code == 5
    assert "log files are not available" in result.output


def test_collect_csv_comparison_ingest_generates_normalized_metrics(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)

    result = invoke(tmp_path, ["collect", task_id])

    assert result.exit_code == 0
    assert "Collected metrics for" in result.output
    assert "Detected fields:" in result.output

    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    csv_path = task_path / "metrics" / "normalized_metrics.csv"
    json_path = task_path / "metrics" / "normalized_metrics.json"
    summary_path = task_path / "metrics" / "collection-summary.json"

    assert csv_path.is_file()
    assert json_path.is_file()
    assert summary_path.is_file()

    rows = read_csv_rows(csv_path)
    assert rows
    assert {"source_file", "epoch", "model", "seed", "val_accuracy", "val_loss"}.issubset(rows[0])
    assert {Path(row["source_file"]).name for row in rows} == {
        "baseline.csv",
        "model_a.csv",
        "model_b.csv",
    }

    json_rows = json.loads(json_path.read_text(encoding="utf-8"))
    assert len(json_rows) == len(rows)
    assert "source_file" in json_rows[0]

    summary = json.loads(summary_path.read_text(encoding="utf-8"))
    assert summary["task_status"] == "completed"
    assert summary["row_count"] == len(rows)
    assert {"epoch", "model", "seed", "val_accuracy", "val_loss"}.issubset(summary["detected_fields"])
    assert summary["processed_files"][0]["source_provenance"]["sha256"]
    assert summary["candidates"][0]["source_provenance"]["sha256"]

    manifest = read_manifest(tmp_path, task_id)
    artifacts = {artifact["artifact_id"]: artifact for artifact in manifest["artifacts"]}
    assert artifacts["metrics_normalized_csv"]["type"] == "table"
    assert artifacts["metrics_normalized_json"]["type"] == "table"
    assert artifacts["metrics_collection_summary"]["type"] == "config"


def test_collect_records_bounded_best_checkpoint_and_anomaly_summary(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    source = tmp_path / "sweep-results"
    source.mkdir()
    (source / "metrics.csv").write_text(
        "\n".join(
            [
                "variant,run_id,config_id,seed,epoch,step,val_accuracy,val_loss,checkpoint,status,error_flag,warning_flag,incomplete_flag,unstable_flag,artifact_present,anomaly_code",
                "baseline,run_001,cfg_00,11,1,100,0.700,0.620,,ok,0,0,0,0,true,none",
                "candidate,run_002,cfg_01,13,1,100,0.820,0.440,checkpoints/run_002/step_100.pt,ok,0,0,0,0,true,none",
                "candidate,run_002,cfg_01,13,2,200,0.910,0.280,checkpoints/run_002/final.pt,ok,0,0,0,0,true,none",
                "broken,run_003,cfg_02,17,1,100,0.510,0.910,,missing_final_metric,1,0,1,0,false,missing_final_metric",
                "broken,run_003,cfg_02,17,2,200,0.500,0.940,,missing_final_metric,1,0,1,0,false,missing_final_metric",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "sweep-results"]).output)

    result = invoke(tmp_path, ["collect", task_id])

    assert result.exit_code == 0
    summary = json.loads(
        (
            tmp_path
            / ".lab-sidecar"
            / "tasks"
            / task_id
            / "metrics"
            / "collection-summary.json"
        ).read_text(encoding="utf-8")
    )
    bounded = summary["bounded_analysis"]
    assert bounded["schema_version"] == "1"
    assert bounded["row_count"] == 5

    best_accuracy = next(item for item in bounded["best_rows"] if item["metric"] == "val_accuracy")
    assert best_accuracy["direction"] == "max"
    assert best_accuracy["value"] == 0.91
    assert best_accuracy["row_number"] == 3
    assert best_accuracy["selected_fields"]["variant"] == "candidate"
    assert best_accuracy["selected_fields"]["run_id"] == "run_002"
    assert best_accuracy["selected_fields"]["checkpoint"] == "checkpoints/run_002/final.pt"
    assert best_accuracy["evidence"] == {
        "artifact_id": "metrics_normalized_csv",
        "path": "metrics/normalized_metrics.csv",
        "row_number": 3,
        "body": "omitted",
    }

    checkpoint = bounded["checkpoint_summary"]
    assert checkpoint["present"] is True
    assert checkpoint["available_checkpoint_count"] == 2
    assert checkpoint["unique_checkpoint_count"] == 2
    assert checkpoint["selected"]["checkpoint"] == "checkpoints/run_002/final.pt"
    assert checkpoint["selected"]["selection_metric"] == "val_accuracy"
    assert checkpoint["selected"]["row_number"] == 3

    anomalies = bounded["anomaly_summary"]
    assert anomalies["present"] is True
    assert anomalies["anomaly_row_count"] == 2
    assert anomalies["anomaly_group_count"] == 1
    reasons = {item["reason"] for item in anomalies["counts_by_reason"]}
    assert {
        "status=missing_final_metric",
        "error_flag=1",
        "incomplete_flag=1",
        "artifact_present=false",
        "anomaly_code=missing_final_metric",
    }.issubset(reasons)
    example = anomalies["examples"][0]
    assert example["row_count"] == 2
    assert example["first_row_number"] == 4
    assert example["last_row_number"] == 5
    assert example["selected_fields"]["run_id"] == "run_003"
    assert example["selected_fields"]["artifact_present"] == "false"
    assert example["evidence"]["body"] == "omitted"


def test_collect_run_working_dir_output_without_wrapper(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    stale_path = tmp_path / "stale.csv"
    stale_path.write_text("epoch,accuracy\n1,0.01\n", encoding="utf-8")
    old_time = time.time() - 30
    os.utime(stale_path, (old_time, old_time))

    command = f'"{sys.executable}" examples/simple-success/train.py --output metrics.csv'
    task_id = extract_task_id(invoke(tmp_path, ["run", command]).output)

    result = invoke(tmp_path, ["collect", task_id])

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    rows = read_csv_rows(task_path / "metrics" / "normalized_metrics.csv")
    assert rows
    assert {Path(row["source_file"]).name for row in rows} == {"metrics.csv"}
    summary = json.loads((task_path / "metrics" / "collection-summary.json").read_text(encoding="utf-8"))
    assert {
        (Path(candidate["source_file"]).name, candidate["origin"])
        for candidate in summary["candidates"]
    } == {("metrics.csv", "run_working_dir")}


def test_collect_json_algorithm_benchmark_ingest_generates_normalized_metrics(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(
        invoke(tmp_path, ["ingest", "examples/algorithm-benchmark/results.json"]).output
    )

    result = invoke(tmp_path, ["collect", task_id])

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    rows = read_csv_rows(task_path / "metrics" / "normalized_metrics.csv")

    assert rows
    assert {
        "source_file",
        "benchmark_name",
        "algorithm",
        "input_size",
        "seed",
        "runtime_ms",
        "memory_mb",
    }.issubset(rows[0])
    assert rows[0]["source_file"].endswith("examples/algorithm-benchmark/results.json")

    json_rows = json.loads((task_path / "metrics" / "normalized_metrics.json").read_text(encoding="utf-8"))
    assert len(json_rows) == len(rows)
    assert json_rows[0]["source_file"].endswith("examples/algorithm-benchmark/results.json")

    summary = json.loads((task_path / "metrics" / "collection-summary.json").read_text(encoding="utf-8"))
    assert {"algorithm", "seed", "runtime_ms", "memory_mb"}.issubset(summary["detected_fields"])


def test_network_experiment_csv_collects_metrics_and_generates_bar_figure(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    source = tmp_path / "network-results"
    source.mkdir()
    (source / "all-results.csv").write_text(
        "\n".join(
            [
                "Stage,Scenario,DurationSec,DATA_TIMER,BACKLOG_FACTOR,AvgUtil,DataTimeoutPerMin,Score,BadCrcTotal,WallSeconds",
                "coarse,default flood,120,1500,3,88.02,24.5,87.285,12,120.8",
                "coarse,default flood,120,1500,4,85.15,29,84.28,14,120.7",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "network-results"]).output)

    result = invoke(tmp_path, ["collect", task_id])

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    rows = read_csv_rows(task_path / "metrics" / "normalized_metrics.csv")
    assert len(rows) == 2
    assert {"DurationSec", "AvgUtil", "DataTimeoutPerMin", "Score", "BadCrcTotal", "WallSeconds"}.issubset(rows[0])
    summary = json.loads((task_path / "metrics" / "collection-summary.json").read_text(encoding="utf-8"))
    assert {
        "DurationSec",
        "AvgUtil",
        "DataTimeoutPerMin",
        "Score",
        "BadCrcTotal",
        "WallSeconds",
    }.issubset(summary["detected_fields"])

    figure_result = invoke(tmp_path, ["figures", task_id])

    assert figure_result.exit_code == 0
    figures_dir = task_path / "figures"
    assert_non_empty_file(figures_dir / "bar_score_by_data_timer.png")
    assert_non_empty_file(figures_dir / "bar_score_by_data_timer.svg")
    figure_summary = json.loads((figures_dir / "figure-summary.json").read_text(encoding="utf-8"))
    assert figure_summary["generated_figures"][0]["chart_type"] == "bar"
    assert figure_summary["generated_figures"][0]["x"] == "DATA_TIMER"
    assert figure_summary["generated_figures"][0]["y"] == "Score"


def test_collect_missing_task_returns_exit_code_3(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0

    result = invoke(tmp_path, ["collect", "task_missing"])

    assert result.exit_code == 3
    assert "was not found" in result.output


def test_collect_without_supported_metrics_returns_exit_code_5(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    source = tmp_path / "notes-only"
    source.mkdir()
    (source / "notes.txt").write_text("no metrics here\n", encoding="utf-8")
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "notes-only"]).output)

    result = invoke(tmp_path, ["collect", task_id])

    assert result.exit_code == 5
    assert "no CSV/JSON metric candidates were found" in result.output
    assert "collection-summary.json" in result.output
    summary_path = tmp_path / ".lab-sidecar" / "tasks" / task_id / "metrics" / "collection-summary.json"
    assert summary_path.is_file()
    summary = json.loads(summary_path.read_text(encoding="utf-8"))
    assert summary["candidate_count"] == 0
    assert summary["row_count"] == 0


def test_collect_bad_and_empty_inputs_record_diagnostics_without_outputs(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    source = tmp_path / "bad-inputs"
    source.mkdir()
    (source / "bad.json").write_text('{"epoch": 1, "accuracy": ', encoding="utf-8")
    (source / "empty.csv").write_text("", encoding="utf-8")
    (source / "missing_metric_columns.csv").write_text("name,notes\nalpha,no metric columns\n", encoding="utf-8")
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "bad-inputs"]).output)

    result = invoke(tmp_path, ["collect", task_id])

    assert result.exit_code == 5
    assert "CSV/JSON candidates were found, but no metrics could be collected" in result.output
    assert "collection-summary.json" in result.output
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    summary_path = task_path / "metrics" / "collection-summary.json"
    assert summary_path.is_file()
    summary = json.loads(summary_path.read_text(encoding="utf-8"))
    assert summary["candidate_count"] == 3
    assert summary["row_count"] == 0
    assert not summary["output_files"]
    assert not (task_path / "metrics" / "scenario-summary.json").exists()
    skipped = {(Path(item["source_file"]).name, item["reason"]) for item in summary["skipped_files"]}
    assert ("bad.json", "parse_failed") in skipped
    assert ("empty.csv", "no_detected_metrics") in skipped
    assert ("missing_metric_columns.csv", "no_detected_metrics") in skipped
    diagnostic_reasons = {item["reason"] for item in summary["diagnostics"]}
    assert "parse_failed" in diagnostic_reasons
    assert any("Failed to parse bad-inputs/bad.json" in warning for warning in summary["warnings"])
    serialized = json.dumps(summary, ensure_ascii=False)
    assert '{"epoch": 1, "accuracy":' not in serialized
    assert "alpha,no metric columns" not in serialized
    assert not (task_path / "metrics" / "normalized_metrics.csv").exists()
    assert not (task_path / "metrics" / "normalized_metrics.json").exists()

    manifest = read_manifest(tmp_path, task_id)
    artifacts = {artifact["artifact_id"]: artifact for artifact in manifest["artifacts"]}
    assert "metrics_collection_summary" in artifacts
    assert "metrics_normalized_csv" not in artifacts
    assert "metrics_normalized_json" not in artifacts


def test_collect_repeated_bad_input_does_not_duplicate_summary_artifact(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    source = tmp_path / "bad-json"
    source.mkdir()
    (source / "results.json").write_text("{not valid json", encoding="utf-8")
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "bad-json"]).output)

    assert invoke(tmp_path, ["collect", task_id]).exit_code == 5
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 5

    manifest = read_manifest(tmp_path, task_id)
    artifact_ids = [artifact["artifact_id"] for artifact in manifest["artifacts"]]
    assert artifact_ids.count("metrics_collection_summary") == 1


def test_collect_after_deleting_sqlite_uses_manifest_and_source_refs(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    index_path = tmp_path / ".lab-sidecar" / "index.sqlite"
    index_path.unlink()

    result = invoke(tmp_path, ["collect", task_id])

    assert result.exit_code == 0
    assert index_path.is_file()
    csv_path = tmp_path / ".lab-sidecar" / "tasks" / task_id / "metrics" / "normalized_metrics.csv"
    assert csv_path.is_file()
    assert read_csv_rows(csv_path)


def test_collect_with_config_maps_nonstandard_fields_selects_sources_and_records_units(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    source = tmp_path / "configured-results"
    source.mkdir()
    (source / "run_a.csv").write_text(
        "\n".join(
            [
                "iter,algo,trial,score_pct,time_ms",
                "1,baseline,1,0.70,45",
                "2,baseline,1,0.75,43",
                "1,candidate,1,0.80,35",
                "2,candidate,1,0.86,31",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    (source / "unrelated.csv").write_text("epoch,accuracy\n1,0.01\n", encoding="utf-8")
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "configured-results"]).output)
    (tmp_path / "metrics.yaml").write_text(
        "\n".join(
            [
                "sources:",
                "  - configured-results/run_a.csv",
                "fields:",
                "  epoch: iter",
                "  method: algo",
                "  seed: trial",
                "  accuracy: score_pct",
                "  latency_ms:",
                "    source: time_ms",
                "    unit: ms",
                "units:",
                "  accuracy: ratio",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    result = invoke(tmp_path, ["collect", task_id, "--config", "metrics.yaml"])

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    rows = read_csv_rows(task_path / "metrics" / "normalized_metrics.csv")
    assert len(rows) == 4
    assert {"source_file", "epoch", "method", "seed", "accuracy", "latency_ms"}.issubset(rows[0])
    assert {Path(row["source_file"]).name for row in rows} == {"run_a.csv"}
    assert {row["method"] for row in rows} == {"baseline", "candidate"}

    summary = json.loads((task_path / "metrics" / "collection-summary.json").read_text(encoding="utf-8"))
    assert summary["config_path"] == "metrics.yaml"
    assert summary["candidate_count"] == 1
    assert summary["candidates"][0]["origin"] == "config"
    assert summary["processed_files"][0]["mapped_fields"] == ["epoch", "method", "seed", "accuracy", "latency_ms"]
    assert summary["units"] == {"accuracy": "ratio", "latency_ms": "ms"}
    assert Path(summary["processed_files"][0]["source_file"]).name == "run_a.csv"
    assert "unrelated.csv" not in json.dumps(summary)


def test_collect_with_config_missing_field_fails_with_diagnostics(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    source = tmp_path / "configured-results"
    source.mkdir()
    (source / "run_a.csv").write_text("iter,algo,score_pct\n1,baseline,0.70\n", encoding="utf-8")
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "configured-results"]).output)
    (tmp_path / "metrics.yaml").write_text(
        "\n".join(
            [
                "sources:",
                "  - configured-results/run_a.csv",
                "fields:",
                "  epoch: iter",
                "  accuracy: score_pct",
                "  method: algo",
                "  seed: missing_seed",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    result = invoke(tmp_path, ["collect", task_id, "--config", "metrics.yaml"])

    assert result.exit_code == 5
    assert "missing_seed" in result.output
    assert "collection-summary.json" in result.output
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    summary = json.loads((task_path / "metrics" / "collection-summary.json").read_text(encoding="utf-8"))
    assert summary["row_count"] == 0
    assert summary["skipped_files"][0]["reason"] == "missing_configured_field"
    assert "missing_seed" in summary["warnings"][0]
    assert not (task_path / "metrics" / "normalized_metrics.csv").exists()


def test_collect_and_figures_with_explicit_config_are_stable_and_reuse_mapped_fields(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    source = tmp_path / "configured-results"
    source.mkdir()
    (source / "run_a.csv").write_text(
        "\n".join(
            [
                "iter,algo,trial,score_pct",
                "1,baseline,1,0.70",
                "2,baseline,1,0.75",
                "1,candidate,1,0.80",
                "2,candidate,1,0.86",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "configured-results"]).output)
    (tmp_path / "metrics.yaml").write_text(
        "\n".join(
            [
                "sources:",
                "  - configured-results/run_a.csv",
                "fields:",
                "  epoch: iter",
                "  method: algo",
                "  seed: trial",
                "  accuracy: score_pct",
                "units:",
                "  accuracy: ratio",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    (tmp_path / "figure.yaml").write_text(
        "\n".join(
            [
                "figure_id: mapped_accuracy",
                "chart_type: line",
                "title: Configured Accuracy",
                "x: epoch",
                "y: accuracy",
                "group_by: method",
                "output:",
                "  - figures/mapped_accuracy.png",
                "  - figures/mapped_accuracy.svg",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    assert invoke(tmp_path, ["collect", task_id, "--config", "metrics.yaml"]).exit_code == 0
    assert invoke(tmp_path, ["collect", task_id, "--config", "metrics.yaml"]).exit_code == 0
    assert invoke(tmp_path, ["figures", task_id, "--spec", "figure.yaml"]).exit_code == 0
    assert invoke(tmp_path, ["figures", task_id, "--spec", "figure.yaml"]).exit_code == 0

    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    assert_non_empty_file(task_path / "figures" / "mapped_accuracy.png")
    assert_non_empty_file(task_path / "figures" / "mapped_accuracy.svg")
    figure_summary = json.loads((task_path / "figures" / "figure-summary.json").read_text(encoding="utf-8"))
    assert figure_summary["generated_figures"][0]["x"] == "epoch"
    assert figure_summary["generated_figures"][0]["y"] == "accuracy"
    assert figure_summary["generated_figures"][0]["group_by"] == "method"

    manifest = read_manifest(tmp_path, task_id)
    artifact_ids = [artifact["artifact_id"] for artifact in manifest["artifacts"]]
    assert len(artifact_ids) == len(set(artifact_ids))
    assert artifact_ids.count("metrics_normalized_csv") == 1
    assert artifact_ids.count("figures_summary") == 1
    assert artifact_ids.count("figure_mapped_accuracy_png") == 1


def test_collect_stage3_messy_nested_results_with_include_exclude_and_aliases(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    write_stage3_messy_results(tmp_path)
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "messy-results", "--name", "stage3 messy"]).output)
    refs = json.loads(
        (
            tmp_path
            / ".lab-sidecar"
            / "tasks"
            / task_id
            / "raw"
            / "source_refs.json"
        ).read_text(encoding="utf-8")
    )
    assert refs["candidate_file_count"] >= 8
    assert any(item["path"].endswith("baseline/seed_1/metrics.csv") for item in refs["candidate_file_refs"])
    assert "messy-results/baseline/seed_1/metrics.csv" not in refs["candidate_files"]
    (tmp_path / "metrics.yaml").write_text(
        "\n".join(
            [
                "sources:",
                "  include:",
                "    - messy-results/**/*.csv",
                "  exclude:",
                "    - messy-results/debug/*.csv",
                "    - messy-results/scratch/*",
                "fields:",
                "  epoch:",
                "    sources: [epoch, step, iter]",
                "  method:",
                "    sources: [model, method, algo, variant]",
                "  seed:",
                "    sources: [seed, trial, run_id]",
                "  accuracy:",
                "    sources: [val_accuracy, score_pct, acc]",
                "    unit: ratio",
                "  latency_ms:",
                "    sources: [runtime_ms, latency_ms, time_ms]",
                "    unit: ms",
                "groups:",
                "  primary: method",
                "  secondary: seed",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    (tmp_path / "figure.yaml").write_text(
        "\n".join(
            [
                "figure_id: messy_accuracy",
                "chart_type: line",
                "title: Messy Accuracy",
                "x: epoch",
                "y: accuracy",
                "group_by: method",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    result = invoke(tmp_path, ["collect", task_id, "--config", "metrics.yaml"])

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    rows = read_csv_rows(task_path / "metrics" / "normalized_metrics.csv")
    assert len(rows) == 12
    assert {"source_file", "epoch", "method", "seed", "accuracy", "latency_ms"}.issubset(rows[0])
    assert {row["method"] for row in rows} == {"baseline", "candidate"}
    assert {row["seed"] for row in rows} == {"1", "2", "3"}
    assert not any("debug" in row["source_file"] or "scratch" in row["source_file"] for row in rows)

    summary = json.loads((task_path / "metrics" / "collection-summary.json").read_text(encoding="utf-8"))
    assert summary["candidate_count"] == 6
    assert summary["units"] == {"accuracy": "ratio", "latency_ms": "ms"}
    assert summary["groups"] == {"primary": "method", "secondary": "seed"}
    assert summary["matched_source_fields"] == {
        "epoch": ["iter"],
        "method": ["algo"],
        "seed": ["trial"],
        "accuracy": ["score_pct"],
        "latency_ms": ["runtime_ms"],
    }
    skipped = {(Path(item["source_file"]).name, item["reason"]) for item in summary["skipped_files"]}
    assert ("debug_metrics.csv", "configured_source_excluded") in skipped
    assert ("scratch.csv", "configured_source_excluded") in skipped
    assert summary["unit_diagnostics"] == []

    figures = invoke(tmp_path, ["figures", task_id, "--spec", "figure.yaml"])

    assert figures.exit_code == 0
    assert_non_empty_file(task_path / "figures" / "messy_accuracy.png")
    assert_non_empty_file(task_path / "figures" / "messy_accuracy.svg")
    figure_summary = json.loads((task_path / "figures" / "figure-summary.json").read_text(encoding="utf-8"))
    assert figure_summary["units"] == {"accuracy": "ratio", "latency_ms": "ms"}
    assert figure_summary["groups"] == {"primary": "method", "secondary": "seed"}
    assert figure_summary["generated_figures"][0]["units"] == {"accuracy": "ratio"}
    assert figure_summary["generated_figures"][0]["group_by"] == "method"


def test_collect_stage3_missing_configured_source_and_not_in_refs_are_diagnosed(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    source = tmp_path / "ingested-results"
    source.mkdir()
    (source / "metrics.csv").write_text("iter,score\n1,0.8\n", encoding="utf-8")
    outside_ingest = tmp_path / "outside-ingest.csv"
    outside_ingest.write_text("iter,score\n1,0.9\n", encoding="utf-8")
    outside_workspace = tmp_path.parent / f"outside-workspace-{tmp_path.name}.csv"
    outside_workspace.write_text("iter,score\n1,0.9\n", encoding="utf-8")
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "ingested-results"]).output)
    (tmp_path / "metrics.yaml").write_text(
        "\n".join(
            [
                "sources:",
                "  include:",
                "    - ingested-results/missing.csv",
                f"    - {outside_workspace.as_posix()}",
                "    - outside-ingest.csv",
                "fields:",
                "  epoch: iter",
                "  accuracy: score",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    try:
        result = invoke(tmp_path, ["collect", task_id, "--config", "metrics.yaml"])
    finally:
        outside_workspace.unlink(missing_ok=True)

    assert result.exit_code == 5
    assert "missing.csv" in result.output
    assert "outside the workspace" in result.output
    assert "source refs" in result.output
    summary = json.loads(
        (
            tmp_path
            / ".lab-sidecar"
            / "tasks"
            / task_id
            / "metrics"
            / "collection-summary.json"
        ).read_text(encoding="utf-8")
    )
    reasons = {item["reason"] for item in summary["skipped_files"]}
    assert {"configured_source_missing", "outside_workspace", "not_in_source_refs"}.issubset(reasons)
    diagnostic_reasons = {item["reason"] for item in summary["diagnostics"]}
    assert {"configured_source_missing", "outside_workspace", "not_in_source_refs"}.issubset(
        diagnostic_reasons
    )
    assert summary["row_count"] == 0


def test_collect_stage3_unit_conflicts_are_recorded(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    source = tmp_path / "unit-results"
    source.mkdir()
    (source / "run_a.csv").write_text(
        "iter,algo,trial,runtime_ms\n1,baseline,1,25\n",
        encoding="utf-8",
    )
    (source / "run_b.csv").write_text(
        "iter,algo,trial,runtime_s\n1,candidate,1,0.02\n",
        encoding="utf-8",
    )
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "unit-results"]).output)
    (tmp_path / "metrics.yaml").write_text(
        "\n".join(
            [
                "sources:",
                "  - unit-results/*.csv",
                "fields:",
                "  epoch: iter",
                "  method: algo",
                "  seed: trial",
                "  latency_ms:",
                "    sources: [runtime_ms, runtime_s]",
                "    unit: ms",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    result = invoke(tmp_path, ["collect", task_id, "--config", "metrics.yaml"])

    assert result.exit_code == 0
    summary = json.loads(
        (
            tmp_path
            / ".lab-sidecar"
            / "tasks"
            / task_id
            / "metrics"
            / "collection-summary.json"
        ).read_text(encoding="utf-8")
    )
    assert summary["units"] == {"latency_ms": "ms"}
    assert summary["matched_source_fields"]["latency_ms"] == ["runtime_ms", "runtime_s"]
    assert summary["unit_diagnostics"]
    assert summary["unit_diagnostics"][0]["reason"] == "mixed_source_units"
    assert "runtime_ms=ms" in summary["unit_diagnostics"][0]["message"]
    assert "runtime_s=s" in summary["unit_diagnostics"][0]["message"]
    assert any("mixed unit suffixes" in warning for warning in summary["warnings"])


def test_figures_csv_comparison_after_collect_generates_png_svg_and_spec(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0

    result = invoke(tmp_path, ["figures", task_id])

    assert result.exit_code == 0
    assert "Generated" in result.output
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    figures_dir = task_path / "figures"
    png_files = sorted(figures_dir.glob("*.png"))
    svg_files = sorted(figures_dir.glob("*.svg"))

    assert png_files
    assert svg_files
    assert len(png_files) == len(svg_files)
    for path in [*png_files, *svg_files]:
        assert_non_empty_file(path)
        assert path.suffix in {".png", ".svg"}

    spec_path = figures_dir / "figure-spec.yaml"
    summary_path = figures_dir / "figure-summary.json"
    assert_non_empty_file(spec_path)
    assert_non_empty_file(summary_path)

    spec = yaml.safe_load(spec_path.read_text(encoding="utf-8"))
    assert spec["source_metrics"].endswith("metrics/normalized_metrics.csv")
    assert spec["figures"]
    assert all(item["chart_type"] == "line" for item in spec["figures"])
    assert len(spec["figures"]) <= 2

    summary = json.loads(summary_path.read_text(encoding="utf-8"))
    assert summary["figure_count"] == len(png_files)

    manifest = read_manifest(tmp_path, task_id)
    artifact_types = {artifact["artifact_id"]: artifact["type"] for artifact in manifest["artifacts"]}
    assert any(kind == "figure" for kind in artifact_types.values())
    assert artifact_types["figures_spec"] == "config"
    assert artifact_types["figures_summary"] == "config"


def test_figures_algorithm_benchmark_after_collect_generates_bar_chart(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(
        invoke(tmp_path, ["ingest", "examples/algorithm-benchmark/results.json"]).output
    )
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0

    result = invoke(tmp_path, ["figures", task_id])

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    figures_dir = task_path / "figures"
    assert_non_empty_file(figures_dir / "bar_runtime_ms_by_algorithm.png")
    assert_non_empty_file(figures_dir / "bar_runtime_ms_by_algorithm.svg")
    spec = yaml.safe_load((figures_dir / "figure-spec.yaml").read_text(encoding="utf-8"))
    assert len(spec["figures"]) == 1
    assert spec["figures"][0]["chart_type"] == "bar"
    assert spec["figures"][0]["x"] == "algorithm"
    assert spec["figures"][0]["y"] == "runtime_ms"


def test_figures_before_collect_returns_exit_code_5(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)

    result = invoke(tmp_path, ["figures", task_id])

    assert result.exit_code == 5
    assert "metrics are not ready" in result.output
    assert "collect" in result.output


def test_figures_missing_task_returns_exit_code_3(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0

    result = invoke(tmp_path, ["figures", "task_missing"])

    assert result.exit_code == 3
    assert "was not found" in result.output


def test_figures_after_deleting_sqlite_uses_manifest_and_metrics(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    index_path = tmp_path / ".lab-sidecar" / "index.sqlite"
    index_path.unlink()

    result = invoke(tmp_path, ["figures", task_id])

    assert result.exit_code == 0
    assert index_path.is_file()
    figures_dir = tmp_path / ".lab-sidecar" / "tasks" / task_id / "figures"
    assert any(path.suffix == ".png" and path.stat().st_size > 0 for path in figures_dir.iterdir())


def test_figures_with_spec_generates_requested_line_chart(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    spec_path = tmp_path / "figure.yaml"
    spec_path.write_text(
        "\n".join(
            [
                "figure_id: accuracy_curve",
                "chart_type: line",
                "title: Validation Accuracy over Epoch",
                "x: epoch",
                "y: val_accuracy",
                "group_by: source_file",
                "output:",
                "  - figures/accuracy_curve.png",
                "  - figures/accuracy_curve.svg",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    result = invoke(tmp_path, ["figures", task_id, "--spec", "figure.yaml"])

    assert result.exit_code == 0
    figures_dir = tmp_path / ".lab-sidecar" / "tasks" / task_id / "figures"
    assert_non_empty_file(figures_dir / "accuracy_curve.png")
    assert_non_empty_file(figures_dir / "accuracy_curve.svg")

    generated_spec = yaml.safe_load((figures_dir / "figure-spec.yaml").read_text(encoding="utf-8"))
    assert generated_spec["spec_input_path"] == "figure.yaml"
    assert generated_spec["figures"][0]["figure_id"] == "accuracy_curve"
    assert generated_spec["figures"][0]["group_by"] == "source_file"

    summary = json.loads((figures_dir / "figure-summary.json").read_text(encoding="utf-8"))
    assert summary["metrics_path"].endswith("metrics/normalized_metrics.csv")
    assert summary["spec_path"] == "figure.yaml"
    assert summary["generated_figures"][0]["png_path"].endswith("figures/accuracy_curve.png")
    assert summary["skipped_candidates"] == []
    assert summary["errors"] == []


def test_figures_with_spec_generates_requested_bar_chart(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(
        invoke(tmp_path, ["ingest", "examples/algorithm-benchmark/results.json"]).output
    )
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    spec_path = tmp_path / "bar.yaml"
    spec_path.write_text(
        "\n".join(
            [
                "figure_id: runtime_by_algorithm",
                "chart_type: bar",
                "title: Runtime by Algorithm",
                "x: algorithm",
                "y: runtime_ms",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    result = invoke(tmp_path, ["figures", task_id, "--spec", "bar.yaml"])

    assert result.exit_code == 0
    figures_dir = tmp_path / ".lab-sidecar" / "tasks" / task_id / "figures"
    assert_non_empty_file(figures_dir / "runtime_by_algorithm.png")
    assert_non_empty_file(figures_dir / "runtime_by_algorithm.svg")
    summary = json.loads((figures_dir / "figure-summary.json").read_text(encoding="utf-8"))
    assert summary["generated_figures"][0]["chart_type"] == "bar"
    assert summary["generated_figures"][0]["x"] == "algorithm"
    assert summary["generated_figures"][0]["y"] == "runtime_ms"


def test_figures_missing_spec_returns_exit_code_2(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0

    result = invoke(tmp_path, ["figures", task_id, "--spec", "missing.yaml"])

    assert result.exit_code == 2
    assert "figure spec is invalid" in result.output
    assert "does not exist" in result.output


def test_figures_invalid_spec_yaml_returns_exit_code_2(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    (tmp_path / "bad.yaml").write_text("figure_id: [unterminated\n", encoding="utf-8")

    result = invoke(tmp_path, ["figures", task_id, "--spec", "bad.yaml"])

    assert result.exit_code == 2
    assert "could not be parsed" in result.output


def test_figures_spec_missing_metrics_field_returns_exit_code_5(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    (tmp_path / "missing-field.yaml").write_text(
        "\n".join(
            [
                "figure_id: missing_metric",
                "chart_type: line",
                "title: Missing Metric",
                "x: epoch",
                "y: does_not_exist",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    result = invoke(tmp_path, ["figures", task_id, "--spec", "missing-field.yaml"])

    assert result.exit_code == 5
    assert "does_not_exist" in result.output
    summary_path = tmp_path / ".lab-sidecar" / "tasks" / task_id / "figures" / "figure-summary.json"
    summary = json.loads(summary_path.read_text(encoding="utf-8"))
    assert summary["generated_figures"] == []
    assert summary["errors"]
    assert "does_not_exist" in summary["skipped_candidates"][0]["reason"]


def test_figures_unsupported_explicit_chart_records_bounded_diagnostics_without_fallback(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    (tmp_path / "scatter.yaml").write_text(
        "\n".join(
            [
                "figure_id: accuracy_scatter",
                "chart_type: scatter",
                "title: Accuracy Scatter",
                "x: epoch",
                "y: val_accuracy",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    result = invoke(tmp_path, ["figures", task_id, "--spec", "scatter.yaml"])

    assert result.exit_code == 5
    assert "unsupported chart diagnostics" in result.output.lower()
    assert "chart_type=scatter" in result.output
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    summary = json.loads((task_path / "figures" / "figure-summary.json").read_text(encoding="utf-8"))
    assert summary["generated_figures"] == []
    assert summary["fallback"]["mode"] == "off"
    assert summary["fallback"]["attempted"] is False
    assert summary["fallback"]["status"] == "not_needed"
    assert summary["unsupported_chart_diagnostics"]
    diagnostic = summary["unsupported_chart_diagnostics"][0]
    assert diagnostic["requested_chart_intent"]["chart_type"] == "scatter"
    assert diagnostic["requested_chart_intent"]["x"] == "epoch"
    assert diagnostic["requested_chart_intent"]["y"] == "val_accuracy"
    assert "epoch" in diagnostic["available_fields"]
    assert "val_accuracy" in diagnostic["available_fields"]
    assert "supported types are bar, box, line" in diagnostic["reason"]
    serialized_summary = json.dumps(summary, ensure_ascii=False)
    assert "Best val_accuracy=0.86" not in serialized_summary
    assert "epoch,train_loss,val_loss" not in serialized_summary
    assert not (task_path / "intelligence").exists()
    traceability = read_traceability(tmp_path, task_id)
    assert traceability["figure_lineage"]["unsupported_chart_diagnostics"]
    assert traceability["figure_lineage"]["fallback"]["mode"] == "off"
    assert_traceability_is_bounded(traceability)


def test_figures_unsupported_explicit_chart_writes_bounded_request_for_fallback_mode(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    (tmp_path / "heatmap.yaml").write_text(
        "\n".join(
            [
                "figure_id: confusion_heatmap",
                "chart_type: heatmap",
                "title: Confusion Heatmap",
                "x: epoch",
                "y: val_accuracy",
                "group_by: source_file",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    result = invoke(tmp_path, ["figures", task_id, "--spec", "heatmap.yaml", "--fallback", "bounded"])

    assert result.exit_code == 5
    assert "Fallback:" in result.output
    assert "status: unavailable" in result.output
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    summary = json.loads((task_path / "figures" / "figure-summary.json").read_text(encoding="utf-8"))
    fallback = summary["fallback"]
    assert fallback["mode"] == "bounded"
    assert fallback["attempted"] is True
    assert fallback["status"] == "unavailable"
    assert fallback["worker_run_id"]
    assert fallback["request_path"].endswith("/figure-request.json")
    assert fallback["validator_result_path"].endswith("/validator-result.json")
    worker_run_dir = task_path / fallback["request_path"].split("/", 1)[0] / fallback["request_path"].split("/", 1)[1]
    request_path = task_path / fallback["request_path"]
    validator_path = task_path / fallback["validator_result_path"]
    assert request_path.is_file()
    assert validator_path.is_file()
    assert (task_path / "intelligence" / fallback["worker_run_id"] / "sandbox").is_dir()
    assert not (task_path / "intelligence" / fallback["worker_run_id"] / "worker-request.json").exists()
    assert not (task_path / "intelligence" / fallback["worker_run_id"] / "worker-result.json").exists()
    assert not (task_path / "intelligence" / fallback["worker_run_id"] / "adoption-record.json").exists()

    request = json.loads(request_path.read_text(encoding="utf-8"))
    assert request["task_id"] == task_id
    assert request["task_status"] == "completed"
    assert request["task_mode"] == "ingest"
    assert request["requested_chart_intent"]["chart_type"] == "heatmap"
    assert request["requested_chart_intent"]["group_by"] == "source_file"
    assert request["row_count"] > 0
    assert request["metric_columns"]
    assert "epoch" in request["metric_columns"]
    assert "val_accuracy" in request["metric_columns"]
    assert "source_file" in request["metric_columns"]
    assert isinstance(request["units"], dict)
    assert isinstance(request["groups"], dict)
    assert isinstance(request["field_sources"], dict)
    assert request["collection_diagnostics"]["processed_files"]
    assert request["collection_diagnostics"]["warnings"] == []
    assert request["artifacts"][0]["artifact_id"] == "metrics_normalized_csv"
    serialized_request = json.dumps(request, ensure_ascii=False)
    assert "RAW_ROW_SENTINEL" not in serialized_request
    assert "epoch,train_loss,val_loss" not in serialized_request
    assert "Best val_accuracy=0.86" not in serialized_request
    assert "ppt/presentation.xml" not in serialized_request
    assert '"prompt":' not in serialized_request.lower()
    assert '"response":' not in serialized_request.lower()

    validator = json.loads(validator_path.read_text(encoding="utf-8"))
    assert validator["accepted"] is False
    assert validator["proposal_type"] == "figure"
    assert validator["checks"][0]["name"] == "worker_available"
    assert "figure_fallback_unavailable" in "\n".join(validator["diagnostics"])
    traceability = read_traceability(tmp_path, task_id)
    serialized_traceability = json.dumps(traceability, ensure_ascii=False)
    assert "epoch,train_loss,val_loss" not in serialized_traceability
    assert "Best val_accuracy=0.86" not in serialized_traceability
    assert traceability["figure_lineage"]["fallback"]["status"] == "unavailable"
    assert_traceability_is_bounded(traceability)


def test_figures_fallback_mock_worker_adopts_validated_official_artifact(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    source = tmp_path / "fallback-metrics"
    source.mkdir()
    (source / "metrics.csv").write_text(
        "\n".join(
            [
                "epoch,val_accuracy,secret_payload",
                "1,0.70,ALPHA4_WORKER_RAW_SENTINEL_ROW_001",
                "2,0.82,ALPHA4_WORKER_RAW_SENTINEL_ROW_002",
                "3,0.91,ALPHA4_WORKER_RAW_SENTINEL_ROW_003",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    (tmp_path / "metrics.yaml").write_text(
        "\n".join(
            [
                "sources:",
                "  - fallback-metrics/metrics.csv",
                "fields:",
                "  epoch: epoch",
                "  val_accuracy:",
                "    source: val_accuracy",
                "    unit: ratio",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "fallback-metrics"]).output)
    assert invoke(tmp_path, ["collect", task_id, "--config", "metrics.yaml"]).exit_code == 0
    (tmp_path / "heatmap.yaml").write_text(
        "\n".join(
            [
                "figure_id: alpha4_heatmap",
                "chart_type: heatmap",
                "title: Alpha4 Heatmap",
                "x: epoch",
                "y: val_accuracy",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    result = invoke(
        tmp_path,
        ["figures", task_id, "--spec", "heatmap.yaml", "--fallback", "bounded", "--fallback-worker", "mock"],
    )

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    summary = json.loads((task_path / "figures" / "figure-summary.json").read_text(encoding="utf-8"))
    fallback = summary["fallback"]
    assert fallback["status"] == "adopted"
    assert fallback["worker_type"] == "mock_chart_fallback"
    assert fallback["validation_status"] == "accepted"
    assert fallback["adoption_record_path"].endswith("/adoption-record.json")
    run_dir = task_path / "intelligence" / fallback["worker_run_id"]
    sandbox = run_dir / "sandbox"
    assert (run_dir / "figure-request.json").is_file()
    assert (run_dir / "worker-request.json").is_file()
    assert (run_dir / "worker-result.json").is_file()
    assert (run_dir / "validator-result.json").is_file()
    assert (sandbox / "figure-fallback-proposal.json").is_file()
    assert (sandbox / "alpha4_heatmap-fallback.png").is_file()
    assert (sandbox / "alpha4_heatmap-fallback.svg").is_file()
    assert (run_dir / "adoption-record.json").is_file()
    official_png = task_path / "figures" / "alpha4_heatmap-fallback.png"
    official_svg = task_path / "figures" / "alpha4_heatmap-fallback.svg"
    assert_nonblank_png(official_png)
    assert_non_empty_file(official_svg)

    request = json.loads((run_dir / "figure-request.json").read_text(encoding="utf-8"))
    worker_request = json.loads((run_dir / "worker-request.json").read_text(encoding="utf-8"))
    worker_result = json.loads((run_dir / "worker-result.json").read_text(encoding="utf-8"))
    validator = json.loads((run_dir / "validator-result.json").read_text(encoding="utf-8"))
    proposal = json.loads((sandbox / "figure-fallback-proposal.json").read_text(encoding="utf-8"))
    adoption = json.loads((run_dir / "adoption-record.json").read_text(encoding="utf-8"))
    serialized = json.dumps(
        {
            "request": request,
            "worker_request": worker_request,
            "worker_result": worker_result,
            "validator": validator,
            "proposal": proposal,
            "adoption": adoption,
            "summary": summary,
        },
        ensure_ascii=False,
    )
    assert validator["accepted"] is True
    assert validator["proposal_type"] == "figure_fallback"
    assert proposal["proposal_type"] == "figure_fallback"
    assert proposal["source_metrics_fields"] == ["epoch", "val_accuracy"]
    assert proposal["output_paths"]["png_path"] == "alpha4_heatmap-fallback.png"
    assert worker_request["sandbox_path"].endswith("/sandbox")
    assert summary["figure_count"] == 1
    figure = summary["generated_figures"][0]
    assert figure["source"] == "fallback"
    assert figure["worker_run_id"] == fallback["worker_run_id"]
    assert figure["validation_status"] == "accepted"
    assert figure["png_path"] == f".lab-sidecar/tasks/{task_id}/figures/alpha4_heatmap-fallback.png"
    assert figure["fallback_lineage"]["source_metrics"].endswith("metrics/normalized_metrics.csv")
    assert figure["fallback_lineage"]["fields_used"] == ["epoch", "val_accuracy"]
    assert figure["field_sources"] == {"epoch": ["epoch"], "val_accuracy": ["val_accuracy"]}
    assert adoption["official_artifacts"][:2] == [
        f".lab-sidecar/tasks/{task_id}/figures/alpha4_heatmap-fallback.png",
        f".lab-sidecar/tasks/{task_id}/figures/alpha4_heatmap-fallback.svg",
    ]
    manifest = read_manifest(tmp_path, task_id)
    artifact_ids = {artifact["artifact_id"] for artifact in manifest["artifacts"]}
    assert "figure_alpha4_heatmap_png" in artifact_ids
    assert "figure_alpha4_heatmap_svg" in artifact_ids
    assert "ALPHA4_WORKER_RAW_SENTINEL_ROW" not in serialized
    assert "1,0.70" not in serialized
    assert '"prompt":' not in serialized.lower()
    assert '"response":' not in serialized.lower()
    traceability = read_traceability(tmp_path, task_id)
    assert traceability["figure_lineage"]["fallback"]["status"] == "adopted"
    assert traceability["figure_lineage"]["fallback"]["adoption_record_path"].endswith("/adoption-record.json")
    trace_figure = traceability["figure_lineage"]["figures"][0]
    assert trace_figure["source"] == "fallback"
    assert trace_figure["worker_run_id"] == fallback["worker_run_id"]
    assert trace_figure["fallback_lineage"]["fields_used"] == ["epoch", "val_accuracy"]
    assert "ALPHA4_WORKER_RAW_SENTINEL_ROW" not in json.dumps(traceability, ensure_ascii=False)
    assert_traceability_is_bounded(traceability)


def test_figures_fallback_mock_worker_rejects_malformed_image_without_adoption(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    (tmp_path / "heatmap.yaml").write_text(
        "\n".join(
            [
                "figure_id: malformed_heatmap",
                "chart_type: heatmap",
                "title: Malformed Heatmap",
                "x: epoch",
                "y: val_accuracy",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    result = invoke(
        tmp_path,
        ["figures", task_id, "--spec", "heatmap.yaml", "--fallback", "bounded", "--fallback-worker", "mock-malformed-image"],
    )

    assert result.exit_code == 5
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    summary = json.loads((task_path / "figures" / "figure-summary.json").read_text(encoding="utf-8"))
    fallback = summary["fallback"]
    assert fallback["status"] == "rejected"
    run_dir = task_path / "intelligence" / fallback["worker_run_id"]
    sandbox = run_dir / "sandbox"
    assert (sandbox / "malformed_heatmap-fallback.png").is_file()
    assert not (task_path / "figures" / "malformed_heatmap-fallback.png").exists()
    assert not (task_path / "figures" / "malformed_heatmap-fallback.svg").exists()
    assert not (run_dir / "adoption-record.json").exists()
    validator = json.loads((run_dir / "validator-result.json").read_text(encoding="utf-8"))
    diagnostics = "\n".join(validator["diagnostics"])
    assert validator["accepted"] is False
    assert "PNG could not be parsed" in diagnostics
    manifest = read_manifest(tmp_path, task_id)
    artifact_ids = {artifact["artifact_id"] for artifact in manifest["artifacts"]}
    assert "figure_malformed_heatmap_png" not in artifact_ids
    traceability = read_traceability(tmp_path, task_id)
    assert traceability["figure_lineage"]["figure_count"] == 0
    assert_traceability_is_bounded(traceability)


def test_figures_fallback_mock_worker_rejects_tiny_image_without_adoption(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    (tmp_path / "heatmap.yaml").write_text(
        "\n".join(
            [
                "figure_id: tiny_heatmap",
                "chart_type: heatmap",
                "title: Tiny Heatmap",
                "x: epoch",
                "y: val_accuracy",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    result = invoke(
        tmp_path,
        ["figures", task_id, "--spec", "heatmap.yaml", "--fallback", "bounded", "--fallback-worker", "mock-tiny-image"],
    )

    assert result.exit_code == 5
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    summary = json.loads((task_path / "figures" / "figure-summary.json").read_text(encoding="utf-8"))
    fallback = summary["fallback"]
    assert fallback["status"] == "rejected"
    run_dir = task_path / "intelligence" / fallback["worker_run_id"]
    assert not (run_dir / "adoption-record.json").exists()
    assert not (task_path / "figures" / "tiny_heatmap-fallback.png").exists()
    validator = json.loads((run_dir / "validator-result.json").read_text(encoding="utf-8"))
    diagnostics = "\n".join(validator["diagnostics"])
    assert validator["accepted"] is False
    assert "PNG dimensions are too small" in diagnostics
    assert "SVG dimensions are too small" in diagnostics
    manifest = read_manifest(tmp_path, task_id)
    artifact_ids = {artifact["artifact_id"] for artifact in manifest["artifacts"]}
    assert "figure_tiny_heatmap_png" not in artifact_ids
    traceability = read_traceability(tmp_path, task_id)
    assert traceability["figure_lineage"]["figure_count"] == 0
    assert_traceability_is_bounded(traceability)


def test_figures_fallback_mock_worker_rejects_missing_field_without_adoption(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    (tmp_path / "scatter.yaml").write_text(
        "\n".join(
            [
                "figure_id: missing_field_scatter",
                "chart_type: scatter",
                "title: Missing Field Scatter",
                "x: epoch",
                "y: val_accuracy",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    result = invoke(
        tmp_path,
        ["figures", task_id, "--spec", "scatter.yaml", "--fallback", "bounded", "--fallback-worker", "mock-missing-field"],
    )

    assert result.exit_code == 5
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    summary = json.loads((task_path / "figures" / "figure-summary.json").read_text(encoding="utf-8"))
    fallback = summary["fallback"]
    assert fallback["status"] == "rejected"
    run_dir = task_path / "intelligence" / fallback["worker_run_id"]
    assert not (run_dir / "adoption-record.json").exists()
    assert not (task_path / "figures" / "missing_field_scatter-fallback.png").exists()
    validator = json.loads((run_dir / "validator-result.json").read_text(encoding="utf-8"))
    diagnostics = "\n".join(validator["diagnostics"])
    assert validator["accepted"] is False
    assert "missing_metric_for_validator" in diagnostics
    manifest = read_manifest(tmp_path, task_id)
    artifact_ids = {artifact["artifact_id"] for artifact in manifest["artifacts"]}
    assert "figure_missing_field_scatter_png" not in artifact_ids
    traceability = read_traceability(tmp_path, task_id)
    assert traceability["figure_lineage"]["fallback"]["status"] == "rejected"
    assert_traceability_is_bounded(traceability)


def test_figures_fallback_mock_worker_rejects_sandbox_escape_without_official_writes(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    (tmp_path / "scatter.yaml").write_text(
        "\n".join(
            [
                "figure_id: escape_scatter",
                "chart_type: scatter",
                "title: Escape Scatter",
                "x: epoch",
                "y: val_accuracy",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    result = invoke(
        tmp_path,
        ["figures", task_id, "--spec", "scatter.yaml", "--fallback", "bounded", "--fallback-worker", "mock-escape"],
    )

    assert result.exit_code == 5
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    summary = json.loads((task_path / "figures" / "figure-summary.json").read_text(encoding="utf-8"))
    fallback = summary["fallback"]
    assert fallback["status"] == "rejected"
    run_dir = task_path / "intelligence" / fallback["worker_run_id"]
    sandbox = run_dir / "sandbox"
    assert (sandbox / "figure-fallback-proposal.json").is_file()
    assert not (task_path / "figures" / "escape.png").exists()
    assert not (task_path / "figures" / "escape.svg").exists()
    assert not (run_dir / "adoption-record.json").exists()
    validator = json.loads((run_dir / "validator-result.json").read_text(encoding="utf-8"))
    diagnostics = "\n".join(validator["diagnostics"])
    assert validator["accepted"] is False
    assert "path escapes sandbox" in diagnostics
    traceability = read_traceability(tmp_path, task_id)
    assert traceability["figure_lineage"]["fallback"]["status"] == "rejected"
    assert_traceability_is_bounded(traceability)


def test_figures_supported_deterministic_spec_does_not_create_fallback_worker_run(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    (tmp_path / "line.yaml").write_text(
        "\n".join(
            [
                "figure_id: accuracy_curve_bounded_flag",
                "chart_type: line",
                "title: Validation Accuracy over Epoch",
                "x: epoch",
                "y: val_accuracy",
            ]
        )
        + "\n",
        encoding="utf-8",
    )

    result = invoke(tmp_path, ["figures", task_id, "--spec", "line.yaml", "--fallback", "bounded"])

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    assert_non_empty_file(task_path / "figures" / "accuracy_curve_bounded_flag.png")
    summary = json.loads((task_path / "figures" / "figure-summary.json").read_text(encoding="utf-8"))
    assert summary["generated_figures"][0]["chart_type"] == "line"
    assert summary["fallback"]["mode"] == "bounded"
    assert summary["fallback"]["attempted"] is False
    assert summary["fallback"]["status"] == "not_needed"
    assert summary["unsupported_chart_diagnostics"] == []
    assert not (task_path / "intelligence").exists()


def test_figures_auto_runtime_alias_uses_runtime_ms(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(
        invoke(tmp_path, ["ingest", "examples/algorithm-benchmark/results.json"]).output
    )
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0

    result = invoke(tmp_path, ["figures", task_id])

    assert result.exit_code == 0
    spec_path = tmp_path / ".lab-sidecar" / "tasks" / task_id / "figures" / "figure-spec.yaml"
    spec = yaml.safe_load(spec_path.read_text(encoding="utf-8"))
    assert spec["figures"][0]["y"] == "runtime_ms"


def test_figures_auto_rejects_bar_chart_with_too_many_categories(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    source = tmp_path / "many-categories.csv"
    with source.open("w", newline="", encoding="utf-8") as fh:
        writer = csv.DictWriter(fh, fieldnames=["model", "accuracy"])
        writer.writeheader()
        for index in range(13):
            writer.writerow({"model": f"model_{index}", "accuracy": 0.8 + index / 100})
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "many-categories.csv"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0

    result = invoke(tmp_path, ["figures", task_id])

    assert result.exit_code == 5
    assert "limit is 12" in result.output
    summary = json.loads(
        (
            tmp_path
            / ".lab-sidecar"
            / "tasks"
            / task_id
            / "figures"
            / "figure-summary.json"
        ).read_text(encoding="utf-8")
    )
    assert summary["generated_figures"] == []
    assert summary["skipped_candidates"]
    assert "limit is 12" in summary["skipped_candidates"][0]["reason"]
    traceability = read_traceability(tmp_path, task_id)
    assert traceability["figure_lineage"]["present"] is True
    assert traceability["figure_lineage"]["summary_path"] == "figures/figure-summary.json"
    assert traceability["figure_lineage"]["spec_path"] is None
    assert traceability["figure_lineage"]["figure_count"] == 0


def test_figures_repeated_run_does_not_duplicate_manifest_artifacts(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0

    assert invoke(tmp_path, ["figures", task_id]).exit_code == 0
    assert invoke(tmp_path, ["figures", task_id]).exit_code == 0

    manifest = read_manifest(tmp_path, task_id)
    artifact_ids = [artifact["artifact_id"] for artifact in manifest["artifacts"]]
    assert len(artifact_ids) == len(set(artifact_ids))


def test_figures_with_spec_after_deleting_sqlite_uses_manifest_and_metrics(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    (tmp_path / "figure.yaml").write_text(
        "\n".join(
            [
                "figure_id: accuracy_curve",
                "chart_type: line",
                "title: Validation Accuracy over Epoch",
                "x: epoch",
                "y: val_accuracy",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    index_path = tmp_path / ".lab-sidecar" / "index.sqlite"
    index_path.unlink()

    result = invoke(tmp_path, ["figures", task_id, "--spec", "figure.yaml"])

    assert result.exit_code == 0
    assert index_path.is_file()
    assert_non_empty_file(
        tmp_path
        / ".lab-sidecar"
        / "tasks"
        / task_id
        / "figures"
        / "accuracy_curve.png"
    )


def test_report_completed_ingest_after_collect_and_figures_generates_markdown(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    assert invoke(tmp_path, ["figures", task_id]).exit_code == 0

    result = invoke(tmp_path, ["report", task_id])

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    report_path = task_path / "reports" / "report-fragment.md"
    summary_path = task_path / "reports" / "report-summary.json"
    assert_non_empty_file(report_path)
    assert_non_empty_file(summary_path)

    text = report_path.read_text(encoding="utf-8")
    assert task_id in text
    assert "status" in text
    assert "completed" in text
    assert "指标行数" in text
    assert "val_accuracy" in text
    assert "../figures/" in text
    assert "working_dir" in text
    assert "source_path" in text

    summary = json.loads(summary_path.read_text(encoding="utf-8"))
    assert summary["template"] == "zh-lab"
    assert summary["provenance"]["task_id"] == task_id
    assert summary["metrics"]["row_count"] > 0
    assert summary["figures"]["figure_count"] > 0
    assert summary["claim_traces"]
    assert any(item["claim_id"] == "report.metrics.row_count" for item in summary["claim_traces"])
    assert any(item["claim_id"].startswith("report.metric.val_accuracy.") for item in summary["claim_traces"])
    assert any(item["claim_id"].startswith("report.metric.val_loss.") for item in summary["claim_traces"])
    assert all(item.get("evidence") for item in summary["claim_traces"] if item["claim_type"] == "numeric_summary")
    assert "manifest.json" in summary["source_artifacts"]
    assert "metrics/normalized_metrics.csv" in summary["generated_from"]
    assert "metrics/collection-summary.json" in summary["generated_from"]
    assert "figures/figure-summary.json" in summary["generated_from"]
    assert "stderr.log" in summary["generated_from"]
    assert "raw/source_refs.json" in summary["generated_from"]

    manifest = read_manifest(tmp_path, task_id)
    artifacts = {artifact["artifact_id"]: artifact for artifact in manifest["artifacts"]}
    assert artifacts["report_fragment_md"]["type"] == "report"
    assert artifacts["report_fragment_md"]["path"] == "reports/report-fragment.md"
    assert artifacts["report_summary_json"]["type"] == "config"
    assert artifacts["report_summary_json"]["sha256"]
    assert artifacts["report_summary_json"]["size_bytes"] > 0
    traceability = read_traceability(tmp_path, task_id)
    assert traceability["task_id"] == task_id
    assert traceability["metric_lineage"]["row_count"] == summary["metrics"]["row_count"]
    assert traceability["metric_lineage"]["columns"] == summary["metrics"]["columns"]
    assert traceability["report_lineage"]["present"] is True
    assert traceability["report_lineage"]["claim_trace_count"] == len(summary["claim_traces"])
    assert any(source["sha256"] for source in traceability["sources"])
    assert any(artifact["artifact_id"] == "report_summary_json" and artifact["sha256"] for artifact in traceability["artifacts"])
    assert_traceability_is_bounded(traceability)


def test_report_with_metrics_but_no_figures_generates_with_hint(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0

    result = invoke(tmp_path, ["report", task_id])

    assert result.exit_code == 0
    report_path = tmp_path / ".lab-sidecar" / "tasks" / task_id / "reports" / "report-fragment.md"
    text = report_path.read_text(encoding="utf-8")
    assert "尚未生成图表" in text
    assert f"labsidecar figures {task_id}" in text


def test_report_completed_ingest_without_metrics_returns_exit_code_5(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)

    result = invoke(tmp_path, ["report", task_id])

    assert result.exit_code == 5
    assert "requires collected metrics" in result.output
    assert "collect" in result.output


def test_report_failed_run_generates_failure_report_without_metrics(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    command = f'"{sys.executable}" examples/simple-failure/fail.py'
    task_id = extract_task_id(invoke(tmp_path, ["run", command]).output)

    result = invoke(tmp_path, ["report", task_id])

    assert result.exit_code == 0
    report_path = tmp_path / ".lab-sidecar" / "tasks" / task_id / "reports" / "report-fragment.md"
    text = report_path.read_text(encoding="utf-8")
    assert "失败" in text
    assert "failure_summary" in text
    assert "FileNotFoundError" in text
    assert "stderr.log 尾部" in text
    assert "reproduce/command.txt" in text
    assert "不写成成功实验结果分析" in text
    summary = json.loads((tmp_path / ".lab-sidecar" / "tasks" / task_id / "reports" / "report-summary.json").read_text(encoding="utf-8"))
    claim_ids = {item["claim_id"] for item in summary["claim_traces"]}
    assert "report.diagnostic.failed_status" in claim_ids
    assert "report.metrics.unavailable" in claim_ids
    assert not any(item["claim_type"] == "numeric_summary" for item in summary["claim_traces"])
    traceability = read_traceability(tmp_path, task_id)
    assert "report.diagnostic.failed_status" in {item["claim_id"] for item in traceability["claim_traces"]}
    assert_traceability_is_bounded(traceability)


def test_report_cancelled_task_generates_cancelled_report(tmp_path: Path) -> None:
    script = tmp_path / "cancel_me.py"
    script.write_text(
        "\n".join(
            [
                "import time",
                "print('ready', flush=True)",
                "time.sleep(30)",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    assert invoke(tmp_path, ["init"]).exit_code == 0
    command = f'"{sys.executable}" cancel_me.py'
    task_id = extract_task_id(invoke(tmp_path, ["run", command, "--background"]).output)
    wait_for_output(tmp_path, task_id, "ready")
    assert invoke(tmp_path, ["cancel", task_id]).exit_code == 0

    result = invoke(tmp_path, ["report", task_id])

    assert result.exit_code == 0
    report_path = tmp_path / ".lab-sidecar" / "tasks" / task_id / "reports" / "report-fragment.md"
    text = report_path.read_text(encoding="utf-8")
    assert "取消" in text
    assert "cancelled" in text
    assert "cancellation note" in text
    assert "started_at" in text
    assert "finished_at" in text
    assert "不写成成功实验结果分析" in text
    summary = json.loads((tmp_path / ".lab-sidecar" / "tasks" / task_id / "reports" / "report-summary.json").read_text(encoding="utf-8"))
    claim_ids = {item["claim_id"] for item in summary["claim_traces"]}
    assert "report.diagnostic.cancelled_status" in claim_ids
    assert "report.metrics.unavailable" in claim_ids
    assert not any(item["claim_type"] == "numeric_summary" for item in summary["claim_traces"])


def test_report_invalid_template_returns_exit_code_2(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0

    result = invoke(tmp_path, ["report", task_id, "--template", "bad-template"])

    assert result.exit_code == 2
    assert "template is invalid" in result.output


def test_report_missing_task_returns_exit_code_3(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0

    result = invoke(tmp_path, ["report", "task_missing"])

    assert result.exit_code == 3
    assert "was not found" in result.output


def test_report_after_deleting_sqlite_uses_manifest_metrics_and_figures(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    assert invoke(tmp_path, ["figures", task_id]).exit_code == 0
    index_path = tmp_path / ".lab-sidecar" / "index.sqlite"
    index_path.unlink()

    result = invoke(tmp_path, ["report", task_id, "--template", "en-paper"])

    assert result.exit_code == 0
    assert index_path.is_file()
    report_path = tmp_path / ".lab-sidecar" / "tasks" / task_id / "reports" / "report-fragment.md"
    text = report_path.read_text(encoding="utf-8")
    assert "Experimental Summary Fragment" in text
    assert "../figures/" in text


def test_report_repeated_run_does_not_duplicate_manifest_artifacts(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0

    assert invoke(tmp_path, ["report", task_id]).exit_code == 0
    assert invoke(tmp_path, ["report", task_id, "--template", "zh-summary"]).exit_code == 0

    manifest = read_manifest(tmp_path, task_id)
    artifact_ids = [artifact["artifact_id"] for artifact in manifest["artifacts"]]
    assert len(artifact_ids) == len(set(artifact_ids))
    assert artifact_ids.count("report_fragment_md") == 1
    assert artifact_ids.count("report_summary_json") == 1


def test_slides_after_collect_figures_report_generates_pptx_and_summary(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    assert invoke(tmp_path, ["figures", task_id]).exit_code == 0
    assert invoke(tmp_path, ["report", task_id]).exit_code == 0

    result = invoke(tmp_path, ["slides", task_id])

    assert result.exit_code == 0
    assert "Presentation draft created:" in result.output
    assert "Template: zh-summary" in result.output
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    pptx_path = task_path / "slides" / "presentation-draft.pptx"
    summary_path = task_path / "slides" / "slides-summary.json"
    assert_non_empty_file(pptx_path)
    assert_non_empty_file(summary_path)

    deck = assert_readable_deck(pptx_path)

    summary = json.loads(summary_path.read_text(encoding="utf-8"))
    assert summary["template"] == "zh-summary"
    assert summary["slide_count"] == len(deck.slides)
    assert len(summary["slides"]) == len(deck.slides)
    assert all({"slide_index", "title", "purpose", "source_artifacts", "evidence"}.issubset(slide) for slide in summary["slides"])
    assert all(slide["evidence"] or slide["empty_source_reason"] for slide in summary["slides"])
    assert summary["included_metrics"]["row_count"] > 0
    assert summary["included_metrics"]["numeric"]
    numeric_columns = [item["column"] for item in summary["included_metrics"]["numeric"]]
    assert "val_accuracy" in numeric_columns
    assert numeric_columns[:2] == ["val_accuracy", "val_loss"]
    assert summary["included_figures"]
    assert any(slide["title"] == "关键对比" for slide in summary["slides"])
    comparison = summary["key_comparisons"][0]
    assert comparison["best_item"]["label"] == "model_a"
    assert comparison["baseline_item"]["label"] == "baseline"
    assert summary["claim_traces"]
    assert any(item["claim_id"] == "slides.metrics.row_count" for item in summary["claim_traces"])
    assert any(item["claim_id"] == "slides.metrics.table_preview" for item in summary["claim_traces"])
    assert any(item["claim_id"] == "slides.key_comparison.status" for item in summary["claim_traces"])
    figure_slides = [slide for slide in summary["slides"] if "figure" in slide["purpose"]]
    assert figure_slides
    assert all(any(evidence["path"].endswith(".png") for evidence in slide["evidence"]) for slide in figure_slides)
    table_slides = [
        slide
        for slide in summary["slides"]
        if any(evidence["role"] == "metrics_table_preview" for evidence in slide["evidence"])
    ]
    assert table_slides
    assert any(evidence["role"] == "metrics_table_preview" for slide in table_slides for evidence in slide["evidence"])
    assert "warnings" in summary
    assert summary["qa_checks"]["slide_count"]["passed"] is True
    assert summary["qa_checks"]["empty_slide_check"]["passed"] is True
    assert summary["qa_checks"]["title_check"]["passed"] is True
    assert summary["qa_checks"]["artifact_duplicate_check"]["passed"] is True
    assert summary["metrics"]["row_count"] > 0
    assert summary["figures"]
    assert "reports/report-fragment.md" in summary["source_artifacts"]
    assert summary["generated_from"] == summary["source_artifacts"]
    assert "manifest.json" in summary["generated_from"]
    assert "metrics/normalized_metrics.csv" in summary["generated_from"]
    assert "metrics/collection-summary.json" in summary["generated_from"]
    assert "figures/figure-summary.json" in summary["generated_from"]
    assert "reports/report-fragment.md" in summary["generated_from"]
    assert "raw/source_refs.json" in summary["generated_from"]
    assert "stdout.log" in summary["generated_from"]
    assert "stderr.log" in summary["generated_from"]
    assert not any(line.lower().startswith(("task_id:", "status:", "mode:")) for line in summary["report_excerpt"])
    traceability = read_traceability(tmp_path, task_id)
    assert traceability["slide_lineage"]["present"] is True
    assert traceability["slide_lineage"]["slide_count"] == summary["slide_count"]
    assert any(item["surface"] == "slides" for item in traceability["claim_traces"])
    assert_traceability_is_bounded(traceability)

    deck_text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "model_a" in deck_text
    assert "baseline" in deck_text
    assert "line_val_accuracy_over_epoch: val_accuracy over epoch by model" in deck_text
    assert "type=line" not in deck_text

    manifest = read_manifest(tmp_path, task_id)
    artifacts = {artifact["artifact_id"]: artifact for artifact in manifest["artifacts"]}
    assert artifacts["slides_presentation_draft_pptx"]["type"] == "presentation"
    assert artifacts["slides_presentation_draft_pptx"]["path"] == "slides/presentation-draft.pptx"
    assert artifacts["slides_presentation_draft_pptx"]["sha256"]
    assert artifacts["slides_presentation_draft_pptx"]["size_bytes"] > 0
    assert artifacts["slides_summary_json"]["type"] == "config"
    assert artifacts["slides_summary_json"]["path"] == "slides/slides-summary.json"
    assert artifacts["slides_summary_json"]["sha256"]


def test_slides_completed_without_artifacts_returns_exit_code_5(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)

    result = invoke(tmp_path, ["slides", task_id])

    assert result.exit_code == 5
    assert "slides generation requires metrics, figures, or report artifacts" in result.output
    assert "collect" in result.output


def test_slides_failed_task_generates_diagnostic_ppt(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    command = f'"{sys.executable}" examples/simple-failure/fail.py'
    task_id = extract_task_id(invoke(tmp_path, ["run", command]).output)

    result = invoke(tmp_path, ["slides", task_id, "--template", "en-summary"])

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    pptx_path = task_path / "slides" / "presentation-draft.pptx"
    summary_path = task_path / "slides" / "slides-summary.json"
    assert_non_empty_file(pptx_path)
    assert_non_empty_file(summary_path)
    deck = assert_readable_deck(pptx_path)
    summary = json.loads(summary_path.read_text(encoding="utf-8"))
    assert summary["task_status"] == "failed"
    assert summary["template"] == "en-summary"
    assert "stderr.log" in summary["source_artifacts"]
    assert any("Failure" in slide["title"] or "Failed" in slide["title"] for slide in summary["slides"])
    assert all(slide["evidence"] or slide["empty_source_reason"] for slide in summary["slides"])
    claim_ids = {item["claim_id"] for item in summary["claim_traces"]}
    assert "slides.diagnostic.failed_status" in claim_ids
    assert "slides.metrics.unavailable" in claim_ids
    assert not any(item["claim_type"] == "numeric_summary" for item in summary["claim_traces"])
    deck_text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "FileNotFoundError" in deck_text
    assert "failed" in deck_text
    traceability = read_traceability(tmp_path, task_id)
    assert "slides.diagnostic.failed_status" in {item["claim_id"] for item in traceability["claim_traces"]}
    assert_traceability_is_bounded(traceability)


def test_slides_cancelled_task_generates_cancelled_diagnostic_ppt(tmp_path: Path) -> None:
    script = tmp_path / "cancel_me.py"
    script.write_text(
        "\n".join(
            [
                "import time",
                "print('ready', flush=True)",
                "time.sleep(30)",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    assert invoke(tmp_path, ["init"]).exit_code == 0
    command = f'"{sys.executable}" cancel_me.py'
    task_id = extract_task_id(invoke(tmp_path, ["run", command, "--background"]).output)
    wait_for_output(tmp_path, task_id, "ready")
    assert invoke(tmp_path, ["cancel", task_id]).exit_code == 0

    result = invoke(tmp_path, ["slides", task_id])

    assert result.exit_code == 0
    pptx_path = tmp_path / ".lab-sidecar" / "tasks" / task_id / "slides" / "presentation-draft.pptx"
    summary_path = tmp_path / ".lab-sidecar" / "tasks" / task_id / "slides" / "slides-summary.json"
    deck = assert_readable_deck(pptx_path)
    summary = json.loads(summary_path.read_text(encoding="utf-8"))
    assert summary["task_status"] == "cancelled"
    assert any("取消" in slide["title"] or "Cancellation" in slide["title"] for slide in summary["slides"])
    assert all(slide["evidence"] or slide["empty_source_reason"] for slide in summary["slides"])
    claim_ids = {item["claim_id"] for item in summary["claim_traces"]}
    assert "slides.diagnostic.cancelled_status" in claim_ids
    assert "slides.metrics.unavailable" in claim_ids
    assert not any(item["claim_type"] == "numeric_summary" for item in summary["claim_traces"])
    deck_text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "cancelled" in deck_text
    assert "cancellation" in deck_text.lower()


def test_slides_with_four_figures_creates_multiple_figure_slides(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id

    for index, (metric, group_by) in enumerate(
        [
            ("val_accuracy", "model"),
            ("val_loss", "model"),
            ("val_accuracy", "seed"),
            ("val_loss", "seed"),
        ],
        start=1,
    ):
        spec_path = tmp_path / f"figure_{index}.yaml"
        spec_path.write_text(
            "\n".join(
                [
                    f"figure_id: multi_{index}",
                    "chart_type: line",
                    f"title: Multi Figure {index}",
                    "x: epoch",
                    f"y: {metric}",
                    f"group_by: {group_by}",
                ]
            )
            + "\n",
            encoding="utf-8",
        )
        assert invoke(tmp_path, ["figures", task_id, "--spec", str(spec_path)]).exit_code == 0

    result = invoke(tmp_path, ["slides", task_id])

    assert result.exit_code == 0
    pptx_path = task_path / "slides" / "presentation-draft.pptx"
    deck = assert_readable_deck(pptx_path, min_slides=8, max_slides=9)
    summary = json.loads((task_path / "slides" / "slides-summary.json").read_text(encoding="utf-8"))
    assert len(summary["included_figures"]) == 4
    figure_slides = [slide for slide in summary["slides"] if "figure" in slide["purpose"]]
    assert len(figure_slides) == 2
    assert all(len([source for source in slide["source_artifacts"] if source.endswith(".png")]) <= 2 for slide in figure_slides)
    deck_text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "figure_id" not in deck_text.lower()
    assert "multi_1" in deck_text
    assert "chart_type" not in deck_text.lower()


def test_slides_repeated_run_does_not_duplicate_manifest_artifacts(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    assert invoke(tmp_path, ["figures", task_id]).exit_code == 0

    assert invoke(tmp_path, ["slides", task_id]).exit_code == 0
    assert invoke(tmp_path, ["slides", task_id, "--template", "en-summary"]).exit_code == 0

    manifest = read_manifest(tmp_path, task_id)
    artifact_ids = [artifact["artifact_id"] for artifact in manifest["artifacts"]]
    assert len(artifact_ids) == len(set(artifact_ids))
    assert artifact_ids.count("slides_presentation_draft_pptx") == 1
    assert artifact_ids.count("slides_summary_json") == 1


def test_slides_after_deleting_sqlite_uses_manifest_artifacts(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    assert invoke(tmp_path, ["figures", task_id]).exit_code == 0
    assert invoke(tmp_path, ["report", task_id]).exit_code == 0
    index_path = tmp_path / ".lab-sidecar" / "index.sqlite"
    index_path.unlink()

    result = invoke(tmp_path, ["slides", task_id])

    assert result.exit_code == 0
    assert index_path.is_file()
    pptx_path = tmp_path / ".lab-sidecar" / "tasks" / task_id / "slides" / "presentation-draft.pptx"
    assert_readable_deck(pptx_path)
    traceability = read_traceability(tmp_path, task_id)
    assert traceability["task_id"] == task_id
    assert traceability["metric_lineage"]["present"] is True
    assert traceability["figure_lineage"]["present"] is True
    assert traceability["report_lineage"]["present"] is True
    assert traceability["slide_lineage"]["present"] is True
    assert any(source["sha256"] for source in traceability["sources"])
    assert any(artifact["path"] == "slides/presentation-draft.pptx" and artifact["sha256"] for artifact in traceability["artifacts"])


def test_slides_long_command_and_paths_are_truncated_in_summary(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0

    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    manifest_path = task_path / "manifest.json"
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    long_command = "py -3 train.py " + " ".join(f"--very-long-option-{index}=value-{index}" for index in range(40))
    long_source = str(tmp_path / ("nested-" + "x" * 160) / ("source-" + "y" * 120))
    long_working_dir = str(tmp_path / ("working-" + "z" * 180))
    manifest["command"] = long_command
    manifest["source_path"] = long_source
    manifest["working_dir"] = long_working_dir
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    result = invoke(tmp_path, ["slides", task_id])

    assert result.exit_code == 0
    summary = json.loads((task_path / "slides" / "slides-summary.json").read_text(encoding="utf-8"))
    assert_readable_deck(task_path / "slides" / "presentation-draft.pptx")
    assert summary["full_text_fields"]["command"] == long_command
    assert summary["full_text_fields"]["source_path"] == long_source
    assert summary["full_text_fields"]["working_dir"] == long_working_dir
    truncated_keys = {item["key"] for item in summary["text_truncations"]}
    assert {"command", "source_path", "working_dir"}.issubset(truncated_keys)


def test_slides_failed_task_long_stderr_is_truncated_and_recorded(tmp_path: Path) -> None:
    script = tmp_path / "long_stderr.py"
    script.write_text(
        "\n".join(
            [
                "import sys",
                "for index in range(40):",
                "    print('stderr-line-' + str(index) + '-' + ('x' * 180), file=sys.stderr)",
                "raise SystemExit(1)",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["run", f'"{sys.executable}" long_stderr.py']).output)
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id

    result = invoke(tmp_path, ["slides", task_id])

    assert result.exit_code == 0
    deck = assert_readable_deck(task_path / "slides" / "presentation-draft.pptx")
    summary = json.loads((task_path / "slides" / "slides-summary.json").read_text(encoding="utf-8"))
    deck_text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "... earlier lines truncated ..." in deck_text
    stderr_truncations = [item for item in summary["text_truncations"] if item["key"] == "stderr_tail"]
    assert stderr_truncations
    stderr_truncation = stderr_truncations[0]
    assert "full" not in stderr_truncation
    assert stderr_truncation["full_omitted_reason"] == "full stderr log body omitted from slides summary"
    assert stderr_truncation["omitted_line_count"] == 40
    assert stderr_truncation["omitted_char_count"] > 40 * 180
    assert stderr_truncation["truncated"] is True
    assert "stderr-line-32" in stderr_truncation["display"]
    assert "... truncated ..." in stderr_truncation["display"]
    assert "stderr-line-0" not in json.dumps(summary, ensure_ascii=False)
    for item in summary["text_truncations"]:
        if item["key"].startswith("stderr_tail"):
            assert "full" not in item
            assert "full_omitted_reason" in item


def test_slides_summary_omits_full_stdout_stderr_logs_after_full_cli_path(tmp_path: Path) -> None:
    script = tmp_path / "long_logs_success.py"
    script.write_text(
        "\n".join(
            [
                "import csv",
                "import sys",
                "from pathlib import Path",
                "for index in range(60):",
                "    print(f'RAW_LOG_BODY_MARKER stdout early {index:04d} FORBIDDEN_STDOUT_EARLY_{index:04d} ' + ('x' * 32))",
                "    print(f'RAW_LOG_BODY_MARKER stderr early {index:04d} FORBIDDEN_STDERR_EARLY_{index:04d} ' + ('y' * 32), file=sys.stderr)",
                "rows = [",
                "    {'epoch': 1, 'train_loss': 1.10, 'val_loss': 1.00, 'val_accuracy': 0.60},",
                "    {'epoch': 2, 'train_loss': 0.90, 'val_loss': 0.82, 'val_accuracy': 0.70},",
                "    {'epoch': 3, 'train_loss': 0.70, 'val_loss': 0.64, 'val_accuracy': 0.78},",
                "    {'epoch': 4, 'train_loss': 0.55, 'val_loss': 0.50, 'val_accuracy': 0.84},",
                "    {'epoch': 5, 'train_loss': 0.45, 'val_loss': 0.42, 'val_accuracy': 0.88},",
                "]",
                "with Path('metrics.csv').open('w', newline='', encoding='utf-8') as fh:",
                "    writer = csv.DictWriter(fh, fieldnames=['epoch', 'train_loss', 'val_loss', 'val_accuracy'])",
                "    writer.writeheader()",
                "    writer.writerows(rows)",
                "for index in range(8):",
                "    print(f'VISIBLE_STDOUT_TAIL_{index:04d} metric_ready')",
                "    print(f'VISIBLE_STDERR_TAIL_{index:04d} warning_ready', file=sys.stderr)",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["run", f'"{sys.executable}" long_logs_success.py']).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    assert invoke(tmp_path, ["figures", task_id]).exit_code == 0
    assert invoke(tmp_path, ["report", task_id]).exit_code == 0

    result = invoke(tmp_path, ["slides", task_id])

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    assert_readable_deck(task_path / "slides" / "presentation-draft.pptx")
    summary = json.loads((task_path / "slides" / "slides-summary.json").read_text(encoding="utf-8"))
    serialized_summary = json.dumps(summary, ensure_ascii=False)
    stdout_text = (task_path / "stdout.log").read_text(encoding="utf-8")
    stderr_text = (task_path / "stderr.log").read_text(encoding="utf-8")
    assert stdout_text not in serialized_summary
    assert stderr_text not in serialized_summary
    assert "RAW_LOG_BODY_MARKER" not in serialized_summary
    assert "FORBIDDEN_STDOUT_EARLY_0000" not in serialized_summary
    assert "FORBIDDEN_STDERR_EARLY_0000" not in serialized_summary

    log_truncations = [
        item
        for item in summary["text_truncations"]
        if item["key"].startswith(("stdout_tail", "stderr_tail"))
    ]
    assert {item["key"] for item in log_truncations}.issuperset({"stdout_tail", "stderr_tail"})
    for item in log_truncations:
        assert "full" not in item
        assert item["full_omitted_reason"]
        assert item["omitted_char_count"] > 0
        assert item["omitted_line_count"] > 0
        assert item["truncated"] is True

    stdout_truncation = next(item for item in log_truncations if item["key"] == "stdout_tail")
    stderr_truncation = next(item for item in log_truncations if item["key"] == "stderr_tail")
    assert stdout_truncation["max_lines"] == 8
    assert stderr_truncation["max_lines"] == 8
    assert "... earlier lines truncated ..." in stdout_truncation["display"]
    assert "... earlier lines truncated ..." in stderr_truncation["display"]
    assert "VISIBLE_STDOUT_TAIL_0007" in stdout_truncation["display"]
    assert "VISIBLE_STDERR_TAIL_0007" in stderr_truncation["display"]
    assert "FORBIDDEN_STDOUT_EARLY" not in stdout_truncation["display"]
    assert "FORBIDDEN_STDERR_EARLY" not in stderr_truncation["display"]

    traceability = read_traceability(tmp_path, task_id)
    serialized_traceability = json.dumps(traceability, ensure_ascii=False)
    assert "RAW_LOG_BODY_MARKER" not in serialized_traceability
    assert "FORBIDDEN_STDOUT_EARLY_0000" not in serialized_traceability
    assert "FORBIDDEN_STDERR_EARLY_0000" not in serialized_traceability
    assert_traceability_is_bounded(traceability)


def test_slides_zh_project_template_generates_project_deck(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/project-presentation-pack"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    assert invoke(tmp_path, ["figures", task_id]).exit_code == 0
    assert invoke(tmp_path, ["report", task_id]).exit_code == 0

    result = invoke(tmp_path, ["slides", task_id, "--template", "zh-project"])

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    deck = assert_readable_deck(task_path / "slides" / "presentation-draft.pptx", min_slides=7, max_slides=9)
    summary = json.loads((task_path / "slides" / "slides-summary.json").read_text(encoding="utf-8"))
    assert summary["template"] == "zh-project"
    assert summary["font_family"] == "Microsoft YaHei"
    assert "SimSun" in summary["font_fallbacks"]
    assert summary["project_goal"]["present"] is True
    slide_titles = {slide["title"] for slide in summary["slides"]}
    assert {"项目概览与来源", "关键对比与消融", "结论与复现"}.issubset(slide_titles)
    deck_text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "项目汇报草稿" in deck_text
    assert len(deck.slides) == summary["slide_count"]


def test_slides_caption_missing_fields_uses_unknown_placeholder(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    assert invoke(tmp_path, ["figures", task_id]).exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    figure_summary_path = task_path / "figures" / "figure-summary.json"
    figure_summary = json.loads(figure_summary_path.read_text(encoding="utf-8"))
    figure_summary["warnings"] = ["synthetic warning for missing metadata"]
    figure_summary["skipped_candidates"] = [{"figure_id": "missing_meta", "reason": "metadata intentionally absent"}]
    for item in figure_summary["generated_figures"]:
        item.pop("x", None)
        item.pop("y", None)
        item.pop("group_by", None)
        item.pop("source_metrics", None)
    figure_summary_path.write_text(json.dumps(figure_summary, ensure_ascii=False, indent=2), encoding="utf-8")

    result = invoke(tmp_path, ["slides", task_id])

    assert result.exit_code == 0
    deck = assert_readable_deck(task_path / "slides" / "presentation-draft.pptx")
    summary = json.loads((task_path / "slides" / "slides-summary.json").read_text(encoding="utf-8"))
    deck_text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "x=未自动推断" in deck_text
    assert "y=未自动推断" in deck_text
    assert "group_by=未自动推断" in deck_text
    assert "source=未自动推断" in deck_text
    assert summary["figure_warnings"] == ["synthetic warning for missing metadata"]
    assert summary["figure_skipped_candidates"][0]["figure_id"] == "missing_meta"


def test_slides_metrics_table_truncates_many_rows_and_columns(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    data_dir = tmp_path / "wide-results"
    data_dir.mkdir()
    csv_path = data_dir / "metrics.csv"
    headers = ["model", "variant", "accuracy", "f1", "latency_ms", "loss", "memory_mb", "runtime_s", "extra_a", "extra_b"]
    rows = []
    for index in range(12):
        rows.append(
            {
                "model": f"model_{index}",
                "variant": f"v{index}",
                "accuracy": str(0.70 + index / 100),
                "f1": str(0.60 + index / 100),
                "latency_ms": str(30 + index),
                "loss": str(1.0 - index / 100),
                "memory_mb": str(100 + index),
                "runtime_s": str(5 + index),
                "extra_a": str(index),
                "extra_b": str(index * 2),
            }
        )
    with csv_path.open("w", newline="", encoding="utf-8") as fh:
        writer = csv.DictWriter(fh, fieldnames=headers)
        writer.writeheader()
        writer.writerows(rows)
    task_id = extract_task_id(invoke(tmp_path, ["ingest", str(data_dir)]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0

    result = invoke(tmp_path, ["slides", task_id])

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    deck = assert_readable_deck(task_path / "slides" / "presentation-draft.pptx", min_slides=6, max_slides=8)
    summary = json.loads((task_path / "slides" / "slides-summary.json").read_text(encoding="utf-8"))
    assert summary["metrics_table"]["displayed_row_count"] <= 6
    assert len(summary["metrics_table"]["columns"]) <= 8
    assert summary["metrics_table"]["shown_columns"]
    assert summary["metrics_table"]["hidden_columns"]
    assert "truncated_cells_count" in summary["metrics_table"]
    assert summary["table_truncations"]
    assert summary["table_truncations"][0]["source_metrics"] == "metrics/normalized_metrics.csv"
    assert summary["table_truncations"][0]["shown_columns"]
    assert summary["table_truncations"][0]["hidden_columns"]
    assert "truncated_cells_count" in summary["table_truncations"][0]
    assert summary["qa_checks"]["table_overflow_guard"]["passed"] is True
    deck_text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "指标表格预览" in deck_text


def test_slides_key_comparison_identifies_best_and_baseline_delta(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    data_dir = tmp_path / "comparison-results"
    data_dir.mkdir()
    csv_path = data_dir / "metrics.csv"
    csv_path.write_text(
        "\n".join(
            [
                "model,accuracy,latency_ms",
                "baseline,0.80,30",
                "candidate_a,0.86,35",
                "candidate_b,0.84,25",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    task_id = extract_task_id(invoke(tmp_path, ["ingest", str(data_dir)]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0

    result = invoke(tmp_path, ["slides", task_id, "--template", "zh-project"])

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    deck = assert_readable_deck(task_path / "slides" / "presentation-draft.pptx", min_slides=7, max_slides=9)
    summary = json.loads((task_path / "slides" / "slides-summary.json").read_text(encoding="utf-8"))
    comparison = summary["key_comparisons"][0]
    assert comparison["metric"] == "accuracy"
    assert comparison["direction"] == "higher"
    assert comparison["best_item"]["label"] == "candidate_a"
    assert comparison["baseline_item"]["label"] == "baseline"
    assert comparison["delta"] is not None
    deck_text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "关键对比" in deck_text
    assert "candidate_a" in deck_text


def test_slides_key_comparison_without_baseline_does_not_invent_delta(tmp_path: Path) -> None:
    assert invoke(tmp_path, ["init"]).exit_code == 0
    data_dir = tmp_path / "no-baseline-results"
    data_dir.mkdir()
    (data_dir / "metrics.csv").write_text(
        "\n".join(
            [
                "method,f1,loss",
                "alpha,0.75,0.45",
                "beta,0.82,0.39",
                "gamma,0.79,0.41",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    task_id = extract_task_id(invoke(tmp_path, ["ingest", str(data_dir)]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0

    result = invoke(tmp_path, ["slides", task_id, "--template", "zh-project"])

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    assert_readable_deck(task_path / "slides" / "presentation-draft.pptx", min_slides=7, max_slides=9)
    summary = json.loads((task_path / "slides" / "slides-summary.json").read_text(encoding="utf-8"))
    comparison = summary["key_comparisons"][0]
    assert comparison["best_item"]["label"] == "beta"
    assert comparison["baseline_item"] is None
    assert comparison["delta"] is None


def test_slides_zh_project_summary_contains_key_comparison_metadata(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/project-presentation-pack"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    assert invoke(tmp_path, ["figures", task_id]).exit_code == 0
    assert invoke(tmp_path, ["report", task_id]).exit_code == 0

    result = invoke(tmp_path, ["slides", task_id, "--template", "zh-project"])

    assert result.exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    assert_readable_deck(task_path / "slides" / "presentation-draft.pptx", min_slides=7, max_slides=9)
    summary = json.loads((task_path / "slides" / "slides-summary.json").read_text(encoding="utf-8"))
    assert summary["key_comparisons"]
    comparison = summary["key_comparisons"][0]
    assert comparison["source_metrics"] == "metrics/normalized_metrics.csv"
    assert comparison["best_item"]
    assert any("关键对比" in slide["title"] for slide in summary["slides"])


def test_slides_long_caption_truncation_is_recorded(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/csv-comparison"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    assert invoke(tmp_path, ["figures", task_id]).exit_code == 0
    task_path = tmp_path / ".lab-sidecar" / "tasks" / task_id
    figure_summary_path = task_path / "figures" / "figure-summary.json"
    figure_summary = json.loads(figure_summary_path.read_text(encoding="utf-8"))
    figure_summary["generated_figures"][0]["figure_id"] = "figure_" + ("very_long_" * 30)
    figure_summary["generated_figures"][0]["group_by"] = "group_" + ("very_long_" * 30)
    figure_summary["generated_figures"][0]["source_metrics"] = "metrics/" + ("very_long_path_" * 30) + "normalized_metrics.csv"
    figure_summary_path.write_text(json.dumps(figure_summary, ensure_ascii=False, indent=2), encoding="utf-8")

    result = invoke(tmp_path, ["slides", task_id])

    assert result.exit_code == 0
    assert_readable_deck(task_path / "slides" / "presentation-draft.pptx", min_slides=6, max_slides=8)
    summary = json.loads((task_path / "slides" / "slides-summary.json").read_text(encoding="utf-8"))
    assert summary["caption_truncations"]
    assert "very_long" in summary["caption_truncations"][0]["full"]
    assert summary["qa_checks"]["caption_overflow_guard"]["passed"] is True


def test_slides_repeated_run_after_table_enhancements_does_not_duplicate_manifest_artifacts(tmp_path: Path) -> None:
    copy_examples(tmp_path)
    assert invoke(tmp_path, ["init"]).exit_code == 0
    task_id = extract_task_id(invoke(tmp_path, ["ingest", "examples/project-presentation-pack"]).output)
    assert invoke(tmp_path, ["collect", task_id]).exit_code == 0
    assert invoke(tmp_path, ["figures", task_id]).exit_code == 0

    assert invoke(tmp_path, ["slides", task_id, "--template", "zh-project"]).exit_code == 0
    assert invoke(tmp_path, ["slides", task_id, "--template", "zh-project"]).exit_code == 0

    manifest = read_manifest(tmp_path, task_id)
    artifact_ids = [artifact["artifact_id"] for artifact in manifest["artifacts"]]
    assert artifact_ids.count("slides_presentation_draft_pptx") == 1
    assert artifact_ids.count("slides_summary_json") == 1
