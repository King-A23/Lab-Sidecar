from __future__ import annotations

import json
import re
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from lab_sidecar.core.artifacts import upsert_artifact
from lab_sidecar.core.manifest import load_task, manifest_path, write_manifest
from lab_sidecar.core.models import ArtifactRecord, TaskRecord, TaskStatus
from lab_sidecar.core.paths import resolve_workspace_path
from lab_sidecar.core.traceability import refresh_traceability
from lab_sidecar.storage.sqlite_index import upsert_task


SUPPORTED_SLIDE_TEMPLATES = {"zh-summary", "en-summary", "zh-project"}
UNKNOWN_ZH = "未自动推断"
UNKNOWN_EN = "Not automatically inferred"
ZH_FONT_FAMILY = "Microsoft YaHei"
ZH_FALLBACK_FONTS = ["Microsoft YaHei", "SimSun", "Calibri"]
EN_FONT_FAMILY = "Calibri"
EN_FALLBACK_FONTS = ["Calibri", "Arial"]
MAX_FIGURES = 4
MAX_FIGURES_PER_SLIDE = 2
MAX_NUMERIC_COLUMNS = 6
MAX_KEY_COLUMNS = 10
MAX_REPORT_BULLETS = 5
MAX_LOG_LINES = 8
MAX_LOG_CHARS = 860
MAX_FIELD_DISPLAY_CHARS = 118
MAX_CAPTION_CHARS = 145
MAX_TABLE_COLUMNS = 8
MAX_TABLE_ROWS = 6
MAX_CELL_CHARS = 42
COMPARISON_GROUP_COLUMNS = ["variant", "model", "method", "algorithm", "source_file"]
HIGHER_IS_BETTER_HINTS = ["accuracy", "acc", "f1", "precision", "recall", "auc", "score"]
LOWER_IS_BETTER_HINTS = ["loss", "runtime", "latency", "time", "memory", "error"]


class InvalidSlidesTemplateError(RuntimeError):
    pass


class SlidesArtifactsRequiredError(RuntimeError):
    pass


class SlidesWriteError(RuntimeError):
    pass


@dataclass
class SlidesGenerationResult:
    record: TaskRecord
    template: str
    pptx_path: Path
    summary_path: Path
    summary: dict[str, Any]


@dataclass
class FigureItem:
    figure_id: str
    chart_type: str
    path: Path
    x: str
    y: str
    group_by: str
    source_metrics: str

    def full_caption(self, unknown: str) -> str:
        chart_type = self.chart_type if self.chart_type in {"line", "bar", "box", "table", "unknown"} else "unknown"
        parts = [
            self.figure_id or unknown,
            f"type={chart_type}",
            f"x={self.x or unknown}",
            f"y={self.y or unknown}",
            f"group_by={self.group_by or unknown}",
            f"source={self.source_metrics or unknown}",
        ]
        return " | ".join(parts)

    def display_caption(self, task_path: Path, unknown: str) -> str:
        caption = self.readable_caption(unknown)
        return _short_text(caption, MAX_CAPTION_CHARS)

    def readable_caption(self, unknown: str) -> str:
        if not self.x or not self.y or not self.group_by or not self.source_metrics:
            return self.full_caption(unknown)
        summary = f"{self.y} over {self.x} by {self.group_by}"
        if self.figure_id:
            return f"{self.figure_id}: {summary}"
        return summary

    def to_summary(self, task_path: Path) -> dict[str, str]:
        return {
            "figure_id": self.figure_id,
            "chart_type": self.chart_type,
            "path": _task_relative(self.path, task_path),
            "x": self.x,
            "y": self.y,
            "group_by": self.group_by,
            "source_metrics": self.source_metrics,
        }


@dataclass
class SlideRecord:
    slide_index: int
    title: str
    purpose: str
    source_artifacts: list[str] = field(default_factory=list)
    evidence: list[dict[str, Any]] = field(default_factory=list)
    empty_source_reason: str | None = None

    def to_dict(self) -> dict[str, Any]:
        return {
            "slide_index": self.slide_index,
            "title": self.title,
            "purpose": self.purpose,
            "source_artifacts": self.source_artifacts,
            "evidence": self.evidence,
            "empty_source_reason": self.empty_source_reason,
        }


@dataclass
class DisplayText:
    display: str
    full: str
    truncated: bool = False


class TextTracker:
    def __init__(self) -> None:
        self.truncations: list[dict[str, Any]] = []

    def fit(
        self,
        key: str,
        value: Any,
        limit: int,
        *,
        include_full: bool = True,
        full_omitted_reason: str | None = None,
        omitted_line_count: int | None = None,
    ) -> DisplayText:
        full = _normalize_inline_text(value)
        if len(full) <= limit:
            return DisplayText(display=full, full=full, truncated=False)
        display = full[: limit - 3].rstrip() + "..."
        truncation = {
            "key": key,
            "display": display,
            "limit": limit,
            "truncated": True,
        }
        if include_full:
            truncation["full"] = full
        else:
            truncation["full_omitted_reason"] = full_omitted_reason or "full text omitted from slides summary"
            truncation["omitted_char_count"] = len(full)
            if omitted_line_count is not None:
                truncation["omitted_line_count"] = omitted_line_count
        self.truncations.append(truncation)
        return DisplayText(display=display, full=full, truncated=True)


class SlidesGenerationService:
    def __init__(self, root: Path):
        self.root = root.resolve()

    def generate(self, task_id: str, template: str = "zh-summary") -> SlidesGenerationResult:
        if template not in SUPPORTED_SLIDE_TEMPLATES:
            allowed = ", ".join(sorted(SUPPORTED_SLIDE_TEMPLATES))
            raise InvalidSlidesTemplateError(f"unsupported template '{template}'. Allowed: {allowed}")

        record = load_task(self.root, task_id)
        task_path = resolve_workspace_path(record.paths.task_dir, self.root)
        slides_dir = task_path / "slides"
        pptx_path = slides_dir / "presentation-draft.pptx"
        summary_path = slides_dir / "slides-summary.json"

        context = self._build_context(record, task_path, template)
        if not context["is_diagnostic"] and not context["has_completed_inputs"]:
            raise SlidesArtifactsRequiredError("metrics, figures, or report artifacts are missing")

        slides_dir.mkdir(parents=True, exist_ok=True)
        builder = _PresentationBuilder(context)
        presentation = builder.build()
        context["slides"] = builder.slide_records
        try:
            presentation.save(pptx_path)
            summary = self._build_summary(record, task_path, pptx_path, summary_path, context, presentation)
            summary_path.write_text(json.dumps(summary, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        except OSError as exc:
            raise SlidesWriteError(f"slides files could not be written: {exc}") from exc

        self._upsert_artifacts(record, context["source_artifacts"])
        record.updated_at = _now_iso()
        record = refresh_traceability(self.root, record)
        write_manifest(manifest_path(self.root, task_id), record)
        upsert_task(self.root, record)

        return SlidesGenerationResult(
            record=record,
            template=template,
            pptx_path=pptx_path,
            summary_path=summary_path,
            summary=summary,
        )

    def _build_context(self, record: TaskRecord, task_path: Path, template: str) -> dict[str, Any]:
        text_tracker = TextTracker()
        metrics_path = task_path / "metrics" / "normalized_metrics.csv"
        scenario_summary_path = task_path / "metrics" / "scenario-summary.json"
        report_path = task_path / "reports" / "report-fragment.md"
        figure_summary_path = task_path / "figures" / "figure-summary.json"
        source_refs_path = task_path / "raw" / "source_refs.json"
        stdout_path = resolve_workspace_path(record.paths.stdout, self.root)
        stderr_path = resolve_workspace_path(record.paths.stderr, self.root)
        figure_items, figure_warnings, figure_metadata = self._find_png_figures(record, task_path, figure_summary_path)
        is_diagnostic = record.status in {TaskStatus.FAILED, TaskStatus.CANCELLED}
        metrics_summary = self._metrics_summary(metrics_path, scenario_summary_path)
        unknown = UNKNOWN_ZH if template.startswith("zh") else UNKNOWN_EN
        metrics_table = self._metrics_table_preview(metrics_path, unknown)
        key_comparison = self._key_comparison(metrics_path, unknown)
        caption_truncations = self._caption_truncations(figure_items, unknown)
        report_excerpt = self._report_excerpt(report_path, template, text_tracker)
        stdout_tail = _bounded_log_tail(stdout_path, "stdout_tail", text_tracker)
        stderr_tail = _bounded_log_tail(stderr_path, "stderr_tail", text_tracker)
        project_goal = self._project_goal_excerpt(record, task_path, text_tracker)
        fields = self._display_fields(record, text_tracker)
        source_artifacts = self._source_artifacts(
            task_path=task_path,
            metrics_path=metrics_path,
            scenario_summary_path=scenario_summary_path,
            report_path=report_path,
            figure_summary_path=figure_summary_path,
            source_refs_path=source_refs_path,
            stdout_path=stdout_path,
            stderr_path=stderr_path,
            figure_items=figure_items,
            project_goal_path=project_goal.get("path_abs"),
        )
        warnings = list(figure_warnings)
        warnings.extend(figure_metadata["warnings"])
        if metrics_summary.get("numeric_omitted_count", 0):
            warnings.append(f"omitted {metrics_summary['numeric_omitted_count']} numeric metric column(s) from slides")
        if metrics_summary.get("omitted_key_column_count", 0):
            warnings.append(f"omitted {metrics_summary['omitted_key_column_count']} metric column name(s) from slides")

        return {
            "template": template,
            "language": "zh" if template.startswith("zh") else "en",
            "unknown": unknown,
            "is_diagnostic": is_diagnostic,
            "is_project_template": template == "zh-project",
            "has_completed_inputs": metrics_path.exists() or bool(figure_items) or report_path.exists(),
            "record": record,
            "task_path": task_path,
            "fields": fields,
            "metrics_path": metrics_path if metrics_path.exists() else None,
            "metrics_summary": metrics_summary,
            "metrics_table": metrics_table,
            "table_truncations": _table_truncations(metrics_table),
            "key_comparison": key_comparison,
            "key_comparisons": [key_comparison] if key_comparison.get("present") else [],
            "figure_items": figure_items,
            "figure_warnings": figure_metadata["warnings"],
            "figure_skipped_candidates": figure_metadata["skipped_candidates"],
            "caption_truncations": caption_truncations,
            "figure_summary_path": figure_summary_path if figure_summary_path.exists() else None,
            "report_path": report_path if report_path.exists() else None,
            "report_excerpt": report_excerpt,
            "project_goal": project_goal,
            "source_refs_path": source_refs_path if source_refs_path.exists() else None,
            "stdout_tail": stdout_tail,
            "stderr_tail": stderr_tail,
            "source_artifacts": source_artifacts,
            "warnings": warnings,
            "text_truncations": text_tracker.truncations,
            "font_family": ZH_FONT_FAMILY if template.startswith("zh") else EN_FONT_FAMILY,
            "font_fallbacks": ZH_FALLBACK_FONTS if template.startswith("zh") else EN_FALLBACK_FONTS,
        }

    def _display_fields(self, record: TaskRecord, text_tracker: TextTracker) -> dict[str, Any]:
        raw = {
            "command": record.command or UNKNOWN_ZH,
            "source_path": record.source_path or UNKNOWN_ZH,
            "working_dir": record.working_dir or UNKNOWN_ZH,
            "failure_summary": record.failure_summary or UNKNOWN_ZH,
            "artifact_dir": record.paths.task_dir,
        }
        display = {
            key: text_tracker.fit(key, value, MAX_FIELD_DISPLAY_CHARS).display
            for key, value in raw.items()
        }
        return {"display": display, "full": raw}

    def _metrics_summary(self, metrics_path: Path, scenario_summary_path: Path) -> dict[str, Any]:
        scenario = _slides_scenario_summary(_read_json(scenario_summary_path) if scenario_summary_path.exists() else {})
        base = {
            "present": False,
            "path": "metrics/normalized_metrics.csv",
            "scenario_summary_path": "metrics/scenario-summary.json" if scenario_summary_path.exists() else None,
            "scenario": scenario,
            "row_count": 0,
            "columns": [],
            "key_columns": [],
            "omitted_key_column_count": 0,
            "numeric": [],
            "numeric_omitted_count": 0,
        }
        if not metrics_path.exists():
            return base
        try:
            df = pd.read_csv(metrics_path)
        except (OSError, pd.errors.EmptyDataError):
            return base

        columns = [str(column) for column in df.columns]
        numeric: list[dict[str, Any]] = []
        for column in df.columns:
            values = pd.to_numeric(df[column], errors="coerce")
            count = int(values.notna().sum())
            if count == 0:
                continue
            final_value = _json_float(values.dropna().iloc[-1]) if count else None
            numeric.append(
                {
                    "column": str(column),
                    "count": count,
                    "mean": _json_float(values.mean()),
                    "min": _json_float(values.min()),
                    "max": _json_float(values.max()),
                    "final": final_value,
                }
            )
        return {
            "present": True,
            "path": "metrics/normalized_metrics.csv",
            "scenario_summary_path": "metrics/scenario-summary.json" if scenario_summary_path.exists() else None,
            "scenario": scenario,
            "row_count": int(len(df)),
            "columns": columns,
            "key_columns": columns[:MAX_KEY_COLUMNS],
            "omitted_key_column_count": max(0, len(columns) - MAX_KEY_COLUMNS),
            "numeric": sorted(numeric, key=lambda item: _metric_priority(str(item["column"])))[:MAX_NUMERIC_COLUMNS],
            "numeric_omitted_count": max(0, len(numeric) - MAX_NUMERIC_COLUMNS),
        }

    def _metrics_table_preview(self, metrics_path: Path, unknown: str) -> dict[str, Any]:
        base = {
            "present": False,
            "source_metrics": "metrics/normalized_metrics.csv",
            "columns": [],
            "shown_columns": [],
            "hidden_columns": [],
            "rows": [],
            "displayed_row_count": 0,
            "total_row_count": 0,
            "total_column_count": 0,
            "truncated": False,
            "omitted_row_count": 0,
            "omitted_column_count": 0,
            "truncated_cells_count": 0,
        }
        if not metrics_path.exists():
            return base
        try:
            df = pd.read_csv(metrics_path)
        except (OSError, pd.errors.EmptyDataError):
            return base
        if df.empty:
            return {
                **base,
                "present": True,
                "total_row_count": 0,
                "total_column_count": len(df.columns),
                "columns": [str(column) for column in df.columns[:MAX_TABLE_COLUMNS]],
                "shown_columns": [str(column) for column in df.columns[:MAX_TABLE_COLUMNS]],
                "hidden_columns": [str(column) for column in df.columns[MAX_TABLE_COLUMNS:]],
            }

        selected_columns = _select_table_columns(df)
        preview_df = df[selected_columns].head(MAX_TABLE_ROWS)
        rows: list[dict[str, str]] = []
        truncated_cells_count = 0
        for _, row in preview_df.iterrows():
            display_row: dict[str, str] = {}
            for column in selected_columns:
                full_value = _format_cell_value(row[column], unknown, limit=None)
                display_value = _format_cell_value(row[column], unknown, limit=MAX_CELL_CHARS)
                if display_value != full_value:
                    truncated_cells_count += 1
                display_row[str(column)] = display_value
            rows.append(display_row)
        omitted_rows = max(0, len(df) - len(rows))
        all_columns = [str(column) for column in df.columns]
        hidden_columns = [column for column in all_columns if column not in selected_columns]
        omitted_columns = len(hidden_columns)
        return {
            "present": True,
            "source_metrics": "metrics/normalized_metrics.csv",
            "columns": [str(column) for column in selected_columns],
            "shown_columns": [str(column) for column in selected_columns],
            "hidden_columns": hidden_columns,
            "rows": rows,
            "displayed_row_count": len(rows),
            "total_row_count": int(len(df)),
            "total_column_count": int(len(df.columns)),
            "truncated": bool(omitted_rows or omitted_columns or truncated_cells_count),
            "omitted_row_count": omitted_rows,
            "omitted_column_count": omitted_columns,
            "truncated_cells_count": truncated_cells_count,
        }

    def _key_comparison(self, metrics_path: Path, unknown: str) -> dict[str, Any]:
        empty = {
            "present": False,
            "metric": None,
            "direction": None,
            "group_column": None,
            "best_item": None,
            "baseline_item": None,
            "delta": None,
            "top_items": [],
            "source_metrics": "metrics/normalized_metrics.csv",
            "reason": "metrics unavailable",
        }
        if not metrics_path.exists():
            return empty
        try:
            df = pd.read_csv(metrics_path)
        except (OSError, pd.errors.EmptyDataError):
            return empty
        if df.empty:
            return {**empty, "reason": "metrics empty"}

        group_column = _first_present(set(str(column) for column in df.columns), COMPARISON_GROUP_COLUMNS)
        if not group_column:
            return {**empty, "reason": "comparison group column not found"}
        metric, direction = _select_comparison_metric(df)
        if not metric or not direction:
            return {**empty, "group_column": group_column, "reason": "priority numeric metric not found"}

        grouped = _aggregate_metric_by_group(df, group_column, metric)
        if not grouped:
            return {**empty, "group_column": group_column, "metric": metric, "direction": direction, "reason": "no numeric comparison values"}
        reverse = direction == "higher"
        sorted_items = sorted(grouped, key=lambda item: item["value"], reverse=reverse)
        best = sorted_items[0]
        baseline = _find_baseline_item(grouped)
        delta = None
        if baseline:
            delta = _json_float(best["value"] - baseline["value"])
        top_items = [
            {"label": item["label"], "value": item["value"], "display_value": _format_number(item["value"], unknown)}
            for item in sorted_items[:3]
        ]
        return {
            "present": True,
            "metric": metric,
            "direction": direction,
            "group_column": group_column,
            "best_item": {"label": best["label"], "value": best["value"], "display_value": _format_number(best["value"], unknown)},
            "baseline_item": {"label": baseline["label"], "value": baseline["value"], "display_value": _format_number(baseline["value"], unknown)} if baseline else None,
            "delta": delta,
            "display_delta": _format_delta(delta, unknown) if delta is not None else None,
            "top_items": top_items,
            "source_metrics": "metrics/normalized_metrics.csv",
            "reason": None,
        }

    def _caption_truncations(self, figure_items: list[FigureItem], unknown: str) -> list[dict[str, Any]]:
        truncations: list[dict[str, Any]] = []
        for item in figure_items:
            full = item.readable_caption(unknown)
            display = _short_text(full, MAX_CAPTION_CHARS)
            if display != full:
                truncations.append(
                    {
                        "figure_id": item.figure_id,
                        "display": display,
                        "full": full,
                        "limit": MAX_CAPTION_CHARS,
                    }
                )
        return truncations

    def _report_excerpt(self, report_path: Path, template: str, text_tracker: TextTracker) -> list[str]:
        if not report_path.exists():
            return []
        text = report_path.read_text(encoding="utf-8", errors="replace")
        lines: list[str] = []
        for raw_line in text.splitlines():
            line = _clean_markdown_line(raw_line)
            if not line:
                continue
            if _is_low_value_report_excerpt(line):
                continue
            lines.append(text_tracker.fit(f"report_excerpt[{len(lines)}]", line, 125).display)
            if len(lines) >= MAX_REPORT_BULLETS:
                break
        return lines or [UNKNOWN_ZH if template.startswith("zh") else UNKNOWN_EN]

    def _project_goal_excerpt(self, record: TaskRecord, task_path: Path, text_tracker: TextTracker) -> dict[str, Any]:
        source_refs_path = task_path / "raw" / "source_refs.json"
        goal_path: Path | None = None
        if source_refs_path.exists():
            refs = _read_json(source_refs_path)
            candidates = refs.get("candidate_files") or []
            source_path = refs.get("source_path")
            possible_paths = [*candidates]
            if source_path:
                source = _resolve_task_or_workspace_path(str(source_path), self.root, task_path)
                if source.is_dir():
                    possible_paths.append((source / "project_goal.md").as_posix())
            for path_text in possible_paths:
                path = _resolve_task_or_workspace_path(str(path_text), self.root, task_path)
                if path.name.lower() == "project_goal.md" and path.exists():
                    goal_path = path
                    break
        if not goal_path:
            return {"present": False, "path": None, "path_abs": None, "excerpt": UNKNOWN_ZH, "full": UNKNOWN_ZH}

        full = goal_path.read_text(encoding="utf-8", errors="replace")
        clean_lines = [_clean_markdown_line(line) for line in full.splitlines()]
        excerpt = " ".join(line for line in clean_lines if line)
        display = text_tracker.fit("project_goal", excerpt or UNKNOWN_ZH, 220)
        return {
            "present": True,
            "path": _portable_source_path(goal_path, self.root, task_path),
            "path_abs": goal_path,
            "excerpt": display.display,
            "full": excerpt,
        }

    def _find_png_figures(
        self,
        record: TaskRecord,
        task_path: Path,
        figure_summary_path: Path,
    ) -> tuple[list[FigureItem], list[str], dict[str, list[Any]]]:
        items: list[FigureItem] = []
        warnings: list[str] = []
        figure_metadata: dict[str, list[Any]] = {"warnings": [], "skipped_candidates": []}
        if figure_summary_path.exists():
            data = _read_json(figure_summary_path)
            figure_metadata["warnings"] = list(data.get("warnings") or [])
            figure_metadata["skipped_candidates"] = list(data.get("skipped_candidates") or [])
            for raw_item in data.get("generated_figures") or data.get("figures") or []:
                path_text = raw_item.get("png_path") or raw_item.get("png")
                if not path_text:
                    continue
                path = _resolve_task_or_workspace_path(path_text, self.root, task_path)
                if path.exists() and path.suffix.lower() == ".png":
                    chart_type = str(raw_item.get("chart_type") or "unknown")
                    if chart_type not in {"line", "bar", "box", "table"}:
                        chart_type = "unknown"
                    items.append(
                        FigureItem(
                            figure_id=str(raw_item.get("figure_id") or path.stem),
                            chart_type=chart_type,
                            path=path,
                            x=str(raw_item.get("x") or ""),
                            y=str(raw_item.get("y") or ""),
                            group_by=str(raw_item.get("group_by") or ""),
                            source_metrics=str(raw_item.get("source_metrics") or raw_item.get("metrics_path") or ""),
                        )
                    )

        known_paths = {item.path.resolve() for item in items}
        for artifact in record.artifacts:
            if artifact.type != "figure" or not artifact.path.lower().endswith(".png"):
                continue
            path = _resolve_task_or_workspace_path(artifact.path, self.root, task_path)
            if path.exists() and path.resolve() not in known_paths:
                items.append(FigureItem(path.stem, "unknown", path, "", "", "", ""))
                known_paths.add(path.resolve())

        figures_dir = task_path / "figures"
        if figures_dir.exists():
            for path in sorted(figures_dir.glob("*.png")):
                if path.resolve() not in known_paths:
                    items.append(FigureItem(path.stem, "unknown", path, "", "", "", ""))
                    known_paths.add(path.resolve())

        if len(items) > MAX_FIGURES:
            warnings.append(f"limited figures to first {MAX_FIGURES} PNG artifact(s); omitted {len(items) - MAX_FIGURES}")
        return items[:MAX_FIGURES], warnings, figure_metadata

    def _source_artifacts(
        self,
        task_path: Path,
        metrics_path: Path,
        scenario_summary_path: Path,
        report_path: Path,
        figure_summary_path: Path,
        source_refs_path: Path,
        stdout_path: Path,
        stderr_path: Path,
        figure_items: list[FigureItem],
        project_goal_path: Path | None,
    ) -> list[str]:
        candidates = [
            task_path / "manifest.json",
            metrics_path,
            task_path / "metrics" / "collection-summary.json",
            scenario_summary_path,
            figure_summary_path,
            report_path,
            source_refs_path,
            stdout_path,
            stderr_path,
            task_path / "reproduce" / "command.txt",
            task_path / "reproduce" / "env.json",
            task_path / "reproduce" / "git.json",
            task_path / "reproduce" / "dependencies.json",
            project_goal_path,
            *[item.path for item in figure_items],
        ]
        existing = [path for path in candidates if path is not None and path.exists()]
        return [_portable_source_path(path, self.root, task_path) for path in _unique_paths(existing)]

    def _build_summary(
        self,
        record: TaskRecord,
        task_path: Path,
        pptx_path: Path,
        summary_path: Path,
        context: dict[str, Any],
        presentation: Presentation,
    ) -> dict[str, Any]:
        included_figures = [item.to_summary(task_path) for item in context["figure_items"]]
        included_metrics = {
            "present": context["metrics_summary"]["present"],
            "path": context["metrics_summary"]["path"],
            "scenario_summary_path": context["metrics_summary"].get("scenario_summary_path"),
            "scenario": context["metrics_summary"].get("scenario", {"present": False}),
            "row_count": context["metrics_summary"]["row_count"],
            "key_columns": context["metrics_summary"]["key_columns"],
            "numeric": context["metrics_summary"]["numeric"],
        }
        return {
            "schema_version": "1",
            "task_id": record.task_id,
            "task_status": record.status.value,
            "template": context["template"],
            "font_family": context["font_family"],
            "font_fallbacks": context["font_fallbacks"],
            "generated_at": _now_iso(),
            "pptx_path": _task_relative(pptx_path, task_path),
            "summary_path": _task_relative(summary_path, task_path),
            "generated_from": context["source_artifacts"],
            "slide_count": len(presentation.slides),
            "included_figures": included_figures,
            "included_metrics": included_metrics,
            "warnings": context["warnings"],
            "figure_warnings": context["figure_warnings"],
            "figure_skipped_candidates": context["figure_skipped_candidates"],
            "text_truncations": context["text_truncations"],
            "table_truncations": context["table_truncations"],
            "key_comparisons": context["key_comparisons"],
            "caption_truncations": context["caption_truncations"],
            "slide_evidence": _slide_claim_traces(context.get("slides", [])),
            "claim_traces": _build_slide_claim_traces(context),
            "qa_checks": self._build_qa_checks(record, context, presentation),
            "slides": [slide.to_dict() for slide in context.get("slides", [])],
            "report_excerpt": context["report_excerpt"],
            "project_goal": {
                "present": context["project_goal"]["present"],
                "path": context["project_goal"]["path"],
                "excerpt": context["project_goal"]["excerpt"],
                "full": context["project_goal"]["full"],
            },
            "full_text_fields": context["fields"]["full"],
            "source_artifacts": context["source_artifacts"],
            # Backward-compatible aliases for Phase 4.1 tests/users.
            "slide_titles": [slide.title for slide in context.get("slides", [])],
            "metrics": context["metrics_summary"],
            "metrics_table": context["metrics_table"],
            "figures": [item.to_summary(task_path)["path"] for item in context["figure_items"]],
        }

    def _build_qa_checks(self, record: TaskRecord, context: dict[str, Any], presentation: Presentation) -> dict[str, Any]:
        slide_records: list[SlideRecord] = context.get("slides", [])
        empty_slide_indices: list[int] = []
        title_failures: list[int] = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text and shape.text.strip()
            ]
            content_texts = [text for text in texts if not text.startswith("Lab-Sidecar |")]
            if not content_texts:
                empty_slide_indices.append(index)
        for slide_record in slide_records:
            if not slide_record.title.strip():
                title_failures.append(slide_record.slide_index)

        virtual_artifacts = [item.artifact_id for item in record.artifacts if item.artifact_id not in {"slides_presentation_draft_pptx", "slides_summary_json"}]
        virtual_artifacts.extend(["slides_presentation_draft_pptx", "slides_summary_json"])
        metrics_table = context["metrics_table"]
        caption_displays = [item.display_caption(context["task_path"], context["unknown"]) for item in context["figure_items"]]
        return {
            "slide_count": {
                "value": len(presentation.slides),
                "passed": len(presentation.slides) == len(slide_records),
            },
            "empty_slide_check": {
                "passed": not empty_slide_indices,
                "empty_slide_indices": empty_slide_indices,
            },
            "title_check": {
                "passed": not title_failures and len(slide_records) == len(presentation.slides),
                "missing_title_slide_indices": title_failures,
            },
            "artifact_duplicate_check": {
                "passed": len(virtual_artifacts) == len(set(virtual_artifacts)),
                "artifact_count_after_upsert": len(virtual_artifacts),
            },
            "table_overflow_guard": {
                "passed": len(metrics_table.get("shown_columns") or []) <= MAX_TABLE_COLUMNS and metrics_table.get("displayed_row_count", 0) <= MAX_TABLE_ROWS,
                "shown_columns": metrics_table.get("shown_columns") or [],
                "hidden_columns": metrics_table.get("hidden_columns") or [],
                "truncated_cells_count": metrics_table.get("truncated_cells_count", 0),
            },
            "caption_overflow_guard": {
                "passed": all(len(text) <= MAX_CAPTION_CHARS for text in caption_displays),
                "max_caption_chars": MAX_CAPTION_CHARS,
                "truncated_caption_count": len(context.get("caption_truncations") or []),
            },
        }

    def _upsert_artifacts(self, record: TaskRecord, source_artifacts: list[str]) -> None:
        upsert_artifact(
            record,
            ArtifactRecord(
                artifact_id="slides_presentation_draft_pptx",
                type="presentation",
                path="slides/presentation-draft.pptx",
                description="Static editable PowerPoint draft",
                source_paths=source_artifacts,
            ),
        )
        upsert_artifact(
            record,
            ArtifactRecord(
                artifact_id="slides_summary_json",
                type="config",
                path="slides/slides-summary.json",
                description="Slides generation summary and provenance",
                source_paths=source_artifacts,
            ),
        )


class _PresentationBuilder:
    def __init__(self, context: dict[str, Any]):
        self.context = context
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.slide_records: list[SlideRecord] = []
        self.zh = context["language"] == "zh"
        self.unknown = context["unknown"]
        self.primary = RGBColor(35, 48, 62)
        self.accent = RGBColor(15, 118, 110)
        self.warn = RGBColor(185, 28, 28)
        self.bg = RGBColor(248, 250, 252)
        self.muted = RGBColor(100, 116, 139)
        self.font_family = context["font_family"]
        self.mono_font = "Consolas"

    def build(self) -> Presentation:
        if self.context["is_diagnostic"]:
            self._build_diagnostic()
        else:
            self._build_completed()
        return self.prs

    def _build_completed(self) -> None:
        record: TaskRecord = self.context["record"]
        fields = self.context["fields"]["display"]
        source = fields["command"] if record.command else fields["source_path"]
        self._title_slide(
            self._deck_title(),
            [f"task_id: {record.task_id}", f"{'来源' if self.zh else 'Source'}: {source}"],
            "task overview",
            ["manifest.json"],
        )
        if self.context["is_project_template"]:
            self._build_project_completed()
            return
        self._settings_slide()
        self._metrics_slide()
        self._metrics_table_slide()
        if self._should_show_key_comparison():
            self._key_comparison_slide()
        self._figure_slides()
        self._result_summary_slide()
        self._reproduce_slide()

    def _should_show_key_comparison(self) -> bool:
        comparison = self.context["key_comparison"]
        if not comparison.get("present"):
            return False
        if comparison.get("baseline_item"):
            return True
        top_items = comparison.get("top_items") or []
        return len(top_items) > 1

    def _build_project_completed(self) -> None:
        self._project_overview_slide()
        self._metrics_slide()
        self._metrics_table_slide()
        self._figure_slides()
        self._project_key_results_slide()
        self._project_conclusion_reproduce_slide()

    def _deck_title(self) -> str:
        if self.context["template"] == "zh-project":
            return "项目汇报草稿"
        return "实验结果演示草稿" if self.zh else "Experiment Results Draft"

    def _build_diagnostic(self) -> None:
        record: TaskRecord = self.context["record"]
        if record.status == TaskStatus.CANCELLED:
            title = "取消任务诊断草稿" if self.zh else "Cancelled Task Diagnostics"
            purpose = "cancelled task overview"
        else:
            title = "失败任务诊断草稿" if self.zh else "Failed Task Diagnostics"
            purpose = "failed task overview"
        self._title_slide(title, [f"task_id: {record.task_id}", f"status: {record.status.value}"], purpose, ["manifest.json"], danger=True)
        self._settings_slide()
        if record.status == TaskStatus.CANCELLED:
            self._cancelled_summary_slide()
        else:
            self._failed_summary_slide()
        self._logs_slide()
        self._reproduce_slide()

    def _title_slide(self, title: str, subtitles: list[str], purpose: str, sources: list[str], danger: bool = False) -> None:
        slide = self._blank_slide(title, purpose, sources, dark=True, danger=danger)
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.25), Inches(11.8), Inches(1.2))
        self._set_text(title_box, title, size=34, bold=True, color=RGBColor(255, 255, 255))
        self._status_bar(slide, Inches(0.78), Inches(2.72), Inches(4.5), self.context["record"].status.value, danger=danger)

        y = 3.42
        for subtitle in subtitles[:3]:
            box = slide.shapes.add_textbox(Inches(0.82), Inches(y), Inches(11.6), Inches(0.42))
            self._set_text(box, _short_text(subtitle, 150), size=15, color=RGBColor(232, 236, 241))
            y += 0.5

    def _settings_slide(self) -> None:
        record: TaskRecord = self.context["record"]
        fields = self.context["fields"]["display"]
        title = "实验设置" if self.zh else "Experiment Settings"
        rows = [
            ("mode", record.mode),
            ("status", record.status.value),
            ("command", fields["command"]),
            ("source_path", fields["source_path"]),
            ("working_dir", fields["working_dir"]),
            ("exit_code", record.exit_code if record.exit_code is not None else self.unknown),
        ]
        self._table_slide(title, rows, "manifest settings and provenance", ["manifest.json"], note="来源: manifest.json" if self.zh else "Source: manifest.json")

    def _project_overview_slide(self) -> None:
        title = "项目概览与来源"
        goal = self.context["project_goal"]
        metrics = self.context["metrics_summary"]
        slide = self._blank_slide(title, "summarize project goal, settings, and source artifacts", ["manifest.json", "raw/source_refs.json", goal["path"] or ""])
        self._add_title(slide, title)
        self._callout_row(
            slide,
            [
                ("mode", self.context["record"].mode),
                ("status", self.context["record"].status.value),
                ("rows", str(metrics["row_count"])),
            ],
            y=Inches(1.14),
        )
        self._add_info_blocks(
            slide,
            [
                f"项目目标: {goal['excerpt'] if goal['present'] else self.unknown}",
                f"来源目录: {self.context['fields']['display']['source_path']}",
                f"artifact: {self.context['fields']['display']['artifact_dir']}",
                f"图表/报告: {len(self.context['figure_items'])} figure(s), {'reports/report-fragment.md' if self.context['report_path'] else self.unknown}",
                "单 task 静态草稿；完整 provenance 写入 slides-summary.json。",
            ],
            top=Inches(2.75),
        )

    def _data_sources_slide(self) -> None:
        title = "数据与来源"
        metrics = self.context["metrics_summary"]
        rows = [
            ("source_path", self.context["fields"]["display"]["source_path"]),
            ("metrics_file", metrics["path"] if metrics["present"] else self.unknown),
            ("rows", metrics["row_count"]),
            ("key_columns", ", ".join(metrics["key_columns"]) or self.unknown),
            ("figures", len(self.context["figure_items"])),
            ("report", "reports/report-fragment.md" if self.context["report_path"] else self.unknown),
        ]
        self._table_slide(title, rows, "summarize project source artifacts", ["manifest.json", "raw/source_refs.json", "metrics/normalized_metrics.csv"])

    def _metrics_slide(self) -> None:
        metrics = self.context["metrics_summary"]
        title = "指标摘要" if self.zh else "Metrics Summary"
        slide = self._blank_slide(title, "summarize normalized metrics", ["metrics/normalized_metrics.csv"])
        self._add_title(slide, title)
        if not metrics["present"]:
            self._callout_row(slide, [("rows", "0"), ("metrics", self.unknown)], y=Inches(1.25))
            self._add_bullets(slide, ["metrics/normalized_metrics.csv not found"], Inches(0.9), Inches(3.2), Inches(11.2), Inches(1.0))
            return

        self._callout_row(
            slide,
            [
                ("rows", str(metrics["row_count"])),
                ("columns", str(len(metrics["columns"]))),
                ("numeric", str(len(metrics.get("numeric", [])) + metrics.get("numeric_omitted_count", 0))),
            ],
            y=Inches(1.14),
        )
        key_columns = ", ".join(metrics["key_columns"]) or self.unknown
        if metrics["omitted_key_column_count"]:
            key_columns += f", ... (+{metrics['omitted_key_column_count']})"
        self._caption(slide, f"key columns: {key_columns}", Inches(0.75), Inches(2.62), Inches(11.7))
        scenario = metrics.get("scenario") or {}
        if scenario.get("present"):
            primary_metric = scenario.get("primary_metric") or {}
            self._caption(
                slide,
                f"scenario: {scenario.get('scenario_type') or self.unknown}; primary_metric={primary_metric.get('name') or self.unknown}; no significance inferred",
                Inches(0.75),
                Inches(2.84),
                Inches(11.7),
                height=Inches(0.32),
            )

        numeric = metrics.get("numeric") or []
        if not numeric:
            self._add_bullets(slide, [self.unknown], Inches(0.9), Inches(3.1), Inches(11.0), Inches(1.0))
            return

        table = slide.shapes.add_table(len(numeric) + 1, 5, Inches(0.75), Inches(3.05), Inches(11.8), Inches(3.05)).table
        headers = ["metric", "min", "mean", "max", "final"]
        for index, header in enumerate(headers):
            cell = table.cell(0, index)
            cell.text = header
            self._style_cell(cell, bold=True, fill=RGBColor(226, 232, 240))
        for row_index, item in enumerate(numeric, start=1):
            values = [
                item["column"],
                _format_number(item["min"], self.unknown),
                _format_number(item["mean"], self.unknown),
                _format_number(item["max"], self.unknown),
                _format_number(item["final"], self.unknown),
            ]
            for col_index, value in enumerate(values):
                cell = table.cell(row_index, col_index)
                cell.text = _short_text(value, 34)
                self._style_cell(cell)

    def _metrics_table_slide(self) -> None:
        table_data = self.context["metrics_table"]
        title = "指标表格预览" if self.zh else "Metrics Table Preview"
        slide = self._blank_slide(title, "preview selected rows and columns from normalized metrics", ["metrics/normalized_metrics.csv"])
        self._add_title(slide, title)
        if not table_data["present"] or not table_data["columns"]:
            self._empty_visual(slide, self.unknown, "metrics/normalized_metrics.csv not found")
            return

        self._callout_row(
            slide,
            [
                ("rows", str(table_data["total_row_count"])),
                ("columns", str(table_data["total_column_count"])),
                ("shown", f"{table_data['displayed_row_count']}x{len(table_data['columns'])}"),
            ],
            y=Inches(1.08),
        )
        self._grid_table(
            slide,
            table_data["columns"],
            table_data["rows"],
            Inches(0.72),
            Inches(2.55),
            Inches(11.85),
            Inches(3.7),
            font_size=7.3,
        )
        note = f"source: {table_data['source_metrics']}"
        if table_data["truncated"]:
            note += f" | truncated rows={table_data['omitted_row_count']}, columns={table_data['omitted_column_count']}"
        self._caption(slide, note, Inches(0.75), Inches(6.5), Inches(11.75))

    def _key_comparison_slide(self) -> None:
        comparison = self.context["key_comparison"]
        title = "关键对比" if self.zh else "Key Comparison"
        slide = self._blank_slide(title, "summarize best, baseline, and top metric comparisons", ["metrics/normalized_metrics.csv"])
        self._add_title(slide, title)
        if not comparison.get("present"):
            self._empty_visual(slide, self.unknown, comparison.get("reason") or "comparison unavailable")
            return

        best = comparison["best_item"]
        baseline = comparison.get("baseline_item")
        delta = comparison.get("display_delta") or self.unknown
        self._callout_row(
            slide,
            [
                ("metric", comparison["metric"]),
                ("best", f"{_short_text(best['label'], 24)} = {best['display_value']}"),
                ("delta", delta),
            ],
            y=Inches(1.08),
        )
        rows = [
            ("direction", comparison["direction"]),
            ("group", comparison["group_column"]),
            ("baseline", f"{baseline['label']} = {baseline['display_value']}" if baseline else self.unknown),
            ("source", comparison["source_metrics"]),
        ]
        self._mini_kv_panel(slide, rows, Inches(0.75), Inches(2.55), Inches(4.15), Inches(3.25))
        top_rows = [
            {"rank": str(index), comparison["group_column"]: item["label"], comparison["metric"]: item["display_value"]}
            for index, item in enumerate(comparison["top_items"], start=1)
        ]
        self._grid_table(
            slide,
            ["rank", comparison["group_column"], comparison["metric"]],
            top_rows,
            Inches(5.25),
            Inches(2.55),
            Inches(7.25),
            Inches(2.65),
            font_size=9,
        )
        self._caption(slide, "规则推断，仅使用 metrics/normalized_metrics.csv；无法可靠判断时保留未自动推断。" if self.zh else "Rule-based; uses metrics/normalized_metrics.csv only.", Inches(5.25), Inches(5.55), Inches(7.1))

    def _project_key_results_slide(self) -> None:
        comparison = self.context["key_comparison"]
        title = "关键对比与消融"
        slide = self._blank_slide(title, "combine key comparison and ablation summary", ["metrics/normalized_metrics.csv"])
        self._add_title(slide, title)
        if not comparison.get("present"):
            metrics = self.context["metrics_summary"]
            group_column = _first_present(set(metrics.get("columns") or []), ["variant", "model", "method", "algorithm", "artifact_role"])
            bullets = [
                f"分组字段: {group_column or self.unknown}",
                f"数值指标: {', '.join(item['column'] for item in (metrics.get('numeric') or [])[:4]) or self.unknown}",
                f"对比结论: {self.unknown}",
                "未识别到可安全排序的主指标；请基于原始 metrics 补充解释。",
            ]
            self._add_info_blocks(slide, bullets, top=Inches(1.2))
            return

        best = comparison["best_item"]
        baseline = comparison.get("baseline_item")
        delta = comparison.get("display_delta") or self.unknown
        self._callout_row(
            slide,
            [
                ("metric", comparison["metric"]),
                ("best", f"{_short_text(best['label'], 24)} = {best['display_value']}"),
                ("delta", delta),
            ],
            y=Inches(1.08),
        )
        rows = [
            ("group", comparison["group_column"]),
            ("direction", "越高越好" if comparison["direction"] == "higher" else "越低越好"),
            ("baseline", f"{baseline['label']} = {baseline['display_value']}" if baseline else self.unknown),
            ("source", "metrics/normalized_metrics.csv"),
        ]
        self._mini_kv_panel(slide, rows, Inches(0.75), Inches(2.52), Inches(4.1), Inches(2.85))
        top_rows = [
            {"rank": str(index), comparison["group_column"]: item["label"], comparison["metric"]: item["display_value"]}
            for index, item in enumerate(comparison["top_items"], start=1)
        ]
        self._grid_table(
            slide,
            ["rank", comparison["group_column"], comparison["metric"]],
            top_rows,
            Inches(5.25),
            Inches(2.52),
            Inches(7.25),
            Inches(2.65),
            font_size=8.8,
        )
        self._caption(slide, "baseline 或 delta 不可靠时保持未自动推断；完整字段见 slides-summary.json。", Inches(5.25), Inches(5.5), Inches(7.1))

    def _figure_slides(self) -> None:
        figure_items: list[FigureItem] = self.context["figure_items"]
        title_base = "图表" if self.zh else "Figures"
        if not figure_items:
            sources = ["figures/figure-summary.json"] if self.context["figure_summary_path"] else []
            slide = self._blank_slide(title_base, "show generated figures", sources)
            self._add_title(slide, title_base)
            self._empty_visual(slide, self.unknown, "figures/*.png not found")
            return

        chunks = [figure_items[index : index + MAX_FIGURES_PER_SLIDE] for index in range(0, len(figure_items), MAX_FIGURES_PER_SLIDE)]
        for page_index, chunk in enumerate(chunks, start=1):
            title = title_base if len(chunks) == 1 else f"{title_base} {page_index}/{len(chunks)}"
            sources = [_task_relative(item.path, self.context["task_path"]) for item in chunk]
            if self.context["figure_summary_path"]:
                sources.insert(0, "figures/figure-summary.json")
            slide = self._blank_slide(title, "show generated figure artifacts", sources)
            self._add_title(slide, title)
            if len(chunk) == 1:
                item = chunk[0]
                self._add_picture_fit(slide, item.path, Inches(0.9), Inches(1.25), Inches(11.55), Inches(5.05))
                self._caption(slide, item.display_caption(self.context["task_path"], self.unknown), Inches(0.9), Inches(6.38), Inches(11.4), height=Inches(0.52))
            else:
                x_positions = [Inches(0.68), Inches(6.82)]
                for index, item in enumerate(chunk):
                    self._add_picture_fit(slide, item.path, x_positions[index], Inches(1.32), Inches(5.78), Inches(4.68))
                    self._caption(slide, item.display_caption(self.context["task_path"], self.unknown), x_positions[index], Inches(6.12), Inches(5.76), height=Inches(0.64))

    def _ablation_or_key_results_slide(self) -> None:
        title = "消融/对比摘要"
        metrics = self.context["metrics_summary"]
        comparison = self.context["key_comparison"]
        columns = set(metrics.get("columns") or [])
        source = ["metrics/normalized_metrics.csv"]
        slide = self._blank_slide(title, "summarize variants, models, or methods when available", source)
        self._add_title(slide, title)
        group_column = _first_present(columns, ["variant", "model", "method", "algorithm", "artifact_role"])
        numeric = metrics.get("numeric") or []
        if not group_column or not numeric:
            self._empty_visual(slide, self.unknown, "variant/model/method not found")
            return
        if comparison.get("present"):
            best = comparison["best_item"]
            baseline = comparison.get("baseline_item")
            bullets = [
                f"分组字段: {comparison['group_column']}",
                f"优先指标: {comparison['metric']} ({comparison['direction']})",
                f"最佳项: {best['label']} = {best['display_value']}",
                f"基线项: {baseline['label']} = {baseline['display_value']}" if baseline else f"基线项: {self.unknown}",
                f"delta: {comparison.get('display_delta') or self.unknown}",
            ]
            self._add_info_blocks(slide, bullets, top=Inches(1.2))
            top_rows = [
                {"rank": str(index), comparison["group_column"]: item["label"], comparison["metric"]: item["display_value"]}
                for index, item in enumerate(comparison["top_items"], start=1)
            ]
            self._grid_table(slide, ["rank", comparison["group_column"], comparison["metric"]], top_rows, Inches(0.82), Inches(5.35), Inches(11.15), Inches(1.05), font_size=8.3)
            return
        bullets = [
            f"分组字段: {group_column}",
            f"数值指标: {', '.join(item['column'] for item in numeric[:4])}",
            f"指标行数: {metrics['row_count']}",
            "更细的消融解释需由用户基于原始 artifact 补充。",
        ]
        self._add_info_blocks(slide, bullets, top=Inches(1.35))

    def _result_summary_slide(self) -> None:
        title = "结果摘要" if self.zh else "Result Summary"
        bullets = self.context["report_excerpt"] or [self.unknown]
        if self.context["report_path"]:
            sources = ["reports/report-fragment.md"]
            suffix = "来源: reports/report-fragment.md" if self.zh else "Source: reports/report-fragment.md"
        else:
            sources = []
            suffix = self.unknown
        slide = self._blank_slide(title, "summarize report fragment", sources)
        self._add_title(slide, title)
        self._add_info_blocks(slide, [*bullets, suffix])

    def _project_conclusion_reproduce_slide(self) -> None:
        record: TaskRecord = self.context["record"]
        fields = self.context["fields"]["display"]
        title = "结论与复现"
        sources = ["manifest.json", "metrics/normalized_metrics.csv"]
        if self.context["report_path"]:
            sources.append("reports/report-fragment.md")
        slide = self._blank_slide(title, "combine conclusion and reproducibility details", sources)
        self._add_title(slide, title)
        bullets = (self.context["report_excerpt"] or [self.unknown])[:3]
        self._add_info_blocks(slide, [*bullets, "结论只来自已有 artifact；未自动推断的内容需用户补充。"], top=Inches(1.15))
        rows = [
            ("command/source", fields["command"] if record.command else fields["source_path"]),
            ("artifact_dir", fields["artifact_dir"]),
            ("exit_code", record.exit_code if record.exit_code is not None else self.unknown),
        ]
        self._mini_kv_panel(slide, rows, Inches(0.82), Inches(5.12), Inches(11.15), Inches(1.55))

    def _failed_summary_slide(self) -> None:
        record: TaskRecord = self.context["record"]
        fields = self.context["fields"]["display"]
        title = "失败诊断" if self.zh else "Failure Diagnostics"
        bullets = [
            f"status: {record.status.value}",
            f"exit_code: {record.exit_code if record.exit_code is not None else self.unknown}",
            f"failure_summary: {fields['failure_summary']}",
            "该任务不会被写成成功实验结果。" if self.zh else "This task is not summarized as a successful experiment.",
        ]
        slide = self._blank_slide(title, "diagnose failed task", ["manifest.json", "stderr.log"], danger=True)
        self._add_title(slide, title)
        self._status_bar(slide, Inches(0.75), Inches(1.18), Inches(3.4), record.status.value, danger=True)
        self._add_info_blocks(slide, bullets, top=Inches(2.05))

    def _cancelled_summary_slide(self) -> None:
        record: TaskRecord = self.context["record"]
        title = "取消诊断" if self.zh else "Cancellation Diagnostics"
        note = _first_non_empty(self.context["stderr_tail"]) or self.unknown
        bullets = [
            f"status: {record.status.value}",
            f"started_at: {record.started_at or self.unknown}",
            f"finished_at: {record.finished_at or self.unknown}",
            f"cancellation note: {_short_text(note, 180)}",
            "该任务不会被写成成功实验结果。" if self.zh else "This task is not summarized as a successful experiment.",
        ]
        slide = self._blank_slide(title, "diagnose cancelled task", ["manifest.json", "stderr.log"], danger=True)
        self._add_title(slide, title)
        self._status_bar(slide, Inches(0.75), Inches(1.18), Inches(3.4), record.status.value, danger=True)
        self._add_info_blocks(slide, bullets, top=Inches(2.05))

    def _logs_slide(self) -> None:
        title = "日志尾部" if self.zh else "Log Tail"
        stderr_lines = self.context["stderr_tail"] or [self.unknown]
        stdout_lines = self.context["stdout_tail"] or [self.unknown]
        lines = [
            "stderr.log:",
            *stderr_lines[:10],
            "",
            "stdout.log:",
            *stdout_lines[:7],
        ]
        slide = self._blank_slide(title, "show bounded stdout/stderr tails", ["stderr.log", "stdout.log"])
        self._add_title(slide, title)
        self._text_panel(slide, lines, Inches(0.75), Inches(1.25), Inches(11.85), Inches(5.35), font_size=8)

    def _reproduce_slide(self) -> None:
        record: TaskRecord = self.context["record"]
        fields = self.context["fields"]["display"]
        title = "复现信息" if self.zh else "Reproducibility"
        sources = ["manifest.json"]
        if "reproduce/command.txt" in self.context["source_artifacts"]:
            sources.append("reproduce/command.txt")
        if self.context["source_refs_path"]:
            sources.append("raw/source_refs.json")
        rows = [
            ("command", fields["command"]),
            ("source_path", fields["source_path"]),
            ("working_dir", fields["working_dir"]),
            ("artifact_dir", fields["artifact_dir"]),
            ("source_refs", "raw/source_refs.json" if self.context["source_refs_path"] else self.unknown),
            ("generated_from", ", ".join(self.context["source_artifacts"][:6]) or self.unknown),
        ]
        self._table_slide(
            title,
            rows,
            "record reproducibility details",
            sources,
            note="SQLite 不是唯一事实来源；本草稿基于 task artifact 文件生成。" if self.zh else "SQLite is not the source of truth; this draft uses task artifact files.",
        )

    def _blank_slide(
        self,
        title: str,
        purpose: str,
        sources: list[str],
        dark: bool = False,
        danger: bool = False,
    ):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_records.append(
            SlideRecord(
                len(self.slide_records) + 1,
                title,
                purpose,
                [source for source in sources if source],
                evidence=_evidence_for_slide(
                    purpose=purpose,
                    sources=sources,
                    context=self.context,
                ),
                empty_source_reason=_empty_source_reason(sources, purpose),
            )
        )
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(30, 41, 59) if dark else self.bg
        accent_color = self.warn if danger else self.accent
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.24), Inches(7.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = accent_color
        accent.line.color.rgb = accent_color
        self._footer(slide, danger=danger, dark=dark)
        return slide

    def _add_title(self, slide, title: str) -> None:
        box = slide.shapes.add_textbox(Inches(0.72), Inches(0.36), Inches(11.8), Inches(0.56))
        self._set_text(box, title, size=25, bold=True, color=self.primary)

    def _footer(self, slide, danger: bool = False, dark: bool = False) -> None:
        record: TaskRecord = self.context["record"]
        text = f"Lab-Sidecar | {record.task_id} | {record.status.value}"
        box = slide.shapes.add_textbox(Inches(0.72), Inches(7.08), Inches(11.9), Inches(0.22))
        self._set_text(box, text, size=8, color=RGBColor(203, 213, 225) if dark else self.muted)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(6.98), Inches(11.85), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = self.warn if danger else RGBColor(203, 213, 225)
        line.line.color.rgb = line.fill.fore_color.rgb

    def _status_bar(self, slide, left, top, width, text: str, danger: bool = False) -> None:
        color = self.warn if danger else self.accent
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.42))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = color
        self._set_text(box, text.upper(), size=13, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)

    def _callout_row(self, slide, values: list[tuple[str, str]], y) -> None:
        width = Inches(3.55)
        gap = Inches(0.34)
        x = Inches(0.75)
        for label, value in values[:3]:
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, Inches(1.18))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(236, 253, 245)
            shape.line.color.rgb = RGBColor(153, 246, 228)
            self._set_text(shape, f"{value}\n{label}", size=16, bold=True, color=self.primary, align=PP_ALIGN.CENTER)
            x += width + gap

    def _add_info_blocks(self, slide, bullets: list[str], top=Inches(1.25)) -> None:
        y = top
        for bullet in bullets[:6]:
            block = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.82), y, Inches(11.15), Inches(0.58))
            block.fill.solid()
            block.fill.fore_color.rgb = RGBColor(241, 245, 249)
            block.line.color.rgb = RGBColor(226, 232, 240)
            self._set_text(block, _short_text(bullet, 165), size=12, color=self.primary)
            y += Inches(0.72)

    def _add_bullets(self, slide, bullets: list[str], left, top, width, height, font_size: int = 14) -> None:
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.word_wrap = True
        frame.clear()
        for index, bullet in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = _short_text(bullet, 160)
            paragraph.level = 0
            paragraph.font.size = Pt(font_size)
            paragraph.font.name = self.font_family
            paragraph.font.color.rgb = self.primary

    def _text_panel(self, slide, lines: list[str], left, top, width, height, font_size: int = 9) -> None:
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(15, 23, 42)
        panel.line.color.rgb = RGBColor(51, 65, 85)
        box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.15), width - Inches(0.36), height - Inches(0.3))
        self._set_text(box, "\n".join(_short_text(line, 150) for line in lines), size=font_size, color=RGBColor(226, 232, 240), font_name=self.mono_font)

    def _mini_kv_panel(self, slide, rows: list[tuple[str, Any]], left, top, width, height) -> None:
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(241, 245, 249)
        panel.line.color.rgb = RGBColor(226, 232, 240)
        y = top + Inches(0.18)
        for key, value in rows[:5]:
            box = slide.shapes.add_textbox(left + Inches(0.18), y, width - Inches(0.36), Inches(0.38))
            self._set_text(box, f"{key}: {_short_text(value, 62)}", size=10, color=self.primary)
            y += Inches(0.5)

    def _grid_table(self, slide, columns: list[str], rows: list[dict[str, Any]], left, top, width, height, font_size: float = 8.5) -> None:
        row_count = max(1, len(rows)) + 1
        col_count = max(1, len(columns))
        table = slide.shapes.add_table(row_count, col_count, left, top, width, height).table
        for col_index, column_width in enumerate(_column_widths(columns, rows, int(width))):
            table.columns[col_index].width = column_width
        for col_index, header in enumerate(columns):
            cell = table.cell(0, col_index)
            cell.text = _short_text(header, 24)
            self._style_cell(cell, bold=True, fill=RGBColor(226, 232, 240), font_size=font_size)
        for row_index, row in enumerate(rows or [{columns[0]: self.unknown}], start=1):
            for col_index, column in enumerate(columns):
                cell = table.cell(row_index, col_index)
                cell.text = _short_text(row.get(column, self.unknown), MAX_CELL_CHARS)
                align = PP_ALIGN.RIGHT if _is_numeric_column(column, rows) else PP_ALIGN.LEFT
                self._style_cell(cell, font_size=font_size, align=align)

    def _table_slide(self, title: str, rows: list[tuple[str, Any]], purpose: str, sources: list[str], note: str | None = None) -> None:
        slide = self._blank_slide(title, purpose, sources)
        self._add_title(slide, title)
        table = slide.shapes.add_table(len(rows) + 1, 2, Inches(0.75), Inches(1.16), Inches(11.85), Inches(4.95)).table
        table.columns[0].width = Inches(2.25)
        table.columns[1].width = Inches(9.6)
        for index, header in enumerate(["field", "value"]):
            cell = table.cell(0, index)
            cell.text = header
            self._style_cell(cell, bold=True, fill=RGBColor(226, 232, 240))
        for row_index, (key, value) in enumerate(rows, start=1):
            cells = [table.cell(row_index, 0), table.cell(row_index, 1)]
            cells[0].text = str(key)
            cells[1].text = _short_text(value, 122)
            self._style_cell(cells[0], bold=True)
            self._style_cell(cells[1])
        if note:
            self._caption(slide, note, Inches(0.75), Inches(6.55), Inches(11.8))

    def _style_cell(self, cell, bold: bool = False, fill: RGBColor | None = None, font_size: float = 9.5, align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
        if fill:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
        cell.margin_left = Inches(0.04)
        cell.margin_right = Inches(0.04)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.name = self.font_family
            paragraph.font.bold = bold
            paragraph.font.color.rgb = self.primary
            paragraph.alignment = align

    def _add_picture_fit(self, slide, path: Path, left, top, width, height) -> None:
        from PIL import Image

        with Image.open(path) as image:
            img_width, img_height = image.size
        img_ratio = img_width / img_height if img_height else 1
        box_ratio = width / height
        if img_ratio >= box_ratio:
            draw_width = width
            draw_height = width / img_ratio
            draw_left = left
            draw_top = top + (height - draw_height) / 2
        else:
            draw_height = height
            draw_width = height * img_ratio
            draw_left = left + (width - draw_width) / 2
            draw_top = top
        slide.shapes.add_picture(str(path), draw_left, draw_top, width=draw_width, height=draw_height)

    def _caption(self, slide, text: str, left, top, width, height=Inches(0.35)) -> None:
        box = slide.shapes.add_textbox(left, top, width, height)
        self._set_text(box, _short_text(text, 145), size=8.5, color=self.muted)

    def _empty_visual(self, slide, title: str, detail: str) -> None:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(2.25), Inches(9.2), Inches(2.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(241, 245, 249)
        shape.line.color.rgb = RGBColor(203, 213, 225)
        self._set_text(shape, f"{title}\n{detail}", size=16, bold=True, color=self.primary, align=PP_ALIGN.CENTER)

    def _set_text(
        self,
        shape,
        text: str,
        size: float,
        color: RGBColor,
        bold: bool = False,
        font_name: str | None = None,
        align: PP_ALIGN = PP_ALIGN.LEFT,
    ) -> None:
        frame = shape.text_frame
        frame.word_wrap = True
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = align
        selected_font = font_name or self.font_family
        for run in paragraph.runs:
            run.font.name = selected_font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def _read_json(path: Path) -> dict[str, Any]:
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return {}
    return data if isinstance(data, dict) else {}


def _slides_scenario_summary(scenario: dict[str, Any]) -> dict[str, Any]:
    if not scenario:
        return {"present": False}
    seed_aggregates = scenario.get("seed_aggregates") if isinstance(scenario.get("seed_aggregates"), dict) else {}
    return {
        "present": True,
        "path": "metrics/scenario-summary.json",
        "scenario_type": scenario.get("scenario_type"),
        "primary_metric": scenario.get("primary_metric") if isinstance(scenario.get("primary_metric"), dict) else {},
        "groups": scenario.get("groups") if isinstance(scenario.get("groups"), dict) else {},
        "best_rows": (scenario.get("best_rows") or [])[:3],
        "last_rows": (scenario.get("last_rows") or [])[:3],
        "seed_aggregates": {
            "present": seed_aggregates.get("present", False),
            "metric": seed_aggregates.get("metric"),
            "direction": seed_aggregates.get("direction"),
            "items": (seed_aggregates.get("items") or [])[:3],
            "claim_limit": seed_aggregates.get("claim_limit"),
        },
        "warnings": (scenario.get("warnings") or [])[:5],
        "omitted": scenario.get("omitted") if isinstance(scenario.get("omitted"), dict) else {},
    }


def _evidence_for_slide(purpose: str, sources: list[str], context: dict[str, Any]) -> list[dict[str, Any]]:
    clean_sources = [source for source in sources if source]
    evidence: list[dict[str, Any]] = []
    lower_purpose = purpose.lower()
    for source in clean_sources:
        item: dict[str, Any] = {
            "path": source,
            "artifact_id": _artifact_id_for_source(source),
            "role": _evidence_role(source, lower_purpose),
        }
        if source == "metrics/normalized_metrics.csv":
            metrics = context.get("metrics_summary") or {}
            item.update(
                {
                    "row_count": metrics.get("row_count"),
                    "columns": metrics.get("columns") or [],
                }
            )
        if source == "figures/figure-summary.json":
            item["figure_ids"] = [figure.figure_id for figure in context.get("figure_items", [])]
        if source.endswith(".png"):
            figure = _find_figure_for_source(source, context.get("figure_items", []), context["task_path"])
            if figure:
                item.update(
                    {
                        "figure_id": figure.figure_id,
                        "source_metrics": figure.source_metrics,
                        "columns": [value for value in [figure.x, figure.y, figure.group_by] if value],
                    }
                )
        if source == "reports/report-fragment.md":
            item["excerpt_count"] = len(context.get("report_excerpt") or [])
        if source in {"stdout.log", "stderr.log"}:
            item["body"] = "omitted"
            item["tail_line_count"] = len(context.get("stdout_tail" if source == "stdout.log" else "stderr_tail") or [])
        evidence.append(item)

    if "metrics table" in lower_purpose or "指标表格" in lower_purpose or "preview selected rows" in lower_purpose:
        table_data = context.get("metrics_table") or {}
        evidence.append(
            {
                "artifact_id": "metrics_normalized_csv",
                "path": table_data.get("source_metrics") or "metrics/normalized_metrics.csv",
                "role": "metrics_table_preview",
                "shown_columns": table_data.get("shown_columns") or [],
                "hidden_columns": table_data.get("hidden_columns") or [],
                "displayed_row_count": table_data.get("displayed_row_count", 0),
                "total_row_count": table_data.get("total_row_count", 0),
                "omitted_row_count": table_data.get("omitted_row_count", 0),
                "omitted_column_count": table_data.get("omitted_column_count", 0),
            }
        )
    if "comparison" in lower_purpose or "对比" in lower_purpose:
        comparison = context.get("key_comparison") or {}
        evidence.append(
            {
                "artifact_id": "metrics_normalized_csv",
                "path": comparison.get("source_metrics") or "metrics/normalized_metrics.csv",
                "role": "key_comparison",
                "metric": comparison.get("metric"),
                "group_column": comparison.get("group_column"),
                "baseline_present": comparison.get("baseline_item") is not None,
                "delta_present": comparison.get("delta") is not None,
                "reason": comparison.get("reason"),
            }
        )
    return [item for item in evidence if item.get("path")]


def _build_slide_claim_traces(context: dict[str, Any]) -> list[dict[str, Any]]:
    traces: list[dict[str, Any]] = []
    metrics = context.get("metrics_summary") or {}
    if metrics.get("present"):
        traces.append(
            {
                "claim_id": "slides.metrics.row_count",
                "surface": "slides",
                "claim_type": "metrics_row_count",
                "value": metrics.get("row_count"),
                "evidence": [
                    {
                        "artifact_id": "metrics_normalized_csv",
                        "path": metrics.get("path") or "metrics/normalized_metrics.csv",
                        "row_count": metrics.get("row_count"),
                    }
                ],
            }
        )
        for item in metrics.get("numeric") or []:
            if not isinstance(item, dict) or not item.get("column"):
                continue
            column = str(item["column"])
            for operation in ["mean", "min", "max", "final"]:
                traces.append(
                    {
                        "claim_id": f"slides.metric.{column}.{operation}",
                        "surface": "slides",
                        "claim_type": "numeric_summary",
                        "operation": operation,
                        "field": column,
                        "value": item.get(operation),
                        "evidence": [
                            {
                                "artifact_id": "metrics_normalized_csv",
                                "path": metrics.get("path") or "metrics/normalized_metrics.csv",
                                "columns": [column],
                                "summary_operation": operation,
                                "numeric_count": item.get("count"),
                                "row_count": metrics.get("row_count"),
                            }
                        ],
                    }
                )
        scenario = metrics.get("scenario") if isinstance(metrics.get("scenario"), dict) else {}
        if scenario.get("present"):
            traces.append(
                {
                    "claim_id": "slides.metrics.scenario_summary",
                    "surface": "slides",
                    "claim_type": "scenario_summary",
                    "value": {
                        "scenario_type": scenario.get("scenario_type"),
                        "primary_metric": scenario.get("primary_metric"),
                    },
                    "evidence": [
                        {
                            "artifact_id": "metrics_scenario_summary",
                            "path": metrics.get("scenario_summary_path") or "metrics/scenario-summary.json",
                            "body": "omitted",
                        }
                    ],
                }
            )
    else:
        traces.append(
            {
                "claim_id": "slides.metrics.unavailable",
                "surface": "slides",
                "claim_type": "unknown_or_unavailable",
                "value": None,
                "evidence": [
                    {
                        "artifact_id": "metrics_normalized_csv",
                        "path": "metrics/normalized_metrics.csv",
                        "reason": "metrics artifact not present",
                    }
                ],
            }
        )

    table_data = context.get("metrics_table") or {}
    if table_data.get("present"):
        traces.append(
            {
                "claim_id": "slides.metrics.table_preview",
                "surface": "slides",
                "claim_type": "metrics_table_preview",
                "value": {
                    "displayed_row_count": table_data.get("displayed_row_count", 0),
                    "total_row_count": table_data.get("total_row_count", 0),
                    "shown_columns": table_data.get("shown_columns") or [],
                    "hidden_column_count": len(table_data.get("hidden_columns") or []),
                },
                "evidence": [
                    {
                        "artifact_id": "metrics_normalized_csv",
                        "path": table_data.get("source_metrics") or "metrics/normalized_metrics.csv",
                        "shown_columns": table_data.get("shown_columns") or [],
                        "hidden_columns": table_data.get("hidden_columns") or [],
                        "displayed_row_count": table_data.get("displayed_row_count", 0),
                        "total_row_count": table_data.get("total_row_count", 0),
                    }
                ],
            }
        )

    comparison = context.get("key_comparison") or {}
    traces.append(
        {
            "claim_id": "slides.key_comparison.status",
            "surface": "slides",
            "claim_type": "comparison_status",
            "value": "available" if comparison.get("present") else "unsupported",
            "evidence": [
                {
                    "artifact_id": "metrics_normalized_csv",
                    "path": comparison.get("source_metrics") or "metrics/normalized_metrics.csv",
                    "metric": comparison.get("metric"),
                    "group_column": comparison.get("group_column"),
                    "baseline_present": comparison.get("baseline_item") is not None,
                    "delta_present": comparison.get("delta") is not None,
                    "reason": comparison.get("reason"),
                }
            ],
        }
    )

    if context.get("is_diagnostic"):
        record: TaskRecord = context["record"]
        traces.append(
            {
                "claim_id": f"slides.diagnostic.{record.status.value}_status",
                "surface": "slides",
                "claim_type": "diagnostic_status",
                "value": record.status.value,
                "evidence": [
                    {
                        "artifact_id": "log_stderr",
                        "path": "stderr.log",
                        "tail_line_count": len(context.get("stderr_tail") or []),
                        "body": "omitted",
                    },
                    {
                        "artifact_id": "log_stdout",
                        "path": "stdout.log",
                        "tail_line_count": len(context.get("stdout_tail") or []),
                        "body": "omitted",
                    },
                ],
            }
        )
    return traces


def _slide_claim_traces(slides: list[SlideRecord]) -> list[dict[str, Any]]:
    return [
        {
            "slide_index": slide.slide_index,
            "title": slide.title,
            "purpose": slide.purpose,
            "evidence_count": len(slide.evidence),
            "empty_source_reason": slide.empty_source_reason,
        }
        for slide in slides
    ]


def _find_figure_for_source(source: str, figures: list[FigureItem], task_path: Path) -> FigureItem | None:
    for item in figures:
        if _task_relative(item.path, task_path) == source:
            return item
    return None


def _artifact_id_for_source(source: str) -> str:
    if source == "manifest.json":
        return "manifest_json"
    if source == "metrics/normalized_metrics.csv":
        return "metrics_normalized_csv"
    if source == "metrics/collection-summary.json":
        return "metrics_collection_summary"
    if source == "figures/figure-summary.json":
        return "figures_summary"
    if source == "reports/report-fragment.md":
        return "report_fragment_md"
    if source == "reports/report-summary.json":
        return "report_summary_json"
    if source == "raw/source_refs.json":
        return "raw_source_refs"
    if source == "stdout.log":
        return "log_stdout"
    if source == "stderr.log":
        return "log_stderr"
    if source.startswith("reproduce/"):
        return "reproduce_" + Path(source).stem
    if source.endswith(".png"):
        return f"figure_{Path(source).stem}_png"
    if source.endswith(".svg"):
        return f"figure_{Path(source).stem}_svg"
    return Path(source).as_posix().replace("/", "_").replace(".", "_")


def _evidence_role(source: str, purpose: str) -> str:
    if source.endswith(".png"):
        return "displayed_figure"
    if source == "metrics/normalized_metrics.csv":
        return "normalized_metrics"
    if source == "reports/report-fragment.md":
        return "report_excerpt"
    if source in {"stdout.log", "stderr.log"}:
        return "bounded_log_tail"
    if "reproduce" in source:
        return "reproduce_metadata"
    if "diagnose" in purpose or "诊断" in purpose:
        return "diagnostic_source"
    return "source_artifact"


def _empty_source_reason(sources: list[str], purpose: str) -> str | None:
    if any(sources):
        return None
    return f"no source artifact was available for slide purpose: {purpose}"


def _bounded_log_tail(path: Path, key: str, text_tracker: TextTracker) -> list[str]:
    if not path.exists():
        return []
    full_text = path.read_text(encoding="utf-8", errors="replace")
    all_lines = full_text.splitlines()
    lines = all_lines[-MAX_LOG_LINES:]
    bounded: list[str] = []
    used = 0
    line_count_truncated = len(all_lines) > MAX_LOG_LINES
    char_count_truncated = False
    stream_name = key.removesuffix("_tail")
    omitted_reason = f"full {stream_name} log body omitted from slides summary"
    for index, line in enumerate(lines):
        display = text_tracker.fit(
            f"{key}[{index}]",
            line,
            150,
            include_full=False,
            full_omitted_reason=f"full {stream_name} log line omitted from slides summary",
            omitted_line_count=1,
        )
        text = display.display
        if used + len(text) > MAX_LOG_CHARS:
            bounded.append("... truncated ...")
            char_count_truncated = True
            break
        bounded.append(text)
        used += len(text)
    if line_count_truncated:
        bounded.insert(0, "... earlier lines truncated ...")
    if line_count_truncated or char_count_truncated:
        text_tracker.truncations.append(
            {
                "key": key,
                "display": "\n".join(bounded),
                "limit": MAX_LOG_CHARS,
                "max_lines": MAX_LOG_LINES,
                "truncated": True,
                "full_omitted_reason": omitted_reason,
                "omitted_char_count": len(full_text),
                "omitted_line_count": len(all_lines),
            }
        )
    return bounded


def _clean_markdown_line(line: str) -> str:
    stripped = line.strip()
    if not stripped or stripped in {"```", "```text"}:
        return ""
    stripped = re.sub(r"^#+\s*", "", stripped)
    stripped = re.sub(r"^[-*]\s*", "", stripped)
    stripped = stripped.replace("`", "")
    stripped = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", stripped)
    stripped = stripped.replace("|", " ").strip()
    return re.sub(r"\s+", " ", stripped)


def _is_low_value_report_excerpt(line: str) -> bool:
    normalized = line.strip().lower()
    if not normalized:
        return True
    if normalized in {
        "实验报告片段",
        "experiment report fragment",
        "实验概览",
        "experiment overview",
        "实验设置与来源",
        "experiment settings and source",
        "字段 值",
        "field value",
        "--- ---",
        "---: ---:",
    }:
        return True
    return normalized.startswith(
        (
            "task_id:",
            "task_id ",
            "status:",
            "status ",
            "mode:",
            "mode ",
            "command ",
            "source_path ",
            "working_dir ",
            "created_at ",
            "started_at ",
            "finished_at ",
            "exit_code ",
            "任务 id:",
            "任务id:",
            "状态:",
            "模式:",
        )
    )


def _resolve_task_or_workspace_path(path_text: str, root: Path, task_path: Path) -> Path:
    path = Path(path_text)
    if path.is_absolute():
        return path
    workspace_candidate = (root / path).resolve()
    if workspace_candidate.exists() or path_text.startswith(".lab-sidecar/"):
        return workspace_candidate
    return (task_path / path).resolve()


def _task_relative(path: Path, task_path: Path) -> str:
    return path.resolve().relative_to(task_path.resolve()).as_posix()


def _portable_source_path(path: Path, root: Path, task_path: Path) -> str:
    resolved = path.resolve()
    try:
        return resolved.relative_to(task_path.resolve()).as_posix()
    except ValueError:
        pass
    try:
        return resolved.relative_to(root.resolve()).as_posix()
    except ValueError:
        return resolved.as_posix()


def _unique_paths(paths: list[Path]) -> list[Path]:
    seen: set[Path] = set()
    result: list[Path] = []
    for path in paths:
        resolved = path.resolve()
        if resolved in seen:
            continue
        seen.add(resolved)
        result.append(resolved)
    return result



def _json_float(value: Any) -> float | None:
    if pd.isna(value):
        return None
    return float(value)


def _format_number(value: Any, unknown: str = UNKNOWN_ZH) -> str:
    if value is None:
        return unknown
    if isinstance(value, int):
        return str(value)
    if isinstance(value, float):
        if pd.isna(value):
            return unknown
        if value.is_integer():
            return str(int(value))
        return f"{value:.4f}".rstrip("0").rstrip(".")
    return str(value)


def _format_delta(value: float | None, unknown: str = UNKNOWN_ZH) -> str:
    if value is None:
        return unknown
    sign = "+" if value > 0 else ""
    return f"{sign}{_format_number(value, unknown)}"


def _format_cell_value(value: Any, unknown: str = UNKNOWN_ZH, limit: int | None = MAX_CELL_CHARS) -> str:
    if pd.isna(value):
        return unknown
    if isinstance(value, (int, float)):
        return _format_number(float(value), unknown)
    text = str(value)
    try:
        numeric = float(text)
    except ValueError:
        return _short_text(text, limit) if limit else _normalize_inline_text(text)
    return _format_number(numeric, unknown)


def _column_widths(columns: list[str], rows: list[dict[str, Any]], total_width: int) -> list[int]:
    weights: list[float] = []
    for column in columns:
        lower = column.lower()
        if _is_numeric_column(column, rows):
            weight = 0.75
        elif any(token in lower for token in ["path", "source", "file", "command", "dir"]):
            weight = 1.65
        elif any(token in lower for token in ["variant", "model", "method", "algorithm"]):
            weight = 1.35
        else:
            weight = 1.0
        weights.append(weight)
    total = sum(weights) or 1.0
    widths = [max(int(total_width * weight / total), Inches(0.72)) for weight in weights]
    diff = total_width - sum(widths)
    if widths:
        widths[-1] += diff
    return widths


def _is_numeric_column(column: str, rows: list[dict[str, Any]]) -> bool:
    lower = column.lower()
    if any(token in lower for token in ["accuracy", "acc", "f1", "loss", "latency", "runtime", "time", "memory", "score", "rank", "seed", "epoch", "step"]):
        return True
    values = [str(row.get(column, "")).strip() for row in rows[:4]]
    values = [value for value in values if value and value not in {UNKNOWN_ZH, UNKNOWN_EN}]
    if not values:
        return False
    numeric = 0
    for value in values:
        try:
            float(value)
        except ValueError:
            continue
        numeric += 1
    return numeric == len(values)


def _select_table_columns(df: pd.DataFrame) -> list[str]:
    columns = [str(column) for column in df.columns]
    selected: list[str] = []
    for candidate in [*COMPARISON_GROUP_COLUMNS, "epoch", "step", "seed", "split", "dataset"]:
        if candidate in columns and candidate not in selected:
            selected.append(candidate)
    numeric_columns = [
        str(column)
        for column in df.columns
        if pd.to_numeric(df[column], errors="coerce").notna().sum() > 0 and str(column) not in selected
    ]
    priority_numeric = sorted(numeric_columns, key=_metric_priority)
    for column in priority_numeric:
        if column not in selected:
            selected.append(column)
        if len(selected) >= MAX_TABLE_COLUMNS:
            break
    for column in columns:
        if len(selected) >= min(MAX_TABLE_COLUMNS, max(4, len(columns))):
            break
        if column not in selected:
            selected.append(column)
    return selected[:MAX_TABLE_COLUMNS]


def _metric_priority(column: str) -> tuple[int, str]:
    lower = column.lower()
    for index, hint in enumerate([*HIGHER_IS_BETTER_HINTS, *LOWER_IS_BETTER_HINTS]):
        if hint in lower:
            return (index, lower)
    return (99, lower)


def _select_comparison_metric(df: pd.DataFrame) -> tuple[str | None, str | None]:
    numeric_columns = [str(column) for column in df.columns if pd.to_numeric(df[column], errors="coerce").notna().sum() > 0]
    for hint in HIGHER_IS_BETTER_HINTS:
        for column in numeric_columns:
            if hint in column.lower():
                return column, "higher"
    for hint in LOWER_IS_BETTER_HINTS:
        for column in numeric_columns:
            if hint in column.lower():
                return column, "lower"
    return None, None


def _aggregate_metric_by_group(df: pd.DataFrame, group_column: str, metric: str) -> list[dict[str, Any]]:
    work = df[[group_column, metric]].copy()
    work[metric] = pd.to_numeric(work[metric], errors="coerce")
    work = work.dropna(subset=[metric])
    if work.empty:
        return []
    grouped = work.groupby(group_column, dropna=False)[metric].mean().reset_index()
    items: list[dict[str, Any]] = []
    for _, row in grouped.iterrows():
        label = str(row[group_column])
        items.append({"label": label, "value": _json_float(row[metric])})
    return [item for item in items if item["value"] is not None]


def _find_baseline_item(items: list[dict[str, Any]]) -> dict[str, Any] | None:
    for item in items:
        label = str(item["label"]).lower()
        if "baseline" in label or label in {"base", "control"}:
            return item
    return None


def _table_truncations(table_data: dict[str, Any]) -> list[dict[str, Any]]:
    if not table_data.get("truncated"):
        return []
    return [
        {
            "source_metrics": table_data["source_metrics"],
            "shown_columns": table_data["shown_columns"],
            "hidden_columns": table_data["hidden_columns"],
            "displayed_rows": table_data["displayed_row_count"],
            "total_rows": table_data["total_row_count"],
            "displayed_columns": len(table_data["columns"]),
            "total_columns": table_data["total_column_count"],
            "omitted_row_count": table_data["omitted_row_count"],
            "omitted_column_count": table_data["omitted_column_count"],
            "truncated_cells_count": table_data.get("truncated_cells_count", 0),
        }
    ]


def _short_text(value: Any, limit: int) -> str:
    text = _normalize_inline_text(value)
    if len(text) <= limit:
        return text
    return text[: limit - 3].rstrip() + "..."


def _normalize_inline_text(value: Any) -> str:
    text = str(value).replace("\r", " ").replace("\n", " ")
    return re.sub(r"\s+", " ", text).strip()


def _first_non_empty(lines: list[str]) -> str | None:
    for line in reversed(lines):
        if line.strip():
            return line
    return None


def _first_present(values: set[str], candidates: list[str]) -> str | None:
    for candidate in candidates:
        if candidate in values:
            return candidate
    return None


def _now_iso() -> str:
    return datetime.now(timezone.utc).astimezone().isoformat(timespec="seconds")
