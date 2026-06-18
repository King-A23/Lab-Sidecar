from __future__ import annotations

import csv
from collections import deque
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation

from lab_sidecar.core.manifest import load_task
from lab_sidecar.core.models import TaskRecord
from lab_sidecar.core.paths import resolve_workspace_path
from lab_sidecar.intelligence.bundle import omitted_contract
from lab_sidecar.intelligence.sandbox import is_path_within


MAX_PREVIEW_ROWS = 20
MAX_PREVIEW_LINES = 80
MAX_LOG_TAIL_LINES = 80
SUPPORTED_SUFFIXES = {".csv", ".md", ".markdown", ".png", ".jpg", ".jpeg", ".pptx", ".log", ".txt"}
DENIED_PREVIEW_NAMES = {
    "ai-provider-prompt.json",
    "ai-provider-response.json",
    "prompt.md",
    "response.json",
}


class ArtifactPreviewError(ValueError):
    pass


def preview_artifact(
    root: Path,
    task_id: str,
    artifact_path: str,
    max_rows: int = 20,
    max_lines: int = 40,
) -> dict[str, Any]:
    record = load_task(root, task_id)
    task_path = resolve_workspace_path(record.paths.task_dir, root)
    resolved, logical_path, artifact_type = _resolve_allowed_artifact(root, record, task_path, artifact_path)
    suffix = resolved.suffix.lower()

    if artifact_type == "log" or suffix == ".log":
        preview = _log_preview(resolved, _bounded(max_lines, 1, MAX_LOG_TAIL_LINES))
        preview_type = "log_tail"
    elif suffix == ".csv":
        preview = _csv_preview(resolved, _bounded(max_rows, 1, MAX_PREVIEW_ROWS))
        preview_type = "csv_rows"
    elif suffix in {".md", ".markdown", ".txt"}:
        if artifact_type not in {"report", "log"} and suffix == ".txt":
            raise ArtifactPreviewError("unsupported artifact preview type")
        preview = _text_lines_preview(resolved, _bounded(max_lines, 1, MAX_PREVIEW_LINES))
        preview_type = "markdown_lines" if suffix in {".md", ".markdown"} else "log_tail"
    elif suffix in {".png", ".jpg", ".jpeg"}:
        preview = _image_preview(resolved)
        preview_type = "image_metadata"
    elif suffix == ".pptx":
        preview = _pptx_preview(resolved)
        preview_type = "pptx_metadata"
    else:
        raise ArtifactPreviewError("unsupported artifact preview type")

    return {
        "schema_version": "2.1",
        "task_id": record.task_id,
        "status": record.status.value,
        "artifact": {
            "path": logical_path,
            "type": artifact_type,
            "size_bytes": resolved.stat().st_size if resolved.exists() else None,
        },
        "preview_type": preview_type,
        "preview": preview,
        "risk_flags": [],
        "omitted": {
            **omitted_contract(),
            "complete_artifact": "omitted_by_default",
            "unbounded_preview": "omitted_by_default",
        },
    }


def _resolve_allowed_artifact(
    root: Path,
    record: TaskRecord,
    task_path: Path,
    artifact_path: str,
) -> tuple[Path, str, str]:
    requested = Path(artifact_path)
    if requested.is_absolute():
        try:
            logical = requested.resolve().relative_to(root.resolve()).as_posix()
        except ValueError as exc:
            raise ArtifactPreviewError("artifact path is outside the workspace") from exc
    else:
        logical = artifact_path

    allowed: dict[str, str] = {artifact.path: artifact.type for artifact in record.artifacts}
    allowed[record.paths.stdout] = "log"
    allowed[record.paths.stderr] = "log"
    task_relative_logs = {
        "stdout.log": "log",
        "stderr.log": "log",
    }
    allowed.update({f"{record.paths.task_dir}/{name}": kind for name, kind in task_relative_logs.items()})

    artifact_type = allowed.get(logical)
    if artifact_type is None:
        task_relative = _task_relative_logical(logical, task_path, root)
        artifact_type = allowed.get(task_relative)
        if artifact_type is not None:
            logical = task_relative

    if artifact_type is None:
        raise ArtifactPreviewError("artifact is not registered for this task")
    if _is_denied_preview(logical, artifact_type):
        raise ArtifactPreviewError("artifact preview is not available for raw or worker audit content")

    resolved = _resolve_registered_artifact_path(logical, root, task_path)
    if not is_path_within(resolved, task_path.resolve()):
        raise ArtifactPreviewError("artifact path is outside the task directory")
    if not resolved.exists() or not resolved.is_file():
        raise ArtifactPreviewError("artifact file does not exist")
    if resolved.suffix.lower() not in SUPPORTED_SUFFIXES:
        raise ArtifactPreviewError("unsupported artifact preview type")
    return resolved, logical, artifact_type


def _is_denied_preview(logical: str, artifact_type: str) -> bool:
    path = Path(logical)
    parts = set(path.parts)
    if artifact_type == "raw":
        return True
    return "intelligence" in parts or path.name in DENIED_PREVIEW_NAMES


def _task_relative_logical(logical: str, task_path: Path, root: Path) -> str:
    path = Path(logical)
    candidate = path if path.is_absolute() else root / path
    try:
        return (task_path / candidate.resolve().relative_to(task_path.resolve())).relative_to(root.resolve()).as_posix()
    except ValueError:
        return logical


def _resolve_registered_artifact_path(logical: str, root: Path, task_path: Path) -> Path:
    path = Path(logical)
    if path.is_absolute():
        return path.resolve()

    workspace_candidate = (root / path).resolve()
    if workspace_candidate.exists() or logical.startswith(".lab-sidecar/"):
        return workspace_candidate
    return (task_path / path).resolve()


def _csv_preview(path: Path, max_rows: int) -> dict[str, Any]:
    rows: list[dict[str, str]] = []
    with path.open(newline="", encoding="utf-8", errors="replace") as fh:
        reader = csv.DictReader(fh)
        columns = list(reader.fieldnames or [])
        for row in reader:
            rows.append(dict(row))
            if len(rows) > max_rows:
                break
    returned, truncated, withheld_complete_body = _bounded_preview_items(rows, max_rows)
    return {
        "columns": columns,
        "rows": returned,
        "row_count_returned": len(returned),
        "truncated": truncated,
        "withheld_complete_body": withheld_complete_body,
    }


def _text_lines_preview(path: Path, max_lines: int) -> dict[str, Any]:
    lines: list[str] = []
    with path.open(encoding="utf-8", errors="replace") as fh:
        for line in fh:
            lines.append(line.rstrip("\n")[:1000])
            if len(lines) > max_lines:
                break
    returned, truncated, withheld_complete_body = _bounded_preview_items(lines, max_lines)
    return {
        "lines": returned,
        "line_count_returned": len(returned),
        "truncated": truncated,
        "withheld_complete_body": withheld_complete_body,
    }


def _log_preview(path: Path, max_lines: int) -> dict[str, Any]:
    buffered: deque[str] = deque(maxlen=max_lines + 1)
    total_lines = 0
    with path.open(encoding="utf-8", errors="replace") as fh:
        for line in fh:
            total_lines += 1
            buffered.append(line.rstrip("\n")[:1000])
    lines = list(buffered)
    if total_lines > max_lines:
        returned = lines[-max_lines:]
        truncated = True
        withheld_complete_body = False
    elif lines:
        returned = lines[1:]
        truncated = True
        withheld_complete_body = True
    else:
        returned = []
        truncated = False
        withheld_complete_body = False
    return {
        "tail": returned,
        "line_count_returned": len(returned),
        "truncated": truncated,
        "withheld_complete_body": withheld_complete_body,
    }


def _image_preview(path: Path) -> dict[str, Any]:
    with Image.open(path) as image:
        return {
            "format": image.format,
            "width": image.width,
            "height": image.height,
            "mode": image.mode,
        }


def _pptx_preview(path: Path) -> dict[str, Any]:
    presentation = Presentation(path)
    titles: list[str] = []
    slide_count = len(presentation.slides)
    for index in range(min(slide_count, 10)):
        slide = presentation.slides[index]
        title = ""
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                title = shape.text.strip().splitlines()[0][:160]
                break
        titles.append(title)
    return {
        "slide_count": slide_count,
        "title_preview": titles,
        "title_preview_count": len(titles),
        "truncated": slide_count > len(titles),
    }


def _bounded(value: int, minimum: int, maximum: int) -> int:
    try:
        parsed = int(value)
    except (TypeError, ValueError):
        parsed = maximum
    return min(max(parsed, minimum), maximum)


def _bounded_preview_items(items: list[Any], maximum: int) -> tuple[list[Any], bool, bool]:
    if len(items) > maximum:
        return items[:maximum], True, False
    if items:
        return items[:-1], True, True
    return [], False, False
